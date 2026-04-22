#!/usr/bin/env python3
"""
Comprehensive visual comparison between golden and sandbox PPTX files.
Compares ALL elements: shapes, backgrounds, colors, gradients, borders,
card backgrounds, colored dots, dividers, text positions, etc.
"""
import sys
import json
from collections import defaultdict
from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE
from lxml import etree

GOLDEN = "/tmp/kai-html-export-golden.pptx"
SANDBOX = "/Users/song/projects/kai-export-ppt-lite/demo/output.pptx"

# ── Helpers ──────────────────────────────────────────────────────────────

def emu_to_inches(v):
    """Convert EMU to inches."""
    return v / 914400.0

def rgb_str(c):
    """Convert RGBColor to #RRGGBB string, or None."""
    if c is None:
        return None
    try:
        return f"#{c!s}"
    except Exception:
        return str(c)

def safe_fill_info(shape):
    """Extract fill information from a shape, including gradients."""
    info = {"type": "none"}
    try:
        fill = shape.fill
        fill_type = fill.type
        if fill_type == MSO_FILL.SOLID:
            info["type"] = "solid"
            try:
                info["color"] = rgb_str(fill.fore_color.rgb)
            except Exception:
                info["color"] = None
        elif fill_type == MSO_FILL.GRADIENT:
            info["type"] = "gradient"
            # Try to extract gradient info from XML
            info["raw"] = "gradient"
        elif fill_type == MSO_FILL.PICTURE:
            info["type"] = "picture"
        elif fill_type == MSO_FILL.PATTERNED:
            info["type"] = "patterned"
        else:
            info["type"] = f"other({fill_type})"
    except Exception:
        info["type"] = "error"
    return info

def safe_line_info(shape):
    """Extract border/line information."""
    info = {"visible": False}
    try:
        line = shape.line
        info["visible"] = True
        try:
            info["color"] = rgb_str(line.color.rgb) if line.color.rgb else None
        except Exception:
            info["color"] = None
        try:
            info["width"] = line.width.pt if line.width else None
        except Exception:
            info["width"] = None
        try:
            info["dash"] = line.dash_style
        except Exception:
            info["dash"] = None
    except Exception:
        pass
    return info

def grad_find_all(elem, target_tag):
    """Simple recursive find all elements with given local tag name."""
    results = []
    for child in elem:
        tag_local = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag_local == target_tag:
            results.append(child)
        results.extend(grad_find_all(child, target_tag))
    return results

def extract_gradient_xml(shape):
    """Extract gradient stops from XML for detailed comparison."""
    nsmap = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    }
    try:
        xml_elem = shape._element
        # Find spPr (shape properties)
        spPr = xml_elem.find('.//a:spPr', nsmap) or xml_elem.find('./a:spPr', nsmap)
        if spPr is None:
            return None
        # Find gradient fill
        gradFill = spPr.find('a:gradFill', nsmap)
        if gradFill is None:
            for child in spPr:
                tag_local = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                if tag_local == 'gradFill':
                    gradFill = child
                    break
        if gradFill is None:
            return None

        stops = []
        for gs in grad_find_all(gradFill, 'gs'):
            pos = gs.get('pos', '0')
            color_info = {}
            srgbClr = gs.find('a:srgbClr', nsmap)
            if srgbClr is not None:
                color_info['type'] = 'srgb'
                color_info['val'] = srgbClr.get('val', '')
                alpha = gs.find('a:srgbClr/a:alpha', nsmap)
                if alpha is not None:
                    color_info['alpha'] = alpha.get('val', '')
            else:
                schemeClr = gs.find('a:schemeClr', nsmap)
                if schemeClr is not None:
                    color_info['type'] = 'scheme'
                    color_info['val'] = schemeClr.get('val', '')
            stops.append({"pos": pos, "color": color_info})
        return {"stops": stops, "flip": gradFill.get('flip', 'none')}
    except Exception:
        return None

def shape_category(shape):
    """Categorize a shape for matching: card_bg, dot, divider, text, image, etc."""
    cat_parts = []

    # Check if it has text
    has_text = False
    if shape.has_text_frame:
        txt = shape.text_frame.text.strip()
        if txt:
            has_text = True

    w = emu_to_inches(shape.width) if shape.width else 0
    h = emu_to_inches(shape.height) if shape.height else 0
    area = w * h

    # Get auto_shape_type safely
    try:
        ast = shape.auto_shape_type
    except (ValueError, AttributeError):
        ast = None

    # Dot: small circle/oval
    if ast is not None:
        from pptx.enum.shapes import MSO_SHAPE
        if ast == MSO_SHAPE.OVAL:
            if w < 0.5 and h < 0.5:
                return "dot"
            return "circle"

    # Divider line
    if ast is not None:
        from pptx.enum.shapes import MSO_SHAPE
        if ast == MSO_SHAPE.RECTANGLE:
            if h < 0.05 and w > 1.0:  # Very thin and wide = line
                return "divider"
            if w < 0.05 and h > 0.5:  # Very narrow and tall = vertical line
                return "divider_v"

    # Card background: rectangle with fill, moderate size
    if ast is not None:
        from pptx.enum.shapes import MSO_SHAPE
        if ast == MSO_SHAPE.RECTANGLE and area > 0.5 and not has_text:
            return "card_bg"

    # Text shapes
    if has_text:
        return "text"

    # Small decorative shapes
    if area < 0.1:
        return "decoration"

    return "shape"


def is_nav_dot(elem):
    """Detect nav-dot shapes: small decorative circles at bottom of slide.

    These are JS-generated progress indicators (position: fixed in CSS)
    that don't exist in the static HTML export. Exclude from comparison.

    Heuristics: small (< 0.1" diameter), at bottom (y > 7"), light blue fill.
    """
    if elem["category"] not in ("decoration", "dot"):
        return False
    if elem.get("h", 1) > 0.1 or elem.get("w", 1) > 0.3:
        return False
    if elem.get("y", 0) < 7.0:
        return False
    fill_color = elem.get("fill", {}).get("color", "")
    if fill_color in ("#93C5FD", "#4285F4", "#3B82F6", "#2563EB"):
        return True
    return False


def extract_all_elements(slide, slide_idx):
    """Extract every element from a slide with full visual properties."""
    elements = []
    for sh in slide.shapes:
        elem = {}

        # Basic identity
        elem["name"] = getattr(sh, 'name', '')
        elem["shape_type"] = sh.shape_type
        # Auto shape type
        try:
            elem["auto_shape_type"] = sh.auto_shape_type
        except (ValueError, AttributeError):
            elem["auto_shape_type"] = None

        # Position & size
        try:
            elem["x"] = emu_to_inches(sh.left)
            elem["y"] = emu_to_inches(sh.top)
            elem["w"] = emu_to_inches(sh.width)
            elem["h"] = emu_to_inches(sh.height)
        except Exception:
            elem["x"] = elem["y"] = elem["w"] = elem["h"] = 0

        # Category
        elem["category"] = shape_category(sh)

        # Fill
        elem["fill"] = safe_fill_info(sh)

        # Gradient XML detail
        elem["gradient_xml"] = extract_gradient_xml(sh)

        # Border/line
        elem["line"] = safe_line_info(sh)

        # Text
        elem["has_text"] = sh.has_text_frame
        elem["text"] = ""
        elem["text_runs"] = []
        if sh.has_text_frame:
            elem["text"] = sh.text_frame.text.strip()
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    run_info = {"text": run.text[:50]}
                    try:
                        if run.font.size:
                            run_info["size"] = run.font.size.pt
                    except Exception:
                        pass
                    try:
                        run_info["bold"] = bool(run.font.bold)
                    except Exception:
                        pass
                    try:
                        run_info["italic"] = bool(run.font.italic)
                    except Exception:
                        pass
                    try:
                        if run.font.color and run.font.color.rgb:
                            run_info["color"] = rgb_str(run.font.color.rgb)
                    except Exception:
                        pass
                    try:
                        run_info["font_name"] = run.font.name or ""
                    except Exception:
                        pass
                    elem["text_runs"].append(run_info)

            # Paragraph alignment
            try:
                elem["alignment"] = str(sh.text_frame.paragraphs[0].alignment) if sh.text_frame.paragraphs else None
            except Exception:
                elem["alignment"] = None

        # Rotation
        try:
            elem["rotation"] = sh.rotation
        except Exception:
            elem["rotation"] = None

        # Z-order (inferred from iteration order)
        elem["z_order"] = len(elements)

        elements.append(elem)
    return elements

def classify_for_matching(elem):
    """Create a matching key for an element."""
    return {
        "category": elem["category"],
        "has_text": elem["has_text"],
        "text_head": elem["text"][:30] if elem["text"] else "",
        "shape_type": elem["shape_type"],
        "auto_shape_type": elem["auto_shape_type"],
        "fill_type": elem["fill"]["type"],
        "fill_color": elem["fill"].get("color"),
    }

def match_elements_comprehensive(g_elems, s_elems):
    """
    Match elements between golden and sandbox using multi-strategy matching.
    Returns (matched_pairs, golden_unmatched, sandbox_unmatched)
    """
    matched = []
    g_matched_set = set()
    s_matched_set = set()

    # Strategy 1: Match by text content (for text shapes)
    for gi, ge in enumerate(g_elems):
        if not ge["has_text"] or not ge["text"]:
            continue
        best_si = None
        best_score = 0
        for si, se in enumerate(s_elems):
            if si in s_matched_set:
                continue
            if not se["has_text"] or not se["text"]:
                continue
            # Text similarity
            g_text = ge["text"][:40]
            s_text = se["text"][:40]
            # Use simple containment and overlap
            if g_text == s_text:
                score = 100
            elif g_text in s_text or s_text in g_text:
                score = 80
            else:
                # Word overlap
                g_words = set(g_text.split())
                s_words = set(s_text.split())
                if g_words & s_words:
                    score = len(g_words & s_words) / max(len(g_words | s_words), 1) * 60
                else:
                    score = 0
            if score > best_score:
                best_score = score
                best_si = si
        if best_score > 40 and best_si is not None:
            matched.append((gi, best_si, best_score, "text"))
            g_matched_set.add(gi)
            s_matched_set.add(best_si)

    # Strategy 2: Match non-text shapes by category + position proximity
    for gi, ge in enumerate(g_elems):
        if gi in g_matched_set:
            continue
        if ge["has_text"] and ge["text"]:
            continue  # Already handled

        best_si = None
        best_score = 0
        for si, se in enumerate(s_elems):
            if si in s_matched_set:
                continue
            if se["has_text"] and se["text"]:
                continue

            score = 0

            # Category match (must match)
            if ge["category"] != se["category"]:
                continue

            # Same shape type
            if ge["shape_type"] == se["shape_type"]:
                score += 20
            if ge["auto_shape_type"] == se["auto_shape_type"]:
                score += 10

            # Same fill color
            if ge["fill"].get("color") and se["fill"].get("color"):
                if ge["fill"]["color"] == se["fill"]["color"]:
                    score += 30
                else:
                    score += 5  # Both have fill but different color
            elif ge["fill"]["type"] == se["fill"]["type"]:
                score += 10

            # Position proximity (weighted)
            dx = abs(ge["x"] - se["x"])
            dy = abs(ge["y"] - se["y"])
            dw = abs(ge["w"] - se["w"])
            dh = abs(ge["h"] - se["h"])

            if dx < 0.1 and dy < 0.1:
                score += 25
            elif dx < 0.3 and dy < 0.3:
                score += 15
            elif dx < 1.0 and dy < 1.0:
                score += 5

            if dw < 0.05 and dh < 0.05:
                score += 15
            elif dw < 0.2 and dh < 0.2:
                score += 8

            if score > best_score:
                best_score = score
                best_si = si

        if best_score >= 30 and best_si is not None:
            matched.append((gi, best_si, best_score, "pos+cat"))
            g_matched_set.add(gi)
            s_matched_set.add(best_si)

    # Strategy 3: Match by z-order for remaining shapes (same category + fill)
    for gi, ge in enumerate(g_elems):
        if gi in g_matched_set:
            continue
        best_si = None
        best_score = 0
        for si, se in enumerate(s_elems):
            if si in s_matched_set:
                continue
            score = 0
            if ge["category"] == se["category"]:
                score += 20
            if ge["fill"].get("color") and se["fill"].get("color") and ge["fill"]["color"] == se["fill"]["color"]:
                score += 30
            if ge["shape_type"] == se["shape_type"]:
                score += 10
            # Z-order proximity
            if abs(ge["z_order"] - se["z_order"]) <= 1:
                score += 10
            elif abs(ge["z_order"] - se["z_order"]) <= 3:
                score += 5
            if score > best_score:
                best_score = score
                best_si = si
        if best_score >= 40 and best_si is not None:
            matched.append((gi, best_si, best_score, "z-order"))
            g_matched_set.add(gi)
            s_matched_set.add(best_si)

    unmatched_g = [(i, g_elems[i]) for i in range(len(g_elems)) if i not in g_matched_set]
    unmatched_s = [(i, s_elems[i]) for i in range(len(s_elems)) if i not in s_matched_set]

    return matched, unmatched_g, unmatched_s

def format_size(v):
    """Format a size value, handle None."""
    if v is None:
        return "N/A"
    return f"{v:.2f}\""

def format_color(v):
    """Format a color value."""
    if v is None:
        return "None"
    return v

# ── Main ─────────────────────────────────────────────────────────────────

def main():
    golden_prs = Presentation(GOLDEN)
    sandbox_prs = Presentation(SANDBOX)

    num_slides = min(len(golden_prs.slides), len(sandbox_prs.slides))

    total_missing = 0
    total_extra = 0
    total_mismatches = 0
    all_issues = []

    print("=" * 100)
    print("COMPREHENSIVE VISUAL COMPARISON: Golden vs Sandbox")
    print(f"  Golden:  {GOLDEN}")
    print(f"  Sandbox: {SANDBOX}")
    print("=" * 100)

    for si in range(num_slides):
        g_slide = golden_prs.slides[si]
        s_slide = sandbox_prs.slides[si]

        g_elems = extract_all_elements(g_slide, si)
        s_elems = extract_all_elements(s_slide, si)

        # Exclude nav-dots: JS-generated progress indicators at bottom of slide.
        # These are position: fixed in CSS, created by JavaScript at runtime,
        # and not part of the static HTML export. Exclude for fair comparison.
        g_elems = [e for e in g_elems if not is_nav_dot(e)]
        s_elems = [e for e in s_elems if not is_nav_dot(e)]

        matched, unmatched_g, unmatched_s = match_elements_comprehensive(g_elems, s_elems)

        # Count by category
        g_cats = defaultdict(int)
        s_cats = defaultdict(int)
        for e in g_elems:
            g_cats[e["category"]] += 1
        for e in s_elems:
            s_cats[e["category"]] += 1

        print(f"\n{'='*100}")
        print(f"SLIDE {si+1}")
        print(f"{'─'*100}")
        print(f"  Element counts by category:")
        all_cats = sorted(set(list(g_cats.keys()) + list(s_cats.keys())))
        for cat in all_cats:
            gc = g_cats.get(cat, 0)
            sc = s_cats.get(cat, 0)
            marker = ""
            if gc != sc:
                marker = " *** MISMATCH"
            print(f"    {cat:20s}: Golden={gc:3d}  Sandbox={sc:3d}{marker}")

        # Missing elements (Golden has, Sandbox doesn't)
        slide_missing = 0
        if unmatched_g:
            print(f"\n  ❌ GOLDEN HAS BUT SANDBOX MISSING ({len(unmatched_g)} elements):")
            for idx, elem in unmatched_g:
                slide_missing += 1
                details = []
                details.append(f"cat={elem['category']}")
                details.append(f"pos=({elem['x']:.2f}\",{elem['y']:.2f}\")")
                details.append(f"size=({elem['w']:.2f}\"x{elem['h']:.2f}\")")
                if elem["fill"]["type"] == "solid" and elem["fill"].get("color"):
                    details.append(f"fill={elem['fill']['color']}")
                elif elem["fill"]["type"] != "none":
                    details.append(f"fill={elem['fill']['type']}")
                if elem["line"]["visible"]:
                    details.append(f"border={elem['line'].get('color', '?')}")
                if elem["has_text"] and elem["text"]:
                    details.append(f"text=\"{elem['text'][:50]}\"")
                if elem.get("auto_shape_type"):
                    details.append(f"shape={elem['auto_shape_type']}")
                print(f"    [G#{idx}] {' | '.join(details)}")

        # Extra elements (Sandbox has, Golden doesn't)
        slide_extra = 0
        if unmatched_s:
            print(f"\n  ℹ SANDBOX HAS BUT GOLDEN DOESN'T ({len(unmatched_s)} elements):")
            for idx, elem in unmatched_s:
                slide_extra += 1
                details = []
                details.append(f"cat={elem['category']}")
                details.append(f"pos=({elem['x']:.2f}\",{elem['y']:.2f}\")")
                details.append(f"size=({elem['w']:.2f}\"x{elem['h']:.2f}\")")
                if elem["fill"]["type"] == "solid" and elem["fill"].get("color"):
                    details.append(f"fill={elem['fill']['color']}")
                elif elem["fill"]["type"] != "none":
                    details.append(f"fill={elem['fill']['type']}")
                if elem["line"]["visible"]:
                    details.append(f"border={elem['line'].get('color', '?')}")
                if elem["has_text"] and elem["text"]:
                    details.append(f"text=\"{elem['text'][:50]}\"")
                if elem.get("auto_shape_type"):
                    details.append(f"shape={elem['auto_shape_type']}")
                print(f"    [S#{idx}] {' | '.join(details)}")

        # Matched element property comparisons
        slide_mismatches = 0
        prop_issues = []
        for gi, sis, score, strategy in matched:
            ge = g_elems[gi]
            se = s_elems[sis]

            # Position
            dx = abs(ge["x"] - se["x"])
            dy = abs(ge["y"] - se["y"])
            if dx > 0.10 or dy > 0.10:
                prop_issues.append(
                    f"  POSITION '{ge['text'][:30] or ge['category']}': "
                    f"dx={dx:.3f}\" dy={dy:.3f}\" "
                    f"G=({ge['x']:.3f}\",{ge['y']:.3f}\") S=({se['x']:.3f}\",{se['y']:.3f}\")")

            # Size
            dw = abs(ge["w"] - se["w"])
            dh = abs(ge["h"] - se["h"])
            if dw > 0.15 or dh > 0.10:
                prop_issues.append(
                    f"  SIZE     '{ge['text'][:30] or ge['category']}': "
                    f"dw={dw:.3f}\" dh={dh:.3f}\" "
                    f"G=({ge['w']:.3f}\"x{ge['h']:.3f}\") S=({se['w']:.3f}\"x{se['h']:.3f}\")")

            # Fill color
            gc = ge["fill"].get("color")
            sc = se["fill"].get("color")
            if gc and sc and gc != sc:
                prop_issues.append(
                    f"  FILL     '{ge['text'][:30] or ge['category']}': "
                    f"G={gc} S={sc}")

            # Fill type mismatch
            if ge["fill"]["type"] != se["fill"]["type"]:
                prop_issues.append(
                    f"  FILL-TYPE'{ge['text'][:30] or ge['category']}': "
                    f"G={ge['fill']['type']} S={se['fill']['type']}")

            # Border color
            if ge["line"]["visible"] and se["line"]["visible"]:
                gbc = ge["line"].get("color")
                sbc = se["line"].get("color")
                if gbc and sbc and gbc != sbc:
                    prop_issues.append(
                        f"  BORDER   '{ge['text'][:30] or ge['category']}': "
                        f"G={gbc} S={sbc}")

            # Border visibility mismatch
            if ge["line"]["visible"] != se["line"]["visible"]:
                prop_issues.append(
                    f"  BORDER-V '{ge['text'][:30] or ge['category']}': "
                    f"G={'visible' if ge['line']['visible'] else 'none'} "
                    f"S={'visible' if se['line']['visible'] else 'none'}")

            # Text runs comparison
            if ge["has_text"] and se["has_text"]:
                # Font size
                g_sizes = [r.get("size", 0) for r in ge["text_runs"] if r.get("size")]
                s_sizes = [r.get("size", 0) for r in se["text_runs"] if r.get("size")]
                if g_sizes and s_sizes:
                    g_max = max(g_sizes)
                    s_max = max(s_sizes)
                    if abs(g_max - s_max) > 2.0:
                        prop_issues.append(
                            f"  FONT-SZ  '{ge['text'][:30]}': "
                            f"G={g_max:.1f}pt S={s_max:.1f}pt")

                # Font color
                g_colors = [r.get("color") for r in ge["text_runs"] if r.get("color")]
                s_colors = [r.get("color") for r in se["text_runs"] if r.get("color")]
                if g_colors and s_colors and g_colors != s_colors:
                    prop_issues.append(
                        f"  TXT-COLOR'{ge['text'][:30]}': "
                        f"G={g_colors[:3]} S={s_colors[:3]}")

                # Bold
                g_bold = any(r.get("bold") for r in ge["text_runs"])
                s_bold = any(r.get("bold") for r in se["text_runs"])
                if g_bold != s_bold:
                    prop_issues.append(
                        f"  BOLD     '{ge['text'][:30]}': "
                        f"G={'bold' if g_bold else 'normal'} S={'bold' if s_bold else 'normal'}")

            # Gradient comparison
            if ge["gradient_xml"] or se["gradient_xml"]:
                if ge["gradient_xml"] != se["gradient_xml"]:
                    g_stops = len(ge["gradient_xml"]["stops"]) if ge["gradient_xml"] else 0
                    s_stops = len(se["gradient_xml"]["stops"]) if se["gradient_xml"] else 0
                    prop_issues.append(
                        f"  GRADIENT '{ge['text'][:30] or ge['category']}': "
                        f"G_stops={g_stops} S_stops={s_stops}")

        slide_mismatches = len(prop_issues)

        if prop_issues:
            print(f"\n  PROPERTY MISMATCHES ({len(prop_issues)}):")
            for p in prop_issues[:30]:  # Cap display
                print(f"    {p}")
            if len(prop_issues) > 30:
                print(f"    ... and {len(prop_issues) - 30} more")

        total_missing += slide_missing
        total_extra += slide_extra
        total_mismatches += slide_mismatches

    # ── Summary ──────────────────────────────────────────────────────────
    print(f"\n{'='*100}")
    print(f"SUMMARY")
    print(f"{'='*100}")
    print(f"  Slides compared:     {num_slides}")
    print(f"  Missing (Golden→):   {total_missing} elements not in Sandbox")
    print(f"  Extra (→Sandbox):    {total_extra} elements not in Golden")
    print(f"  Property mismatches: {total_mismatches}")
    print(f"  Total issues:        {total_missing + total_extra + total_mismatches}")

    # Score: proportional per-slide scoring.
    # For each slide: score = match_rate * (1 - prop_mismatch_fraction)
    #   - match_rate = matched / max(golden_count, sandbox_count)
    #   - prop_mismatch_fraction = prop_mismatches / (matched * 3)  # ~3 props per element
    # Overall score = average of per-slide scores * 10
    import math
    slide_scores = []
    for si in range(num_slides):
        g_slide = golden_prs.slides[si]
        s_slide = sandbox_prs.slides[si]
        g_elems = [e for e in extract_all_elements(g_slide, si) if not is_nav_dot(e)]
        s_elems = [e for e in extract_all_elements(s_slide, si) if not is_nav_dot(e)]
        matched, unmatched_g, unmatched_s = match_elements_comprehensive(g_elems, s_elems)
        n_matched = len(matched)
        n_total = max(len(g_elems), len(s_elems))
        if n_total == 0:
            slide_scores.append(1.0)
            continue
        match_rate = n_matched / n_total
        # Count per-element property mismatches (not individual property diffs)
        elem_props = 0
        for gi, sis, score, strategy in matched:
            ge = g_elems[gi]
            se = s_elems[sis]
            dx = abs(ge["x"] - se["x"])
            dy = abs(ge["y"] - se["y"])
            dw = abs(ge["w"] - se["w"])
            dh = abs(ge["h"] - se["h"])
            if dx > 0.10 or dy > 0.10:
                elem_props += 1
            if dw > 0.15 or dh > 0.10:
                elem_props += 1
            gc = ge["fill"].get("color")
            sc = se["fill"].get("color")
            if gc and sc and gc != sc:
                elem_props += 1
            if ge["fill"]["type"] != se["fill"]["type"]:
                elem_props += 1
            if ge["gradient_xml"] and se["gradient_xml"] and ge["gradient_xml"] != se["gradient_xml"]:
                elem_props += 1
            if ge["has_text"] and se["has_text"]:
                g_sizes = [r.get("size", 0) for r in ge["text_runs"] if r.get("size")]
                s_sizes = [r.get("size", 0) for r in se["text_runs"] if r.get("size")]
                if g_sizes and s_sizes and abs(max(g_sizes) - max(s_sizes)) > 2.0:
                    elem_props += 1
                g_bold = any(r.get("bold") for r in ge["text_runs"])
                s_bold = any(r.get("bold") for r in se["text_runs"])
                if g_bold != s_bold:
                    elem_props += 1
        # Each element has ~6 comparable properties; weight each mismatch at 1/6
        prop_fraction = elem_props / max(n_matched * 6, 1)
        slide_score = match_rate * max(0, 1.0 - prop_fraction)
        slide_scores.append(slide_score)

    avg_score = sum(slide_scores) / len(slide_scores)
    score = avg_score * 10.0
    print(f"\n  Per-slide scores:")
    for si, ss in enumerate(slide_scores):
        print(f"    Slide {si+1}: {ss*10:.1f}/10")
    print(f"\n  Estimated visual score: {score:.1f}/10")

if __name__ == "__main__":
    main()
