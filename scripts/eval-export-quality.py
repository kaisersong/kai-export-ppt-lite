#!/usr/bin/env python3
"""
eval-export-quality.py — Evaluate HTML→PPTX export quality.

Compares sandbox output against golden reference (Playwright version).
Scores on a 10-point scale across 5 dimensions.

Usage:
    python3 scripts/eval-export-quality.py <golden.pptx> <sandbox.pptx>
"""

import sys
import re
from pathlib import Path
from dataclasses import dataclass, field
from typing import List, Dict, Any, Optional, Tuple

try:
    from pptx import Presentation
    from pptx.enum.dml import MSO_FILL
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from PIL import Image, ImageChops
except ImportError as e:
    print(f"Missing dependency: {e}")
    print("Install with: pip install python-pptx Pillow")
    sys.exit(1)


@dataclass
class TextElement:
    text: str
    x: float  # normalized 0-1
    y: float
    w: float
    h: float
    font_size: float  # pt
    is_bold: bool
    color: Tuple[int, int, int] = (0, 0, 0)
    shape_type: str = "text"


@dataclass
class SlideData:
    index: int
    bg_color: Optional[Tuple[int, int, int]]
    elements: List[TextElement]
    element_count: int


@dataclass
class EvalResult:
    text_accuracy: float = 0.0      # 0-3
    layout_accuracy: float = 0.0     # 0-2
    visual_accuracy: float = 0.0     # 0-2
    special_elements: float = 0.0    # 0-2
    readability: float = 0.0         # 0-1
    total: float = 0.0
    notes: List[str] = field(default_factory=list)


def get_slide_bg_color(slide) -> Optional[Tuple[int, int, int]]:
    try:
        ft = slide.background.fill.type
        if ft == MSO_FILL.SOLID:
            return slide.background.fill.fore_color.rgb
    except Exception:
        pass
    return None


def extract_slide_data(slide, slide_idx: int, slide_w_emu=None, slide_h_emu=None) -> SlideData:
    """Extract structured data from a PPTX slide."""
    slide_w = slide_w_emu or 12192000  # default 13.33in in EMU
    slide_h = slide_h_emu or 6858000   # default 7.5in in EMU

    bg = get_slide_bg_color(slide)
    elements = []

    for shape in slide.shapes:
        try:
            x = shape.left / slide_w
            y = shape.top / slide_h
            w = shape.width / slide_w
            h = shape.height / slide_h
        except Exception:
            continue

        if w < 0.001 or h < 0.001:
            continue

        # Skip nav dots (very small)
        if w < 0.02 and h < 0.02:
            continue

        # Counter (page number)
        if w < 0.08 and h < 0.04 and shape.has_text_frame:
            txt = shape.text.strip()
            if re.match(r'\d{2}\s*/\s*\d{2}', txt):
                continue

        if shape.has_text_frame and shape.text.strip():
            font_size = 12.0
            is_bold = False
            color = (0, 0, 0)

            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        font_size = max(font_size, run.font.size.pt)
                    if run.font.bold:
                        is_bold = True
                    try:
                        if run.font.color and run.font.color.rgb:
                            color = run.font.color.rgb
                    except Exception:
                        pass

            elements.append(TextElement(
                text=shape.text.strip(),
                x=x, y=y, w=w, h=h,
                font_size=font_size,
                is_bold=is_bold,
                color=color,
                shape_type="text"
            ))
        elif shape.has_text_frame and not shape.text.strip():
            # AUTO_SHAPE with no text content — treat as shape
            fill_color = None
            try:
                if shape.fill.type == MSO_FILL.SOLID:
                    fill_color = shape.fill.fore_color.rgb
            except Exception:
                pass
            elements.append(TextElement(
                text="",
                x=x, y=y, w=w, h=h,
                font_size=0,
                is_bold=False,
                color=fill_color or (0, 0, 0),
                shape_type="shape"
            ))
        elif not shape.has_text_frame:
            # Shape (background, card, etc.)
            fill_color = None
            try:
                if shape.fill.type == MSO_FILL.SOLID:
                    fill_color = shape.fill.fore_color.rgb
            except Exception:
                pass
            elements.append(TextElement(
                text="",
                x=x, y=y, w=w, h=h,
                font_size=0,
                is_bold=False,
                color=fill_color or (0, 0, 0),
                shape_type="shape"
            ))

    return SlideData(
        index=slide_idx,
        bg_color=bg,
        elements=elements,
        element_count=len(elements)
    )


def levenshtein(s1, s2):
    """Compute Levenshtein distance."""
    if len(s1) < len(s2):
        return levenshtein(s2, s1)
    if len(s2) == 0:
        return len(s1)
    prev_row = range(len(s2) + 1)
    for i, c1 in enumerate(s1):
        curr_row = [i + 1]
        for j, c2 in enumerate(s2):
            insertions = prev_row[j + 1] + 1
            deletions = curr_row[j] + 1
            substitutions = prev_row[j] + (c1 != c2)
            curr_row.append(min(insertions, deletions, substitutions))
        prev_row = curr_row
    return prev_row[-1]


def text_similarity(s1, s2):
    """Compute text similarity ratio (0-1)."""
    s1_clean = re.sub(r'\s+', ' ', s1.strip())
    s2_clean = re.sub(r'\s+', ' ', s2.strip())
    if not s1_clean and not s2_clean:
        return 1.0
    if not s1_clean or not s2_clean:
        return 0.0
    max_len = max(len(s1_clean), len(s2_clean))
    dist = levenshtein(s1_clean, s2_clean)
    return max(0.0, 1.0 - dist / max_len)


def evaluate_text_accuracy(golden_slides: List[SlideData],
                           sandbox_slides: List[SlideData]) -> Tuple[float, List[str]]:
    """Score text correctness (0-3)."""
    notes = []
    total_score = 0.0
    max_score = 0.0

    for gs, ss in zip(golden_slides, sandbox_slides):
        g_texts = sorted([e.text for e in gs.elements if e.text])
        s_texts = sorted([e.text for e in ss.elements if e.text])

        # Check text coverage
        matched = 0
        for gt in g_texts:
            best_sim = 0.0
            for st in s_texts:
                sim = text_similarity(gt, st)
                if sim > best_sim:
                    best_sim = sim
            if best_sim > 0.7:
                matched += 1
            elif best_sim > 0.4:
                matched += best_sim

        coverage = matched / max(len(g_texts), 1)
        total_score += coverage * 3.0 / len(golden_slides)
        max_score += 3.0 / len(golden_slides)

        if coverage < 0.8:
            missing = [t for t in g_texts if not any(text_similarity(t, s) > 0.7 for s in s_texts)]
            if missing:
                notes.append(f"Slide {gs.index + 1}: missing {len(missing)} text elements: {missing[:3]}")

    return min(3.0, total_score), notes


def evaluate_layout_accuracy(golden_slides: List[SlideData],
                             sandbox_slides: List[SlideData]) -> Tuple[float, List[str]]:
    """Score layout correctness (0-2)."""
    notes = []
    total_score = 0.0

    for gs, ss in zip(golden_slides, sandbox_slides):
        g_texts = [e for e in gs.elements if e.text]
        s_texts = [e for e in ss.elements if e.text]

        if not g_texts:
            total_score += 2.0 / len(golden_slides)
            continue

        # Match elements by text similarity, then compare positions
        used = set()
        pos_scores = []
        for ge in g_texts:
            best_score = 0.0
            for si, se in enumerate(s_texts):
                if si in used:
                    continue
                sim = text_similarity(ge.text, se.text)
                if sim > 0.7:
                    # Compare normalized positions
                    dx = abs(ge.x - se.x)
                    dy = abs(ge.y - se.y)
                    # Allow some tolerance for layout differences
                    pos_score = max(0, 1.0 - (dx + dy) * 2)
                    if pos_score > best_score:
                        best_score = pos_score
                        used.add(si)

            pos_scores.append(best_score)

        avg_pos = sum(pos_scores) / max(len(pos_scores), 1)
        # Also penalize wrong element count
        count_ratio = min(len(s_texts), len(g_texts)) / max(len(g_texts), 1)
        slide_score = (avg_pos * 0.6 + count_ratio * 0.4) * 2.0
        total_score += slide_score / len(golden_slides)

        if avg_pos < 0.5:
            notes.append(f"Slide {gs.index + 1}: layout positions significantly off (avg similarity: {avg_pos:.2f})")

    return min(2.0, total_score), notes


def evaluate_visual_accuracy(golden_slides: List[SlideData],
                             sandbox_slides: List[SlideData]) -> Tuple[float, List[str]]:
    """Score visual correctness - colors, fonts, bold (0-2)."""
    notes = []
    total_score = 0.0

    for gs, ss in zip(golden_slides, sandbox_slides):
        g_texts = [e for e in gs.elements if e.text]
        s_texts = [e for e in ss.elements if e.text]

        if not g_texts:
            total_score += 2.0 / len(golden_slides)
            continue

        font_scores = []
        bold_scores = []
        color_scores = []

        used = set()
        for ge in g_texts:
            best_font = 0.0
            best_bold = 0.0
            best_color = 0.0
            for si, se in enumerate(s_texts):
                if si in used:
                    continue
                sim = text_similarity(ge.text, se.text)
                if sim > 0.7:
                    # Font size (allow ±3pt tolerance)
                    if ge.font_size > 0 and se.font_size > 0:
                        font_sim = max(0, 1.0 - abs(ge.font_size - se.font_size) / max(ge.font_size, 1))
                    else:
                        font_sim = 1.0 if ge.font_size == 0 and se.font_size == 0 else 0.0

                    # Bold match
                    bold_sim = 1.0 if ge.is_bold == se.is_bold else 0.0

                    # Color match (allow some tolerance)
                    if ge.color and se.color:
                        color_diff = sum(abs(a - b) for a, b in zip(ge.color, se.color)) / (3 * 255)
                        color_sim = max(0, 1.0 - color_diff)
                    else:
                        color_sim = 1.0

                    if font_sim > best_font:
                        best_font = font_sim
                        best_bold = bold_sim
                        best_color = color_sim
                        used.add(si)

            font_scores.append(best_font)
            bold_scores.append(best_bold)
            color_scores.append(best_color)

        avg_font = sum(font_scores) / max(len(font_scores), 1)
        avg_bold = sum(bold_scores) / max(len(bold_scores), 1)
        avg_color = sum(color_scores) / max(len(color_scores), 1)

        slide_score = (avg_font * 0.4 + avg_bold * 0.3 + avg_color * 0.3) * 2.0
        total_score += slide_score / len(golden_slides)

        if avg_font < 0.6:
            notes.append(f"Slide {gs.index + 1}: font sizes significantly off (avg: {avg_font:.2f})")

    return min(2.0, total_score), notes


def evaluate_special_elements(golden_slides: List[SlideData],
                              sandbox_slides: List[SlideData]) -> Tuple[float, List[str]]:
    """Score tables, cards, backgrounds, rounded corners (0-2)."""
    notes = []
    total_score = 0.0

    for gs, ss in zip(golden_slides, sandbox_slides):
        g_shapes = [e for e in gs.elements if e.shape_type == "shape"]
        s_shapes = [e for e in ss.elements if e.shape_type == "shape"]

        # Check shape count match
        shape_count_score = min(len(s_shapes), len(g_shapes)) / max(len(g_shapes), 1)

        # Check table elements (golden may have table elements encoded differently)
        # For now, just check that we have comparable element counts
        g_total = len(gs.elements)
        s_total = len(ss.elements)
        count_score = min(s_total, g_total) / max(g_total, 1)

        slide_score = (shape_count_score * 0.5 + count_score * 0.5) * 2.0
        total_score += slide_score / len(golden_slides)

        if shape_count_score < 0.5:
            notes.append(f"Slide {gs.index + 1}: shape count mismatch (golden={len(g_shapes)}, sandbox={len(s_shapes)})")

    return min(2.0, total_score), notes


def evaluate_readability(golden_slides: List[SlideData],
                         sandbox_slides: List[SlideData]) -> Tuple[float, List[str]]:
    """Score overall readability (0-1).

    Check: no severe overlap, text is legible, no truncation.
    """
    notes = []
    total_score = 0.0

    for gs, ss in zip(golden_slides, sandbox_slides):
        s_texts = [e for e in ss.elements if e.text]

        # Check for severe height underestimation
        # In golden, text elements typically have h > 0.02 (normalized)
        # if they have meaningful text
        height_issues = 0
        for e in s_texts:
            if len(e.text) > 20 and e.h < 0.03:
                height_issues += 1

        # Check text truncation (very short elements with long text)
        truncation_issues = 0
        for e in s_texts:
            if len(e.text) > 50 and e.h < 0.04:
                truncation_issues += 1

        issue_ratio = (height_issues + truncation_issues) / max(len(s_texts), 1)
        slide_score = max(0, 1.0 - issue_ratio)
        total_score += slide_score / len(golden_slides)

        if height_issues > 0:
            notes.append(f"Slide {gs.index + 1}: {height_issues} elements with severely underestimated height")

    return min(1.0, total_score), notes


def evaluate(golden_path: Path, sandbox_path: Path) -> EvalResult:
    """Full evaluation."""
    golden_prs = Presentation(str(golden_path))
    sandbox_prs = Presentation(str(sandbox_path))
    golden_w = golden_prs.slide_width
    golden_h = golden_prs.slide_height
    sandbox_w = sandbox_prs.slide_width
    sandbox_h = sandbox_prs.slide_height

    num_slides = min(len(golden_prs.slides), len(sandbox_prs.slides))

    golden_slides = []
    for i in range(num_slides):
        golden_slides.append(extract_slide_data(golden_prs.slides[i], i, golden_w, golden_h))

    sandbox_slides = []
    for i in range(num_slides):
        sandbox_slides.append(extract_slide_data(sandbox_prs.slides[i], i, sandbox_w, sandbox_h))

    result = EvalResult()

    # Text accuracy (0-3)
    result.text_accuracy, notes = evaluate_text_accuracy(golden_slides, sandbox_slides)
    result.notes.extend(notes)

    # Layout accuracy (0-2)
    result.layout_accuracy, notes = evaluate_layout_accuracy(golden_slides, sandbox_slides)
    result.notes.extend(notes)

    # Visual accuracy (0-2)
    result.visual_accuracy, notes = evaluate_visual_accuracy(golden_slides, sandbox_slides)
    result.notes.extend(notes)

    # Special elements (0-2)
    result.special_elements, notes = evaluate_special_elements(golden_slides, sandbox_slides)
    result.notes.extend(notes)

    # Readability (0-1)
    result.readability, notes = evaluate_readability(golden_slides, sandbox_slides)
    result.notes.extend(notes)

    result.total = (result.text_accuracy + result.layout_accuracy +
                   result.visual_accuracy + result.special_elements + result.readability)

    return result


def print_report(result: EvalResult):
    print("=" * 60)
    print("HTML→PPTX Export Quality Evaluation")
    print("=" * 60)
    print(f"\n  {'Dimension':<20s} {'Score':>6s} {'Max':>4s}")
    print(f"  {'─' * 17:<20s} {'─' * 4:>6s} {'─' * 3:>4s}")
    print(f"  {'文字正确性':<20s} {result.text_accuracy:>6.2f} {3.0:>4.1f}")
    print(f"  {'布局正确性':<20s} {result.layout_accuracy:>6.2f} {2.0:>4.1f}")
    print(f"  {'视觉正确性':<20s} {result.visual_accuracy:>6.2f} {2.0:>4.1f}")
    print(f"  {'特殊元素':<20s} {result.special_elements:>6.2f} {2.0:>4.1f}")
    print(f"  {'整体可读性':<20s} {result.readability:>6.2f} {1.0:>4.1f}")
    print(f"  {'─' * 17:<20s} {'─' * 4:>6s} {'─' * 3:>4s}")
    print(f"  {'TOTAL':<20s} {result.total:>6.2f} {10.0:>4.1f}")

    if result.notes:
        print(f"\n  Issues:")
        for note in result.notes:
            print(f"    ⚠ {note}")

    # Grade
    if result.total >= 9.0:
        grade = "EXCELLENT ✓"
    elif result.total >= 7.0:
        grade = "GOOD"
    elif result.total >= 5.0:
        grade = "FAIR"
    else:
        grade = "NEEDS WORK"
    print(f"\n  Grade: {grade} ({result.total:.1f}/10)")
    print("=" * 60)


def main():
    if len(sys.argv) < 3:
        print(f"Usage: {sys.argv[0]} <golden.pptx> <sandbox.pptx>")
        sys.exit(1)

    golden_path = Path(sys.argv[1])
    sandbox_path = Path(sys.argv[2])

    if not golden_path.exists():
        print(f"Golden reference not found: {golden_path}")
        sys.exit(1)
    if not sandbox_path.exists():
        print(f"Sandbox output not found: {sandbox_path}")
        sys.exit(1)

    result = evaluate(golden_path, sandbox_path)
    print_report(result)

    return 0 if result.total >= 9.0 else 1


if __name__ == "__main__":
    sys.exit(main())
