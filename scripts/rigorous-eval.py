#!/usr/bin/env python3
"""
rigorous-eval.py — Multi-dimensional HTML→PPTX export quality evaluation.

Runs 5 independent checks in parallel:
1. compare-visual-comprehensive.py — pixel/property match vs golden
2. overflow-check — text wider than containing shape/card
3. overlap-check — adjacent elements overlapping beyond tolerance
4. position-drift — systematic offset analysis (Y shift, X shift distribution)
5. element-count-gap — missing/extra element categorization

Usage:
    python3 scripts/rigorous-eval.py
"""

import argparse
import sys
import subprocess
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import List, Dict, Tuple, Optional

GOLDEN = "/tmp/kai-html-export-golden.pptx"
SANDBOX = "demo/output.pptx"

EMU_PER_IN = 914400


# ─── 1. Visual comprehensive (wrapper) ─────────────────────────────────

def run_visual_comprehensive(golden_path: str = GOLDEN, sandbox_path: str = SANDBOX) -> str:
    """Run compare-visual-comprehensive.py and return summary."""
    if Path(golden_path).resolve() != Path(GOLDEN).resolve() or Path(sandbox_path).resolve() != Path(SANDBOX).resolve():
        return "  SKIP: visual comprehensive only supports the default GOLDEN/SANDBOX pair"
    result = subprocess.run(
        [sys.executable, "scripts/compare-visual-comprehensive.py"],
        capture_output=True, text=True, timeout=60
    )
    # Extract per-slide scores and total
    lines = result.stdout.split("\n")
    summary = []
    in_summary = False
    for line in lines:
        if "SUMMARY" in line:
            in_summary = True
        if in_summary:
            summary.append(line)
    return "\n".join(summary) if summary else result.stdout[-500:]


# ─── 2. Overflow check ─────────────────────────────────────────────────

def _is_nav_dot(shape) -> bool:
    try:
        if shape.width / EMU_PER_IN < 0.1 and shape.height / EMU_PER_IN < 0.1:
            return True
    except Exception:
        pass
    return False


def _get_fill_color(shape) -> Optional[str]:
    try:
        if shape.fill.type == MSO_FILL.SOLID:
            return f"#{shape.fill.fore_color.rgb}"
    except Exception:
        pass
    return None


def _text_runs_info(shape) -> List[Dict]:
    """Extract text run info: text, font_size_pt, is_bold, color."""
    runs = []
    if not hasattr(shape, "text_frame"):
        return runs
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            runs.append({
                "text": run.text,
                "size": run.font.size.pt if run.font.size else None,
                "bold": bool(run.font.bold),
            })
    return runs


def check_overflow(pptx_path: str, label: str) -> List[str]:
    """Check if text elements overflow their container width."""
    prs = Presentation(pptx_path)
    issues = []
    slide_w = prs.slide_width / EMU_PER_IN

    for si, slide in enumerate(prs.slides):
        shapes = list(slide.shapes)
        for i, shape in enumerate(shapes):
            if _is_nav_dot(shape):
                continue
            if not hasattr(shape, "text_frame") or not shape.text:
                continue
            text = shape.text.strip()
            if len(text) < 5:
                continue

            w_in = shape.width / EMU_PER_IN
            h_in = shape.height / EMU_PER_IN
            x_in = shape.left / EMU_PER_IN

            # Check: text element extends beyond slide right edge
            right_edge = x_in + w_in
            if right_edge > slide_w + 0.05:
                issues.append(
                    f"  [{label}] Slide {si+1}: element '{text[:40]}' "
                    f"extends beyond slide right edge "
                    f"(x={x_in:.2f}\", w={w_in:.2f}\", right={right_edge:.2f}\", slide_w={slide_w:.2f}\")"
                )

            # Check: text height severely underestimated
            # Estimate minimum height needed for the text
            runs = _text_runs_info(shape)
            if runs:
                max_fs = max(r["size"] or 12 for r in runs)
                lines_est = max(1, len(text) / max(w_in * 6 / (max_fs / 12), 1))
                min_h = lines_est * max_fs * 1.3 / 72
                if h_in < min_h * 0.5 and len(text) > 30:
                    issues.append(
                        f"  [{label}] Slide {si+1}: text '{text[:30]}...' "
                        f"height {h_in:.3f}\" << estimated min {min_h:.3f}\" "
                        f"(likely truncated)"
                    )

    return issues


# ─── 3. Overlap check ──────────────────────────────────────────────────

def check_overlap(pptx_path: str, label: str) -> List[str]:
    """Check if non-background elements overlap each other excessively."""
    prs = Presentation(pptx_path)
    issues = []

    for si, slide in enumerate(prs.slides):
        shapes = [s for s in slide.shapes if not _is_nav_dot(s)]
        # Only check text elements against each other
        text_shapes = []
        for s in shapes:
            if hasattr(s, "text_frame") and s.text.strip():
                text_shapes.append({
                    "text": s.text.strip()[:40],
                    "x": s.left / EMU_PER_IN,
                    "y": s.top / EMU_PER_IN,
                    "w": s.width / EMU_PER_IN,
                    "h": s.height / EMU_PER_IN,
                })

        # Check pairwise overlap
        for i in range(len(text_shapes)):
            for j in range(i + 1, len(text_shapes)):
                a, b = text_shapes[i], text_shapes[j]
                # Compute overlap
                x_overlap = min(a["x"] + a["w"], b["x"] + b["w"]) - max(a["x"], b["x"])
                y_overlap = min(a["y"] + a["h"], b["y"] + b["h"]) - max(a["y"], b["y"])

                if x_overlap > 0.05 and y_overlap > 0.05:
                    # Significant overlap — check if one isn't a bg shape
                    area_a = a["w"] * a["h"]
                    area_b = b["w"] * b["h"]
                    overlap_area = x_overlap * y_overlap
                    overlap_ratio = overlap_area / min(area_a, area_b) if min(area_a, area_b) > 0 else 0

                    if overlap_ratio > 0.3:
                        issues.append(
                            f"  [{label}] Slide {si+1}: overlap {overlap_ratio:.0%} between "
                            f"'{a['text'][:25]}' and '{b['text'][:25]}'"
                        )

    return issues


# ─── 4. Position drift analysis ────────────────────────────────────────

def check_position_drift(golden_path: str = GOLDEN, sandbox_path: str = SANDBOX) -> List[str]:
    """Analyze systematic position offsets between golden and sandbox."""
    from collections import defaultdict

    try:
        from compare_visual_comprehensive import match_elements_comprehensive
    except ImportError:
        # Inline minimal matching
        pass

    prs_g = Presentation(golden_path)
    prs_s = Presentation(sandbox_path)
    num_slides = min(len(prs_g.slides), len(prs_s.slides))

    issues = []
    all_dx = []
    all_dy = []
    per_slide_drift = []

    for si in range(num_slides):
        g_shapes = [(s, s.text.strip()[:30] if hasattr(s, "text") and s.text.strip() else None)
                    for s in prs_g.slides[si].shapes if not _is_nav_dot(s)]
        s_shapes = [(s, s.text.strip()[:30] if hasattr(s, "text") and s.text.strip() else None)
                    for s in prs_s.slides[si].shapes if not _is_nav_dot(s)]

        dxs, dys = [], []
        # Match by text content
        for gs, gt in g_shapes:
            if not gt:
                continue
            for ss, st in s_shapes:
                if st and gt == st:
                    dx = (ss.left - gs.left) / EMU_PER_IN
                    dy = (ss.top - gs.top) / EMU_PER_IN
                    dxs.append(dx)
                    dys.append(dy)
                    break

        if dxs:
            avg_dx = sum(dxs) / len(dxs)
            avg_dy = sum(dys) / len(dys)
            max_dx = max(abs(d) for d in dxs)
            max_dy = max(abs(d) for d in dys)
            per_slide_drift.append((si + 1, avg_dx, avg_dy, max_dx, max_dy))
            all_dx.extend(dxs)
            all_dy.extend(dys)

    if all_dx:
        issues.append(f"\n  [POSITION DRIFT] Analysis across {num_slides} slides:")
        issues.append(f"    Mean X drift: {sum(all_dx)/len(all_dx):+.3f}\"  |  Mean Y drift: {sum(all_dy)/len(all_dy):+.3f}\"")
        issues.append(f"    Max |X| drift: {max(abs(d) for d in all_dx):.3f}\"  |  Max |Y| drift: {max(abs(d) for d in all_dy):.3f}\"")

        for slide_num, avg_dx, avg_dy, max_dx, max_dy in per_slide_drift:
            flag = ""
            if abs(avg_dy) > 0.15 or abs(avg_dx) > 0.15:
                flag = " ⚠ SYSTEMATIC"
            issues.append(
                f"    Slide {slide_num}: avg(dx={avg_dx:+.3f}\", dy={avg_dy:+.3f}\") "
                f"max(|dx|={max_dx:.3f}\", |dy|={max_dy:.3f}\"){flag}"
            )

    return issues


# ─── 5. Element count gap analysis ──────────────────────────────────────

def check_element_gaps(golden_path: str = GOLDEN, sandbox_path: str = SANDBOX) -> List[str]:
    """Categorize missing and extra elements."""
    prs_g = Presentation(golden_path)
    prs_s = Presentation(sandbox_path)
    num_slides = min(len(prs_g.slides), len(prs_s.slides))

    issues = []
    for si in range(num_slides):
        g_texts = set()
        s_texts = set()
        g_shapes = 0
        s_shapes = 0

        for s in prs_g.slides[si].shapes:
            if _is_nav_dot(s):
                continue
            if hasattr(s, "text") and s.text.strip():
                g_texts.add(s.text.strip()[:50])
            else:
                g_shapes += 1

        for s in prs_s.slides[si].shapes:
            if _is_nav_dot(s):
                continue
            if hasattr(s, "text") and s.text.strip():
                s_texts.add(s.text.strip()[:50])
            else:
                s_shapes += 1

        missing_texts = g_texts - s_texts
        extra_texts = s_texts - g_texts

        if missing_texts:
            for t in sorted(missing_texts)[:5]:
                issues.append(f"  [MISSING] Slide {si+1}: text '{t[:50]}'")
        if extra_texts:
            for t in sorted(extra_texts)[:5]:
                issues.append(f"  [EXTRA]   Slide {si+1}: text '{t[:50]}'")

        shape_diff = s_shapes - g_shapes
        if abs(shape_diff) > 1:
            direction = "extra" if shape_diff > 0 else "missing"
            issues.append(f"  [SHAPES]  Slide {si+1}: {abs(shape_diff)} {direction} shapes (golden={g_shapes}, sandbox={s_shapes})")

    return issues


# ─── 6. Card containment analysis ───────────────────────────────────────

def _shape_box(shape) -> Dict[str, float]:
    return {
        "x": shape.left / EMU_PER_IN,
        "y": shape.top / EMU_PER_IN,
        "w": shape.width / EMU_PER_IN,
        "h": shape.height / EMU_PER_IN,
    }


def _text_box(shape) -> Dict[str, float]:
    box = _shape_box(shape)
    box["text"] = shape.text.strip()
    return box


def _is_card_like_shape(shape) -> bool:
    if hasattr(shape, "text") and shape.text.strip():
        return False
    box = _shape_box(shape)
    return 2.0 <= box["w"] <= 12.0 and 0.6 <= box["h"] <= 3.0


def _find_containing_card(text_box: Dict[str, float], cards: List[Dict[str, float]]) -> Optional[Dict[str, float]]:
    best = None
    best_overlap = 0.0
    tx1, tx2 = text_box["x"], text_box["x"] + text_box["w"]
    for card in cards:
        cx1, cx2 = card["x"], card["x"] + card["w"]
        overlap = min(tx2, cx2) - max(tx1, cx1)
        if overlap <= 0:
            continue
        overlap_ratio = overlap / max(text_box["w"], 0.01)
        if overlap_ratio < 0.6:
            continue
        if text_box["y"] < card["y"] - 0.15:
            continue
        if best is None or overlap_ratio > best_overlap:
            best = card
            best_overlap = overlap_ratio
    return best


def check_card_containment(golden_path: str = GOLDEN, sandbox_path: str = SANDBOX) -> List[str]:
    """Detect text that should sit inside a card but is stacked below it."""
    prs_g = Presentation(golden_path)
    prs_s = Presentation(sandbox_path)
    num_slides = min(len(prs_g.slides), len(prs_s.slides))

    issues = []
    for si in range(num_slides):
        golden_cards = [_shape_box(s) for s in prs_g.slides[si].shapes if _is_card_like_shape(s)]
        sandbox_cards = [_shape_box(s) for s in prs_s.slides[si].shapes if _is_card_like_shape(s)]
        if not golden_cards or not sandbox_cards:
            continue

        golden_texts = {}
        for s in prs_g.slides[si].shapes:
            if hasattr(s, "text") and s.text.strip():
                golden_texts[s.text.strip()[:60]] = _text_box(s)

        for s in prs_s.slides[si].shapes:
            if not (hasattr(s, "text") and s.text.strip()):
                continue
            key = s.text.strip()[:60]
            if key not in golden_texts:
                continue
            sandbox_text = _text_box(s)
            golden_text = golden_texts[key]
            golden_card = _find_containing_card(golden_text, golden_cards)
            sandbox_card = _find_containing_card(sandbox_text, sandbox_cards)
            if not golden_card or not sandbox_card:
                continue

            golden_inside = golden_text["y"] <= golden_card["y"] + golden_card["h"] + 0.1
            sandbox_outside = sandbox_text["y"] > sandbox_card["y"] + sandbox_card["h"] + 0.05
            if golden_inside and sandbox_outside:
                issues.append(
                    f"  [CARD] Slide {si+1}: text '{key[:40]}' is below its card in sandbox "
                    f"(text_y={sandbox_text['y']:.2f}\", card_bottom={sandbox_card['y'] + sandbox_card['h']:.2f}\")"
                )

    return issues


# ─── 7. Color diff analysis ─────────────────────────────────────────────

def check_color_diffs(golden_path: str = GOLDEN, sandbox_path: str = SANDBOX) -> List[str]:
    """Find color mismatches between golden and sandbox."""
    prs_g = Presentation(golden_path)
    prs_s = Presentation(sandbox_path)
    num_slides = min(len(prs_g.slides), len(prs_s.slides))

    issues = []
    for si in range(num_slides):
        g_map = {}
        s_map = {}

        for s in prs_g.slides[si].shapes:
            if _is_nav_dot(s):
                continue
            if hasattr(s, "text") and s.text.strip():
                g_map[s.text.strip()[:40]] = _get_fill_color(s)

        for s in prs_s.slides[si].shapes:
            if _is_nav_dot(s):
                continue
            if hasattr(s, "text") and s.text.strip():
                s_map[s.text.strip()[:40]] = _get_fill_color(s)

        for key in g_map:
            if key in s_map and g_map[key] and s_map[key] and g_map[key] != s_map[key]:
                issues.append(
                    f"  [COLOR] Slide {si+1}: '{key[:40]}' golden={g_map[key]} sandbox={s_map[key]}"
                )

    return issues


def collect_eval_summary(
    golden_path: str = GOLDEN,
    sandbox_path: str = SANDBOX,
    include_visual: bool = True,
) -> Dict[str, object]:
    """Return a structured evaluation summary for a golden/sandbox PPTX pair."""
    visual_summary = None
    if include_visual:
        visual_summary = run_visual_comprehensive(golden_path, sandbox_path)

    golden_overflows = check_overflow(golden_path, "GOLDEN")
    sandbox_overflows = check_overflow(sandbox_path, "SANDBOX")
    golden_overlaps = check_overlap(golden_path, "GOLDEN")
    sandbox_overlaps = check_overlap(sandbox_path, "SANDBOX")
    drift_issues = check_position_drift(golden_path, sandbox_path)
    gap_issues = check_element_gaps(golden_path, sandbox_path)
    card_issues = check_card_containment(golden_path, sandbox_path)
    color_issues = check_color_diffs(golden_path, sandbox_path)
    total_issues = (
        len(sandbox_overflows)
        + len(sandbox_overlaps)
        + len(gap_issues)
        + len(card_issues)
        + len(color_issues)
    )

    return {
        "golden_path": golden_path,
        "sandbox_path": sandbox_path,
        "visual_summary": visual_summary,
        "golden_overflow_issues": golden_overflows,
        "sandbox_overflow_issues": sandbox_overflows,
        "golden_overlap_issues": golden_overlaps,
        "sandbox_overlap_issues": sandbox_overlaps,
        "position_drift_issues": drift_issues,
        "element_gap_issues": gap_issues,
        "card_containment_issues": card_issues,
        "color_diff_issues": color_issues,
        "golden_overflow_count": len(golden_overflows),
        "sandbox_overflow_count": len(sandbox_overflows),
        "golden_overlap_count": len(golden_overlaps),
        "sandbox_overlap_count": len(sandbox_overlaps),
        "element_gap_count": len(gap_issues),
        "card_containment_count": len(card_issues),
        "color_diff_count": len(color_issues),
        "total_actionable": total_issues,
    }


# ─── Main ──────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="Rigorous HTML→PPTX export evaluation")
    parser.add_argument("--golden", default=GOLDEN, help="Golden PPTX path")
    parser.add_argument("--sandbox", default=SANDBOX, help="Sandbox PPTX path")
    parser.add_argument("--skip-visual", action="store_true", help="Skip compare-visual-comprehensive wrapper")
    args = parser.parse_args()

    summary = collect_eval_summary(
        golden_path=args.golden,
        sandbox_path=args.sandbox,
        include_visual=not args.skip_visual,
    )

    print("=" * 80)
    print("RIGOROUS MULTI-DIMENSIONAL EVALUATION")
    print(f"  Golden:  {summary['golden_path']}")
    print(f"  Sandbox: {summary['sandbox_path']}")
    print("=" * 80)

    # 1. Visual comprehensive
    print("\n" + "=" * 80)
    print("DIMENSION 1: Visual Comprehensive (compare-visual-comprehensive.py)")
    print("=" * 80)
    print(summary["visual_summary"] or "  SKIP: visual comprehensive disabled")

    # 2. Overflow check
    print("\n" + "=" * 80)
    print("DIMENSION 2: Text Overflow Check")
    print("=" * 80)
    golden_overflows = summary["golden_overflow_issues"]
    sandbox_overflows = summary["sandbox_overflow_issues"]
    all_overflows = golden_overflows + sandbox_overflows
    if all_overflows:
        for issue in all_overflows:
            print(issue)
    else:
        print("  ✓ No overflow issues detected")
    print(f"\n  Summary: {len(sandbox_overflows)} sandbox overflow(s), {len(golden_overflows)} golden overflow(s)")

    # 3. Overlap check
    print("\n" + "=" * 80)
    print("DIMENSION 3: Element Overlap Check")
    print("=" * 80)
    sandbox_overlaps = summary["sandbox_overlap_issues"]
    golden_overlaps = summary["golden_overlap_issues"]
    all_overlaps = golden_overlaps + sandbox_overlaps
    if all_overlaps:
        for issue in all_overlaps:
            print(issue)
    else:
        print("  ✓ No overlap issues detected")
    print(f"\n  Summary: {len(sandbox_overlaps)} sandbox overlap(s), {len(golden_overlaps)} golden overlap(s)")

    # 4. Position drift
    print("\n" + "=" * 80)
    print("DIMENSION 4: Position Drift Analysis")
    print("=" * 80)
    drift_issues = summary["position_drift_issues"]
    for issue in drift_issues:
        print(issue)

    # 5. Element gaps
    print("\n" + "=" * 80)
    print("DIMENSION 5: Element Gap Analysis")
    print("=" * 80)
    gap_issues = summary["element_gap_issues"]
    if gap_issues:
        for issue in gap_issues:
            print(issue)
    else:
        print("  ✓ No element gaps detected")

    # 6. Card containment
    print("\n" + "=" * 80)
    print("DIMENSION 6: Card Containment Analysis")
    print("=" * 80)
    card_issues = summary["card_containment_issues"]
    if card_issues:
        for issue in card_issues:
            print(issue)
    else:
        print("  ✓ No card containment issues detected")

    # 7. Color diffs
    print("\n" + "=" * 80)
    print("DIMENSION 7: Color Difference Analysis")
    print("=" * 80)
    color_issues = summary["color_diff_issues"]
    if color_issues:
        for issue in color_issues:
            print(issue)
    else:
        print("  ✓ No color differences detected")

    # Final summary
    print("\n" + "=" * 80)
    print("FINAL SUMMARY")
    print("=" * 80)
    print(f"  Overflow issues:     {summary['sandbox_overflow_count']}")
    print(f"  Overlap issues:      {summary['sandbox_overlap_count']}")
    print(f"  Position drift:      see above")
    print(f"  Element gaps:        {summary['element_gap_count']}")
    print(f"  Card containment:    {summary['card_containment_count']}")
    print(f"  Color differences:   {summary['color_diff_count']}")
    print(f"  Total actionable:    {summary['total_actionable']}")
    print("=" * 80)

    return 0


if __name__ == "__main__":
    sys.exit(main())
