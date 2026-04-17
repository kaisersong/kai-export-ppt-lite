#!/usr/bin/env python3
"""
逐页详细对比 sandbox 版本和 kai-html-export golden 版本的差异。
"""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE

golden_path = Path("/tmp/kai-html-export-golden.pptx")
sandbox_path = Path("/tmp/kai-test-v3.pptx")

golden_prs = Presentation(str(golden_path))
sandbox_prs = Presentation(str(sandbox_path))

num_slides = min(len(golden_prs.slides), len(sandbox_prs.slides))

def extract_elements(slide, label):
    """提取幻灯片中所有元素信息"""
    elems = []
    for sh in slide.shapes:
        try:
            x = sh.left.inches
            y = sh.top.inches
            w = sh.width.inches
            h = sh.height.inches
        except:
            continue
        txt = (sh.text or '').strip()[:40] if sh.has_text_frame else ''
        font_size = 0
        font_color = None
        is_bold = False
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        font_size = max(font_size, run.font.size.pt)
                    if run.font.bold:
                        is_bold = True
                    try:
                        if run.font.color and run.font.color.rgb is not None:
                            font_color = f'#{run.font.color.rgb}'
                    except AttributeError:
                        pass
        fill_color = None
        try:
            if sh.fill.type == MSO_FILL.SOLID:
                fill_color = f'#{sh.fill.fore_color.rgb}'
        except:
            pass
        elems.append({
            'x': x, 'y': y, 'w': w, 'h': h,
            'text': txt, 'font_size': font_size, 'font_color': font_color,
            'is_bold': is_bold, 'fill': fill_color
        })
    return elems

def match_elements(golden_elems, sandbox_elems):
    """通过文本内容匹配 golden 和 sandbox 的元素"""
    matched = []
    unmatched_g = []
    unmatched_s = []
    used_s = set()

    for ge in golden_elems:
        if not ge['text']:
            continue
        best_match = None
        best_score = 0
        for si, se in enumerate(sandbox_elems):
            if si in used_s:
                continue
            if not se['text']:
                continue
            # Simple text overlap score
            g_words = set(ge['text'][:20])
            s_words = set(se['text'][:20])
            if g_words & s_words:
                score = len(g_words & s_words) / max(len(g_words | s_words), 1)
            else:
                score = 0
            if score > best_score:
                best_score = score
                best_match = (se, si)
        if best_match and best_score > 0.3:
            matched.append((ge, best_match[0], best_score))
            used_s.add(best_match[1])
        else:
            unmatched_g.append(ge)

    for si, se in enumerate(sandbox_elems):
        if si not in used_s and se['text']:
            unmatched_s.append(se)

    return matched, unmatched_g, unmatched_s

print("=" * 80)
print("逐页详细对比: kai-html-export (golden) vs kai-export-ppt-lite (sandbox)")
print("=" * 80)

total_issues = 0
total_warnings = 0

for si in range(num_slides):
    g_slide = golden_prs.slides[si]
    s_slide = sandbox_prs.slides[si]

    g_elems = extract_elements(g_slide, f"Slide {si+1} Golden")
    s_elems = extract_elements(s_slide, f"Slide {si+1} Sandbox")

    g_text = [e for e in g_elems if e['text']]
    s_text = [e for e in s_elems if e['text']]

    matched, unmatched_g, unmatched_s = match_elements(g_elems, s_elems)

    print(f"\n{'─' * 76}")
    print(f"第 {si+1} 页")
    print(f"  Golden: {len(g_text)} 个文字元素, {len(g_elems)} 个总元素")
    print(f"  Sandbox: {len(s_text)} 个文字元素, {len(s_elems)} 个总元素")
    print(f"  匹配: {len(matched)} 个, Golden独有: {len(unmatched_g)} 个, Sandbox独有: {len(unmatched_s)} 个")

    # Detailed comparison
    issues = []

    # Check matched elements for position/size/font differences
    for ge, se, score in matched:
        dx = abs(ge['x'] - se['x'])
        dy = abs(ge['y'] - se['y'])
        dw = abs(ge['w'] - se['w'])
        dh = abs(ge['h'] - se['h'])
        font_diff = abs(ge['font_size'] - se['font_size'])

        if dx > 0.5 or dy > 0.5:
            issues.append(f"    ⚠ 位置偏移 '{ge['text'][:20]}': dx={dx:.2f}\" dy={dy:.2f}\" "
                        f"(G:({ge['x']:.2f},{ge['y']:.2f}) S:({se['x']:.2f},{se['y']:.2f}))")

        if font_diff > 3:
            issues.append(f"    ⚠ 字号差异 '{ge['text'][:20]}': 差{font_diff:.1f}pt "
                        f"(G:{ge['font_size']:.1f} S:{se['font_size']:.1f})")

        if ge['font_color'] and se['font_color'] and ge['font_color'] != se['font_color']:
            issues.append(f"    ⚠ 颜色差异 '{ge['text'][:20]}': G:{ge['font_color']} S:{se['font_color']}")

    if issues:
        total_issues += len(issues)
        print(f"  问题 ({len(issues)} 个):")
        for iss in issues[:8]:  # Max 8 per slide
            print(iss)
        if len(issues) > 8:
            print(f"    ... 还有 {len(issues) - 8} 个")
    else:
        print(f"  ✅ 无明显差异")

    # Check unmatched golden elements
    if unmatched_g:
        for u in unmatched_g[:3]:
            print(f"  ❌ Golden有但Sandbox缺失: '{u['text'][:30]}' at ({u['x']:.2f},{u['y']:.2f})")
        total_issues += len(unmatched_g)

    # Check unmatched sandbox elements (excluding nav dots and counters)
    extra = []
    for u in unmatched_s:
        if u['w'] > 0.15 or u['h'] > 0.15 or u['text']:
            extra.append(u)
    if extra:
        for u in extra[:3]:
            print(f"  ℹ️  Sandbox额外: '{u['text'][:30]}' at ({u['x']:.2f},{u['y']:.2f})")
        total_warnings += len(extra)

print(f"\n{'─' * 76}")
print(f"\n总计: {total_issues} 个问题, {total_warnings} 个额外元素")

# Score calculation
total_possible = num_slides * 10
text_coverage = sum(1 for si in range(num_slides) for _ in golden_prs.slides[si].shapes if _.has_text_frame and (_.text or '').strip())
sandbox_coverage = sum(1 for si in range(num_slides) for _ in sandbox_prs.slides[si].shapes if _.has_text_frame and (_.text or '').strip())

# Weighted score: coverage (40%), position accuracy (40%), no extra elements (20%)
coverage_score = min(sandbox_coverage, text_coverage) / text_coverage * 4.0
position_score = max(0, 4.0 - total_issues * 0.01)  # Each issue costs 0.01
extra_score = max(0, 2.0 - total_warnings * 0.1)
score = coverage_score + position_score + extra_score
score = min(10, max(0, score))

print(f"\n估算评分: {score:.1f}/10")
print(f"  覆盖率: {coverage_score:.1f}/4.0 (Golden文字: {text_coverage}, Sandbox文字: {sandbox_coverage})")
print(f"  位置分: {position_score:.1f}/4.0 (问题数: {total_issues})")
print(f"  额外分: {extra_score:.1f}/2.0 (额外元素: {total_warnings})")
