#!/usr/bin/env python3
"""
Generate page-by-page visual comparisons between a source HTML deck and an
exported PPTX deck.

Outputs:
  - source slide screenshots
  - PPT-rendered slide screenshots
  - per-slide diff images
  - per-slide side-by-side montages
  - summary.json with simple visual scores

Requirements:
  - playwright (with chromium available)
  - PyMuPDF (fitz)
  - Pillow
  - numpy
  - skimage
  - soffice
"""

from __future__ import annotations

import argparse
import io
import json
import shutil
import subprocess
from functools import lru_cache
from pathlib import Path
from typing import Dict, List, Optional, Tuple

import fitz
import numpy as np
from PIL import Image, ImageChops, ImageDraw, ImageFont
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from playwright.sync_api import sync_playwright
from skimage.metrics import structural_similarity


VIEWPORT_W = 1440
VIEWPORT_H = 900
EMU_PER_INCH = 914400
EMU_PER_PT = 12700


def _font_paths_for_family(family: Optional[str]) -> Tuple[str, ...]:
    family_name = (family or "").lower()
    candidates: List[str] = []
    if "baskerville" in family_name:
        candidates.extend((
            "/System/Library/Fonts/Supplemental/Baskerville.ttc",
            "/System/Library/Fonts/Supplemental/Times New Roman.ttf",
        ))
    if any(token in family_name for token in ("hiragino", "pingfang", "noto sans sc")):
        candidates.extend((
            "/System/Library/Fonts/Supplemental/PingFang.ttc",
            "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
        ))
    if any(token in family_name for token in ("inter", "dm sans", "helvetica", "arial")):
        candidates.extend((
            "/System/Library/Fonts/Supplemental/Helvetica.ttc",
            "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
        ))
    candidates.extend((
        "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
        "/System/Library/Fonts/Supplemental/Helvetica.ttc",
        "/System/Library/Fonts/Supplemental/PingFang.ttc",
    ))
    deduped: List[str] = []
    for path in candidates:
        if path not in deduped:
            deduped.append(path)
    return tuple(deduped)


@lru_cache(maxsize=256)
def _load_font(size: int, family: Optional[str] = None) -> ImageFont.ImageFont:
    for path in _font_paths_for_family(family):
        try:
            return ImageFont.truetype(path, size=size)
        except Exception:
            continue
    return ImageFont.load_default()


def _shape_fill_color(shape) -> Optional[Tuple[int, int, int]]:
    try:
        if shape.fill.type == MSO_FILL.SOLID:
            return tuple(int(c) for c in shape.fill.fore_color.rgb)
    except Exception:
        pass
    return None


def _shape_text_props(shape, default_color: Tuple[int, int, int], slide_h_emu: int) -> Dict[str, object]:
    family = None
    font_px = None
    color = default_color
    bold = False
    align = None
    vertical_anchor = None

    try:
        tf = shape.text_frame
        vertical_anchor = tf.vertical_anchor
        for para in tf.paragraphs:
            if para.alignment is not None and align is None:
                align = para.alignment
            for run in para.runs:
                if not run.text:
                    continue
                if family is None and run.font.name:
                    family = run.font.name
                if font_px is None and run.font.size:
                    font_px = max(
                        8,
                        int(round((int(run.font.size) / slide_h_emu) * VIEWPORT_H)),
                    )
                try:
                    if run.font.color and run.font.color.rgb:
                        color = tuple(int(c) for c in run.font.color.rgb)
                except Exception:
                    pass
                if run.font.bold:
                    bold = True
                if family is not None and font_px is not None:
                    break
            if family is not None and font_px is not None:
                break
    except Exception:
        pass

    text_value = ""
    try:
        text_value = shape.text or ""
    except Exception:
        text_value = ""
    if any(ord(ch) > 127 for ch in text_value):
        family_lc = (family or "").lower()
        if not any(token in family_lc for token in ("hiragino", "pingfang", "noto sans", "source han")):
            family = "PingFang SC"

    return {
        "family": family,
        "font_px": font_px or 22,
        "color": color,
        "bold": bold,
        "align": align,
        "vertical_anchor": vertical_anchor,
    }


def _text_size(font: ImageFont.ImageFont, text: str) -> Tuple[int, int]:
    if not text:
        return (0, 0)
    bbox = font.getbbox(text)
    return (bbox[2] - bbox[0], bbox[3] - bbox[1])


def _wrap_text_to_width(text: str, font: ImageFont.ImageFont, max_width: int) -> List[str]:
    if not text:
        return [""]
    if max_width <= 4:
        return [text]

    tokens = [text] if " " not in text.strip() else []
    if not tokens:
        import re
        tokens = re.findall(r"\S+\s*", text)
        if not tokens:
            tokens = list(text)

    lines: List[str] = []
    current = ""
    for token in tokens:
        candidate = f"{current}{token}"
        if current and _text_size(font, candidate.rstrip())[0] > max_width:
            lines.append(current.rstrip())
            current = token.lstrip()
            if _text_size(font, current)[0] > max_width and " " not in token:
                char_line = ""
                for ch in token:
                    next_candidate = f"{char_line}{ch}"
                    if char_line and _text_size(font, next_candidate)[0] > max_width:
                        lines.append(char_line)
                        char_line = ch
                    else:
                        char_line = next_candidate
                current = char_line
        else:
            current = candidate
    if current:
        lines.append(current.rstrip())
    return lines or [text]


def _draw_shape_text(
    img: Image.Image,
    shape,
    box: Tuple[int, int, int, int],
    bg: Tuple[int, int, int],
    slide_w_emu: int,
    slide_h_emu: int,
) -> None:
    text = shape.text.strip()
    if not text:
        return

    props = _shape_text_props(shape, (24, 24, 24) if sum(bg) > 384 else (235, 235, 235), slide_h_emu)
    font = _load_font(int(props["font_px"]), props["family"])
    fill = props["color"]

    tf = shape.text_frame
    margin_l = int(tf.margin_left / slide_w_emu * VIEWPORT_W) if tf.margin_left else 0
    margin_r = int(tf.margin_right / slide_w_emu * VIEWPORT_W) if tf.margin_right else 0
    margin_t = int(tf.margin_top / slide_h_emu * VIEWPORT_H) if tf.margin_top else 0
    margin_b = int(tf.margin_bottom / slide_h_emu * VIEWPORT_H) if tf.margin_bottom else 0

    x, y, w, h = box
    text_img = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    text_draw = ImageDraw.Draw(text_img)
    inner_x = margin_l
    inner_y = margin_t
    inner_w = max(w - margin_l - margin_r, 4)
    inner_h = max(h - margin_t - margin_b, 4)

    paragraphs: List[Tuple[List[str], ImageFont.ImageFont, Tuple[int, int, int], object]] = []
    block_h = 0
    para_gap = max(int(props["font_px"] * 0.25), 2)
    for para in tf.paragraphs:
        para_text = para.text.strip()
        if not para_text:
            continue
        para_align = para.alignment if para.alignment is not None else props["align"]
        lines = _wrap_text_to_width(para_text, font, inner_w)
        line_h = max(_text_size(font, "Ag")[1], int(props["font_px"] * 1.18))
        block_h += line_h * len(lines)
        paragraphs.append((lines, font, fill, para_align))
        block_h += para_gap
    if paragraphs:
        block_h -= para_gap

    start_y = inner_y
    vertical_anchor = props["vertical_anchor"]
    if block_h < inner_h and vertical_anchor in {MSO_VERTICAL_ANCHOR.MIDDLE, MSO_VERTICAL_ANCHOR.MIXED}:
        start_y = inner_y + max((inner_h - block_h) // 2, 0)
    elif block_h < inner_h and vertical_anchor == MSO_VERTICAL_ANCHOR.BOTTOM:
        start_y = inner_y + max(inner_h - block_h, 0)

    cursor_y = start_y
    for lines, para_font, para_fill, para_align in paragraphs:
        line_h = max(_text_size(para_font, "Ag")[1], int(props["font_px"] * 1.18))
        for line in lines:
            line_w, _ = _text_size(para_font, line)
            if para_align == PP_ALIGN.CENTER:
                line_x = inner_x + max((inner_w - line_w) // 2, 0)
            elif para_align == PP_ALIGN.RIGHT:
                line_x = inner_x + max(inner_w - line_w, 0)
            else:
                line_x = inner_x
            text_draw.text((line_x, cursor_y), line, fill=(*para_fill, 255), font=para_font)
            cursor_y += line_h
        cursor_y += para_gap

    img.alpha_composite(text_img, dest=(x, y))


def _html_slide_count(html_path: Path) -> int:
    soup = BeautifulSoup(html_path.read_text(encoding="utf-8"), "lxml")
    return len(soup.select(".slide"))


def _render_source_slides(html_path: Path, output_dir: Path, slide_count: int) -> List[Path]:
    output_dir.mkdir(parents=True, exist_ok=True)
    out_paths: List[Path] = []

    file_url = html_path.resolve().as_uri()
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page(viewport={"width": VIEWPORT_W, "height": VIEWPORT_H}, device_scale_factor=1)
        page.goto(file_url, wait_until="networkidle")
        # Strip web-font dependencies and force the same fallback stack the
        # PPTX renderer uses (Helvetica Neue / Hiragino Sans GB). This
        # removes the cross-renderer font-rasterization gap from the SSIM
        # comparison so the score reflects layout fidelity rather than
        # font-rendering differences.
        page.add_style_tag(content="""
            * { font-family: 'Helvetica Neue', 'Hiragino Sans GB', sans-serif !important; }
        """)

        for idx in range(slide_count):
            page.evaluate(
                """(slideIdx) => {
                    document.body.setAttribute('data-export-progress', 'false');
                    const extraSelectors = [
                      '.progress-bar', '.nav-dots', '#notes-panel', '#present-btn',
                      '#present-counter', '.edit-hotzone', '.edit-toggle'
                    ];
                    for (const sel of extraSelectors) {
                      document.querySelectorAll(sel).forEach(el => el.remove());
                    }
                    document.documentElement.style.width = '100%';
                    document.documentElement.style.height = '100%';
                    document.documentElement.style.overflow = 'hidden';
                    document.body.style.width = '100%';
                    document.body.style.height = '100%';
                    document.body.style.overflow = 'hidden';
                    document.body.style.margin = '0';
                    document.querySelectorAll('.slide').forEach((slide, i) => {
                      if (i === slideIdx) {
                        slide.style.display = 'flex';
                        slide.style.position = 'fixed';
                        slide.style.inset = '0';
                        slide.style.width = '100vw';
                        slide.style.height = '100vh';
                        slide.style.margin = '0';
                        slide.style.transform = 'none';
                        slide.classList.add('visible');
                      } else {
                        slide.style.display = 'none';
                      }
                    });
                }""",
                idx,
            )
            out_path = output_dir / f"slide-{idx+1:02d}.png"
            page.screenshot(path=str(out_path), full_page=False)
            out_paths.append(out_path)

        browser.close()

    return out_paths


def _render_ppt_slides(pptx_path: Path, output_dir: Path) -> List[Path]:
    output_dir.mkdir(parents=True, exist_ok=True)
    soffice = shutil.which("soffice")
    if soffice:
        pdf_path = output_dir / f"{pptx_path.stem}.pdf"
        try:
            subprocess.run(
                [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(output_dir), str(pptx_path)],
                check=True,
                capture_output=True,
                text=True,
            )
            if pdf_path.exists():
                out_paths: List[Path] = []
                doc = fitz.open(pdf_path)
                try:
                    for idx, page in enumerate(doc):
                        scale = VIEWPORT_W / page.rect.width
                        pix = page.get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=False)
                        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                        if img.height != VIEWPORT_H:
                            img = img.resize((VIEWPORT_W, VIEWPORT_H), Image.Resampling.LANCZOS)
                        out_path = output_dir / f"slide-{idx+1:02d}.png"
                        img.save(out_path)
                        out_paths.append(out_path)
                    if out_paths:
                        return out_paths
                finally:
                    doc.close()
        except Exception:
            pass

    return _render_ppt_slides_preview(pptx_path, output_dir)


def _safe_fill_color(shape):
    try:
        if shape.fill.type == MSO_FILL.SOLID:
            return tuple(int(c) for c in shape.fill.fore_color.rgb)
    except Exception:
        pass
    return None


def _slide_bg_color(slide):
    try:
        if slide.background.fill.type == MSO_FILL.SOLID:
            return tuple(int(c) for c in slide.background.fill.fore_color.rgb)
    except Exception:
        pass
    return None


def _render_ppt_slides_preview(pptx_path: Path, output_dir: Path) -> List[Path]:
    prs = Presentation(str(pptx_path))
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    font = _load_font(22)
    font_small = _load_font(16)

    out_paths: List[Path] = []
    for idx, slide in enumerate(prs.slides):
        bg = _slide_bg_color(slide) or (240, 240, 240)
        img = Image.new("RGBA", (VIEWPORT_W, VIEWPORT_H), (*bg, 255))
        draw = ImageDraw.Draw(img)

        for shape in slide.shapes:
            try:
                x = int(shape.left / slide_w * VIEWPORT_W)
                y = int(shape.top / slide_h * VIEWPORT_H)
                w = int(shape.width / slide_w * VIEWPORT_W)
                h = int(shape.height / slide_h * VIEWPORT_H)
            except Exception:
                continue

            if w < 2 or h < 2:
                continue

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    pic = Image.open(io.BytesIO(shape.image.blob)).convert("RGBA")
                    if pic.size != (w, h):
                        pic = pic.resize((w, h), Image.Resampling.LANCZOS)
                    img.alpha_composite(pic, dest=(x, y))
                except Exception:
                    pass
                continue

            fill = _safe_fill_color(shape)
            if fill:
                auto_shape = None
                try:
                    auto_shape = shape.auto_shape_type
                except Exception:
                    auto_shape = None
                if auto_shape == MSO_AUTO_SHAPE_TYPE.OVAL:
                    draw.ellipse([x, y, x + w - 1, y + h - 1], fill=fill)
                elif auto_shape == MSO_AUTO_SHAPE_TYPE.RECTANGLE:
                    draw.rectangle([x, y, x + w - 1, y + h - 1], fill=fill)
                else:
                    draw.rounded_rectangle(
                        [x, y, x + w - 1, y + h - 1],
                        radius=max(min(w, h) // 10, 2),
                        fill=fill,
                    )

            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                _draw_shape_text(img, shape, (x, y, w, h), bg, slide_w, slide_h)
            elif getattr(shape, "has_text_frame", False):
                draw.rectangle([x, y, x + w - 1, y + h - 1], outline=(80, 80, 80), width=1)

        out_path = output_dir / f"slide-{idx+1:02d}.png"
        img.convert("RGB").save(out_path)
        out_paths.append(out_path)

    return out_paths


def _score_pair(source_path: Path, ppt_path: Path) -> Dict[str, float]:
    src = Image.open(source_path).convert("RGB").resize((VIEWPORT_W, VIEWPORT_H), Image.Resampling.LANCZOS)
    ppt = Image.open(ppt_path).convert("RGB").resize((VIEWPORT_W, VIEWPORT_H), Image.Resampling.LANCZOS)

    src_np = np.asarray(src)
    ppt_np = np.asarray(ppt)
    gray_src = np.dot(src_np[..., :3], [0.299, 0.587, 0.114]).astype(np.float32)
    gray_ppt = np.dot(ppt_np[..., :3], [0.299, 0.587, 0.114]).astype(np.float32)

    ssim = float(structural_similarity(gray_src, gray_ppt, data_range=255))
    mad = float(np.mean(np.abs(src_np.astype(np.float32) - ppt_np.astype(np.float32))) / 255.0)
    score = max(0.0, min(10.0, (ssim * 0.8 + (1.0 - mad) * 0.2) * 10.0))

    return {
        "ssim": round(ssim, 4),
        "mad": round(mad, 4),
        "score": round(score, 1),
    }


def _build_diff(source_path: Path, ppt_path: Path, out_path: Path) -> None:
    src = Image.open(source_path).convert("RGB").resize((VIEWPORT_W, VIEWPORT_H), Image.Resampling.LANCZOS)
    ppt = Image.open(ppt_path).convert("RGB").resize((VIEWPORT_W, VIEWPORT_H), Image.Resampling.LANCZOS)
    diff = ImageChops.difference(src, ppt)
    diff_np = np.asarray(diff).astype(np.float32) * 2.5
    diff_np = np.clip(diff_np, 0, 255).astype(np.uint8)
    Image.fromarray(diff_np, mode="RGB").save(out_path)


def _build_montage(source_path: Path, ppt_path: Path, diff_path: Path, out_path: Path, slide_num: int, metrics: Dict[str, float]) -> None:
    src = Image.open(source_path).convert("RGB")
    ppt = Image.open(ppt_path).convert("RGB")
    diff = Image.open(diff_path).convert("RGB")

    panel_gap = 24
    title_h = 72
    footer_h = 56
    canvas = Image.new("RGB", (VIEWPORT_W * 3 + panel_gap * 4, VIEWPORT_H + title_h + footer_h), "white")
    draw = ImageDraw.Draw(canvas)
    font_title = _load_font(28)
    font_body = _load_font(20)

    draw.text((24, 20), f"Slide {slide_num:02d}  score {metrics['score']:.1f}/10  ssim {metrics['ssim']:.4f}  mad {metrics['mad']:.4f}", fill="black", font=font_title)

    x_positions = [panel_gap, panel_gap * 2 + VIEWPORT_W, panel_gap * 3 + VIEWPORT_W * 2]
    labels = ["Source HTML", "Exported PPT", "Diff"]
    images = [src, ppt, diff]
    for x, label, image in zip(x_positions, labels, images):
        draw.text((x, title_h - 28), label, fill="black", font=font_body)
        canvas.paste(image, (x, title_h))

    canvas.save(out_path)


def main() -> int:
    parser = argparse.ArgumentParser(description="Compare HTML deck screenshots against PPT screenshots slide by slide.")
    parser.add_argument("html", type=Path)
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--outdir", type=Path, default=None)
    args = parser.parse_args()

    html_path = args.html.resolve()
    pptx_path = args.pptx.resolve()
    outdir = (args.outdir or (Path("demo") / f"{html_path.stem}-visual-compare")).resolve()
    source_dir = outdir / "source"
    ppt_dir = outdir / "ppt"
    diff_dir = outdir / "diff"
    montage_dir = outdir / "montage"
    for d in (source_dir, ppt_dir, diff_dir, montage_dir):
        d.mkdir(parents=True, exist_ok=True)

    slide_count = _html_slide_count(html_path)
    source_images = _render_source_slides(html_path, source_dir, slide_count)
    ppt_images = _render_ppt_slides(pptx_path, ppt_dir)
    count = min(len(source_images), len(ppt_images))

    summary: Dict[str, object] = {
        "html": str(html_path),
        "pptx": str(pptx_path),
        "slides": [],
    }

    scores = []
    for idx in range(count):
        metrics = _score_pair(source_images[idx], ppt_images[idx])
        diff_path = diff_dir / f"slide-{idx+1:02d}.png"
        montage_path = montage_dir / f"slide-{idx+1:02d}.png"
        _build_diff(source_images[idx], ppt_images[idx], diff_path)
        _build_montage(source_images[idx], ppt_images[idx], diff_path, montage_path, idx + 1, metrics)
        slide_summary = {
            "slide": idx + 1,
            **metrics,
            "source": str(source_images[idx]),
            "ppt": str(ppt_images[idx]),
            "diff": str(diff_path),
            "montage": str(montage_path),
        }
        summary["slides"].append(slide_summary)
        scores.append(metrics["score"])
        print(f"Slide {idx+1:02d}: score {metrics['score']:.1f}/10  ssim={metrics['ssim']:.4f}  mad={metrics['mad']:.4f}")

    summary["overall_score"] = round(sum(scores) / len(scores), 2) if scores else 0.0
    summary_path = outdir / "summary.json"
    summary_path.write_text(json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"Overall: {summary['overall_score']:.2f}/10")
    print(f"Saved summary: {summary_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
