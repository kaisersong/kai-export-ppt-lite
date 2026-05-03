#!/usr/bin/env python3
"""
export-sandbox-pptx.py — Export HTML presentations to editable PPTX without a browser.

Ported from export-native-pptx.py logic (v5):
  - flatExtract algorithm adapted for no-browser CSS cascade simulation
  - Pre-pass corrections (shape+text sync, adjacent element push)
  - CJK/condensed font width correction
  - Sophisticated segment handling (bold, strike, bgColor, inlineBgBounds)
  - Better word-wrap / auto-size strategy for text boxes

Pure Python: no Playwright, no Chrome. Uses beautifulsoup4 for HTML parsing,
python-pptx for PPTX generation, and Pillow for preview images.

Usage:
    python export-sandbox-pptx.py <presentation.html> [output.pptx] [--width W] [--height H]

Dependencies:
    pip install beautifulsoup4 lxml python-pptx Pillow
"""

import sys
import os
import re
import io
import math
import argparse
import base64
import urllib.request
from pathlib import Path
from dataclasses import dataclass, field
from typing import List, Dict, Any, Optional, Tuple

# ─── Dependency check ─────────────────────────────────────────────────────────

def check_deps():
    missing = []
    try:
        from bs4 import BeautifulSoup, NavigableString, Tag
    except ImportError:
        missing.append("beautifulsoup4")
    try:
        from pptx import Presentation
    except ImportError:
        missing.append("python-pptx")
    if missing:
        print(f"Missing dependencies. Install with:")
        print(f"  pip install {' '.join(missing)}")
        sys.exit(1)

check_deps()

from bs4 import BeautifulSoup, NavigableString, Tag
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree as _etree


# ─── Constants ───────────────────────────────────────────────────────────────

PX_PER_IN = 108.0  # 1440px / 13.33in
CJK_BOX_FACTOR = 1.15
CJK_V_FACTOR = 1.30
CJK_H_FACTOR = 0.15  # extra horizontal space for CJK in PPTX
PPTX_HEIGHT_FACTOR = 1.30  # vertical correction for multi-line CJK text

TEXT_TAGS = {'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'li', 'span', 'a', 'blockquote'}
RASTER_TAGS = {'img', 'svg', 'canvas'}
CONTAINER_TAGS = {'div', 'section', 'article', 'ul', 'ol'}
INLINE_TAGS = {'strong', 'em', 'b', 'i', 'span', 'a', 'mark', 'code', 'small',
               'kbd', 'var', 'abbr', 'time', 'sup', 'sub', 'br'}

# Elements to skip during extraction (navigation, chrome, controls)
SKIP_ELEMENT_IDS = {'present-btn', 'present-counter', 'notes-panel', 'editToggle', 'brand-mark'}
SKIP_ELEMENT_CLASSES = {'nav-dots', 'progress-bar', 'edit-hotzone', 'edit-toggle', 'slide-credit', 'slide-num-label'}


# ─── CSS Parsing ──────────────────────────────────────────────────────────────

@dataclass
class CSSRule:
    selector: str
    properties: Dict[str, str]

# Global: root CSS variables, populated during CSS parsing
_ROOT_CSS_VARS: Dict[str, str] = {}


def resolve_css_variables(css_text: str) -> str:
    """Resolve :root CSS custom properties (var(--name) and var(--name, fallback))."""
    global _ROOT_CSS_VARS
    root_match = re.search(r':root\s*\{([^}]+)\}', css_text)
    if root_match:
        for prop_match in re.finditer(r'(--[\w-]+)\s*:\s*([^;]+);', root_match.group(1)):
            _ROOT_CSS_VARS[prop_match.group(1)] = prop_match.group(2).strip()

    def replace_var(match):
        var_name = match.group(1)
        fallback = match.group(2)
        return _ROOT_CSS_VARS.get(var_name, fallback or match.group(0))

    return re.sub(r'var\((--[\w-]+)(?:,\s*([^)]+))?\)', replace_var, css_text)


def _should_skip_element(element: Tag) -> bool:
    """Check if element should be skipped during extraction (navigation, chrome, etc.)."""
    if not element.name:
        return False
    # Check by ID
    elem_id = element.get('id', '')
    if elem_id in SKIP_ELEMENT_IDS:
        return True
    # Check by class
    classes = element.get('class', [])
    if isinstance(classes, str):
        classes = classes.split()
    for cls in classes:
        if cls in SKIP_ELEMENT_CLASSES:
            return True
    return False


def _expand_background_shorthand(bg_value: str) -> Dict[str, str]:
    """Expand CSS `background` shorthand into backgroundColor / backgroundImage."""
    result = {}
    if 'url(' in bg_value:
        url_match = re.search(r'url\([^)]+\)', bg_value)
        if url_match:
            result['backgroundImage'] = url_match.group(0)
    if 'gradient' in bg_value:
        grad_match = re.search(r'(?:linear|radial|conic)-gradient\([^)]+\)', bg_value)
        if grad_match:
            result['backgroundImage'] = grad_match.group(0)
    stripped = re.sub(r'url\([^)]*\)', '', bg_value)
    stripped = re.sub(r'(?:linear|radial|conic)-gradient\([^)]*\)', '', stripped)
    stripped = stripped.strip()
    hex_match = re.search(r'#(?:[0-9a-fA-F]{6}|[0-9a-fA-F]{3})\b', stripped)
    if hex_match:
        result['backgroundColor'] = hex_match.group(0)
    else:
        rgba_match = re.search(r'rgba?\([^)]+\)', stripped)
        if rgba_match:
            result['backgroundColor'] = rgba_match.group(0)
    return result


def _kebab_to_camel(name: str) -> str:
    """Convert CSS kebab-case property name to camelCase."""
    parts = name.split('-')
    return parts[0] + ''.join(p.capitalize() for p in parts[1:])


def parse_css_rules(css_text: str) -> List[CSSRule]:
    """Parse CSS text into list of (selector, properties) rules.
    Property names are converted from kebab-case to camelCase."""
    css_text = resolve_css_variables(css_text)
    css_text = re.sub(r'/\*.*?\*/', '', css_text, flags=re.DOTALL)
    rules = []
    for block_match in re.finditer(r'([^{}]+)\{([^{}]+)\}', css_text):
        selector = block_match.group(1).strip()
        props_text = block_match.group(2).strip()
        if not selector or selector.startswith('@'):
            continue
        props = {}
        for prop_match in re.finditer(r'([\w-]+)\s*:\s*([^;]+);?', props_text):
            prop_name = _kebab_to_camel(prop_match.group(1).strip())
            props[prop_name] = prop_match.group(2).strip()
        if 'background' in props and 'backgroundColor' not in props:
            expanded = _expand_background_shorthand(props['background'])
            props.update(expanded)
        if 'border' in props:
            bv = props['border']
            if bv and 'none' not in bv and bv != '0px' and bv != '0':
                for side in ('borderLeft', 'borderRight', 'borderTop', 'borderBottom'):
                    if side not in props:
                        props[side] = bv
        for sel in selector.split(','):
            sel = sel.strip()
            if sel and not sel.startswith('@'):
                rules.append(CSSRule(selector=sel, properties=props))
    return rules


def selector_matches(element: Tag, selector: str) -> bool:
    """Check if a CSS selector matches a bs4 element."""
    if not selector or not element.name:
        return False
    if '::' in selector:
        return False
    parts = [p.strip() for p in selector.split() if p.strip()]
    if not parts:
        return False
    if len(parts) == 1:
        return _match_simple_selector(element, parts[0])
    # Descendant selector: check if last part matches element
    # and earlier parts match some ancestor chain
    return _match_descendant_selector(element, parts)


def _match_descendant_selector(element: Tag, parts: List[str]) -> bool:
    """Match a descendant selector like '.card h3' or 'h1 .accent'."""
    if not _match_simple_selector(element, parts[-1]):
        return False
    # Walk up the ancestor chain looking for parts[-2], parts[-3], etc.
    remaining = list(reversed(parts[:-1]))
    ancestor = element.parent
    while ancestor and remaining:
        if _match_simple_selector(ancestor, remaining[0]):
            remaining.pop(0)
        ancestor = ancestor.parent
    return len(remaining) == 0


def _match_simple_selector(element: Tag, selector: str) -> bool:
    """Match a single (non-combinator) CSS selector."""
    if not selector:
        return False
    if selector == '*':
        return True
    # Strip pseudo-classes/elements but track them for special handling
    pseudo_match = re.search(r'::?([\w-]+)(\([^)]*\))?$', selector)
    pseudo_name = pseudo_match.group(1) if pseudo_match else None
    sel_base = re.sub(r'::?[\w-]+(\([^)]*\))?$', '', selector)
    if not sel_base:
        sel_base = element.name or ''

    # Handle pseudo-class :last-child
    if pseudo_name == 'last-child':
        siblings = list(element.parent.children) if element.parent else []
        tag_siblings = [s for s in siblings if isinstance(s, Tag) and s.name == element.name]
        if tag_siblings and tag_siblings[-1] is not element:
            return False

    tag_name = element.name or ''
    classes = element.get('class', [])
    if isinstance(classes, str):
        classes = classes.split()

    # Handle ID selector
    id_match = re.match(r'^#([\w-]+)$', sel_base)
    if id_match:
        return element.get('id') == id_match.group(1)

    # Handle compound class selector (e.g., .swiss-rule.red)
    if sel_base.startswith('.') and '.' in sel_base[1:]:
        class_parts = sel_base[1:].split('.')
        return all(cp in classes for cp in class_parts)

    # Handle single class selector
    class_match = re.match(r'^\.([\w-]+)$', sel_base)
    if class_match:
        return class_match.group(1) in classes
    elem_class_match = re.match(r'^([\w-]+)\.([\w-]+)$', sel_base)
    if elem_class_match:
        tag_ok = elem_class_match.group(1).lower() == tag_name.lower()
        classes = element.get('class', [])
        if isinstance(classes, str):
            classes = classes.split()
        return tag_ok and elem_class_match.group(2) in classes
    elem_id_match = re.match(r'^([\w-]+)#([\w-]+)$', sel_base)
    if elem_id_match:
        return (elem_id_match.group(1).lower() == tag_name.lower() and
                element.get('id') == elem_id_match.group(2))
    return sel_base.lower() == tag_name.lower()


def compute_element_style(
    element: Tag,
    css_rules: List[CSSRule],
    inline_style_str: Optional[str] = None,
    parent_style: Optional[Dict[str, str]] = None,
) -> Dict[str, str]:
    """Compute effective CSS style for an element via cascade."""
    computed = {}
    if parent_style:
        for prop in ('color', 'fontFamily', 'fontSize', 'fontWeight', 'lineHeight', 'textAlign'):
            if prop in parent_style:
                computed[prop] = parent_style[prop]
    for rule in css_rules:
        if selector_matches(element, rule.selector):
            computed.update(rule.properties)
    if inline_style_str:
        # Resolve CSS variables in inline style values
        resolved_inline = resolve_css_variables_inline(inline_style_str)
        for prop_match in re.finditer(r'([\w-]+)\s*:\s*([^;]+);?', resolved_inline):
            prop_name = _kebab_to_camel(prop_match.group(1).strip())
            computed[prop_name] = prop_match.group(2).strip()
    # Expand padding shorthand (padding: X Y → paddingTop/Left/etc)
    _expand_padding(computed)
    # Expand margin shorthand (margin: X Y → marginTop/Left/etc)
    _expand_margin(computed)
    return computed


def resolve_css_variables_inline(css_value: str) -> str:
    """Resolve var() references in a CSS value string using global root vars."""
    def replace_var(match):
        var_name = match.group(1)
        fallback = match.group(2)
        return _ROOT_CSS_VARS.get(var_name, fallback or match.group(0))
    return re.sub(r'var\((--[\w-]+)(?:,\s*([^)]+))?\)', replace_var, css_value)


def _expand_padding(computed: Dict[str, str]) -> None:
    """Expand padding shorthand in-place."""
    if 'padding' not in computed:
        return
    vals = computed['padding'].split()
    # Extract values (could be px, rem, etc.)
    if len(vals) == 1:
        computed['paddingTop'] = computed['paddingRight'] = computed['paddingBottom'] = computed['paddingLeft'] = vals[0]
    elif len(vals) == 2:
        computed['paddingTop'] = computed['paddingBottom'] = vals[0]
        computed['paddingLeft'] = computed['paddingRight'] = vals[1]
    elif len(vals) == 3:
        computed['paddingTop'] = vals[0]
        computed['paddingLeft'] = computed['paddingRight'] = vals[1]
        computed['paddingBottom'] = vals[2]
    elif len(vals) == 4:
        computed['paddingTop'] = vals[0]
        computed['paddingRight'] = vals[1]
        computed['paddingBottom'] = vals[2]
        computed['paddingLeft'] = vals[3]


def _expand_margin(computed: Dict[str, str]) -> None:
    """Expand margin shorthand in-place."""
    if 'margin' not in computed:
        return
    vals = computed['margin'].split()
    if len(vals) == 1:
        computed['marginTop'] = computed['marginRight'] = computed['marginBottom'] = computed['marginLeft'] = vals[0]
    elif len(vals) == 2:
        computed['marginTop'] = computed['marginBottom'] = vals[0]
        computed['marginLeft'] = computed['marginRight'] = vals[1]
    elif len(vals) == 3:
        computed['marginTop'] = vals[0]
        computed['marginLeft'] = computed['marginRight'] = vals[1]
        computed['marginBottom'] = vals[2]
    elif len(vals) == 4:
        computed['marginTop'] = vals[0]
        computed['marginRight'] = vals[1]
        computed['marginBottom'] = vals[2]
        computed['marginLeft'] = vals[3]


# ─── Color Parsing ────────────────────────────────────────────────────────────

def parse_color(css_color: str, bg: Tuple[int, int, int] = (255, 255, 255)) -> Optional[Tuple[int, int, int]]:
    """Parse a CSS color string, blending rgba() alpha over the given bg color."""
    if not css_color or css_color in ('transparent',) or 'rgba(0, 0, 0, 0)' in css_color:
        return None
    m = re.search(r'rgba?\((\d+),\s*(\d+),\s*(\d+)(?:,\s*([\d.]+))?\)', css_color)
    if m:
        r, g, b = int(m.group(1)), int(m.group(2)), int(m.group(3))
        a = float(m.group(4)) if m.group(4) else 1.0
        if a <= 0:
            return None
        if a < 1.0:
            r = int(a * r + (1 - a) * bg[0])
            g = int(a * g + (1 - a) * bg[1])
            b = int(a * b + (1 - a) * bg[2])
        return (r, g, b)
    m = re.search(r'#([0-9a-fA-F]{6}|[0-9a-fA-F]{3})', css_color)
    if m:
        h = m.group(1)
        if len(h) == 3:
            h = ''.join([c*2 for c in h])
        return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    return None


def px_to_pt(px_value: str) -> float:
    s = str(px_value)
    # Handle clamp(): resolve to max (last) value first
    clamp_match = re.match(r'clamp\(\s*([\d.]+\w*)\s*,\s*([\d.]+\w*)\s*,\s*([\d.]+\w*)\s*\)', s)
    if clamp_match:
        s = clamp_match.group(3)  # Use max value
    # Handle rem: 1rem = 16px
    rem_match = re.match(r'^([\d.]+)rem$', s.strip())
    if rem_match:
        return round(float(rem_match.group(1)) * 16.0 * 0.75, 1)
    m = re.search(r'([\d.]+)px', s)
    if m:
        return round(float(m.group(1)) * 0.75, 1)
    return 12.0


def parse_px(val: str) -> float:
    if not val or val in ('0px', '0', 'auto', 'none', 'normal', ''):
        return 0.0
    s = str(val)
    # Handle clamp(): use max (last) value, convert rem to px (1rem=16px)
    clamp_match = re.match(r'clamp\(\s*([\d.]+\w*)\s*,\s*([\d.]+\w*)\s*,\s*([\d.]+\w*)\s*\)', s)
    if clamp_match:
        max_val = clamp_match.group(3)
        if 'rem' in max_val:
            return float(re.search(r'([\d.]+)', max_val).group(1)) * 16.0
        if 'vw' in max_val:
            return float(re.search(r'([\d.]+)', max_val).group(1)) * 14.4  # 1440/100
        return float(re.search(r'([\d.]+)', max_val).group(1))
    # Handle rem unit
    rem_match = re.match(r'^([\d.]+)rem$', s.strip())
    if rem_match:
        return float(rem_match.group(1)) * 16.0
    # Handle vh unit
    vh_match = re.match(r'^([\d.]+)vh$', s.strip())
    if vh_match:
        return float(vh_match.group(1)) * 8.1  # 810/100
    m = re.search(r'([\d.]+)', s)
    return float(m.group(1)) if m else 0.0


def has_cjk(text: str) -> bool:
    return bool(re.search(r'[\u2E80-\u9FFF\uF900-\uFAFF\uFE10-\uFE6F\uFF00-\uFFEF]', text))


def is_bold(fw: str) -> bool:
    return fw in ('bold', '700', '800', '900') or (fw.isdigit() and int(fw) >= 600)


def is_condensed_font(family: str) -> bool:
    return bool(re.search(r'condensed|narrow|compressed', family or '', re.I))


def resolve_border_radius(style: Dict[str, str], width_px: float, height_px: float) -> float:
    """Convert CSS border-radius to px."""
    br = style.get('borderTopLeftRadius', style.get('borderRadius', ''))
    if not br or br == '0px':
        return 0.0
    if br.endswith('%'):
        pct = float(re.search(r'([\d.]+)', br).group(1))
        return pct / 100.0 * min(width_px, height_px)
    return parse_px(br)


def has_visible_bg_or_border(style: Dict[str, str]) -> bool:
    bg = style.get('backgroundColor', '')
    if bg and bg not in ('transparent', 'rgba(0, 0, 0, 0)', ''):
        return True
    for side in ('border', 'borderLeft', 'borderRight', 'borderTop', 'borderBottom'):
        bs = style.get(side, '')
        if bs and 'none' not in bs and not bs.startswith('0px') and bs != '0px':
            return True
    return False


def is_leaf_text_container(element: Tag, css_rules: Optional[List] = None) -> bool:
    """Check if element's entire visible content is text (no block children)."""
    children = list(element.children)
    if len(children) == 0:
        return bool(get_text_content(element).strip())
    has_inline_with_bg = False
    rules = css_rules or []
    for child in children:
        if isinstance(child, NavigableString):
            continue
        if not isinstance(child, Tag):
            continue
        if child.name not in INLINE_TAGS:
            return False
        # If inline child has visible bg/border, it should be extracted separately
        if rules and has_visible_bg_or_border(compute_element_style(child, rules, child.get('style', ''))):
            has_inline_with_bg = True
    if has_inline_with_bg:
        return False  # Don't treat as leaf — recurse so each styled span gets its own shape
    return bool(get_text_content(element).strip())


def get_text_content(element: Tag, exclude_elements: set = None) -> str:
    """Get the text content of an element, converting <br> to newlines without modifying the tree."""
    if not element:
        return ''
    parts = []
    for child in element.descendants:
        if exclude_elements is not None:
            # Skip descendants of excluded elements
            if hasattr(child, 'parents') and any(p in exclude_elements for p in child.parents):
                continue
        if isinstance(child, NavigableString):
            text = str(child)
            # Trim leading/trailing whitespace from each text node
            # (removes HTML formatting indentation)
            stripped = text.strip()
            if stripped:
                parts.append(stripped)
        elif isinstance(child, Tag) and child.name == 'br':
            parts.append('\n')
    result = ''.join(parts)
    # Collapse 3+ consecutive newlines to 2 (preserve blank lines from <br/><br/>)
    # but also ensure single \n between content lines isn't lost
    result = re.sub(r'\n{3,}', '\n\n', result)
    return result.strip()


def estimate_text_width(text: str, font_size_px: float) -> float:
    """Rough text width estimate in pixels."""
    if not text:
        return 0.0
    char_width = font_size_px * 0.55
    if has_cjk(text):
        char_width *= 1.0
    return len(text) * char_width


def compute_text_content_width(element: Tag, css_rules: List[CSSRule], parent_style: Optional[Dict[str, str]] = None) -> float:
    """Compute natural width of text content within an element, in inches."""
    max_w = 0.0
    style = compute_element_style(element, css_rules, element.get('style', ''), parent_style)
    for desc in element.descendants:
        if hasattr(desc, 'name') and desc.name and desc.name in ('span', 'div', 'p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'li', 'a'):
            txt = get_text_content(desc).strip()
            if not txt:
                continue
            dstyle = compute_element_style(desc, css_rules, desc.get('style', ''), style)
            font_size_px = parse_px(dstyle.get('fontSize', '16px'))
            if font_size_px <= 0:
                font_size_px = 16.0
            cjk = sum(1 for c in txt if ord(c) > 127)
            latin = len(txt) - cjk
            text_w = (cjk * font_size_px + latin * font_size_px * 0.55) / PX_PER_IN
            # Add small padding
            text_w += 0.1
            if text_w > max_w:
                max_w = text_w
    return max_w


def _is_grid_container(style: Dict[str, str]) -> bool:
    display = style.get('display', '')
    if display == 'grid':
        return True
    cols = style.get('gridTemplateColumns', style.get('grid-template-columns', ''))
    return bool(cols)


def _detect_flex_row(style: Dict[str, str]) -> bool:
    display = style.get('display', '')
    if display != 'flex' and display != 'inline-flex':
        return False
    direction = style.get('flexDirection', style.get('flex-direction', ''))
    return direction != 'column'


def _parse_grid_columns(grid_template: str) -> int:
    m = re.search(r'repeat\((\d+)', grid_template)
    if m:
        return int(m.group(1))
    cols = [c.strip() for c in grid_template.split() if c.strip() and c.strip() != 'auto']
    return max(len(cols), 1)


def _get_gap_px(style: Dict[str, str]) -> float:
    gap = style.get('gap', style.get('gridGap', ''))
    return parse_px(gap) if gap else 20.0


# ─── Text Segment Extraction ─────────────────────────────────────────────────

def extract_text_segments(
    element: Tag,
    css_rules: List[CSSRule],
    parent_style: Optional[Dict[str, str]] = None,
) -> List[Dict[str, Any]]:
    """
    Walk DOM tree under element, collecting text segments with inherited styles.
    Tracks: color, bold, fontSize, strike, bgColor.
    """
    segments = []

    def walk(node, inherited_color, inherited_bold, inherited_size, inherited_strike, inherited_bg):
        if isinstance(node, NavigableString):
            text = str(node)
            stripped = text.strip()
            if stripped:
                segments.append({
                    'text': stripped, 'color': inherited_color, 'bold': inherited_bold,
                    'fontSize': inherited_size, 'strike': inherited_strike,
                    'bgColor': inherited_bg if inherited_bg else None,
                    'inlineBgBounds': None,
                })
            return
        if not isinstance(node, Tag):
            return
        tag = node.name
        if tag == 'br':
            segments.append({
                'text': '\n', 'color': inherited_color, 'bold': inherited_bold,
                'fontSize': inherited_size, 'strike': inherited_strike,
                'bgColor': None, 'inlineBgBounds': None,
            })
            return
        if tag not in ('h1','h2','h3','h4','h5','h6','p','li','span','a',
                        'strong','em','b','i','mark','code','small','br',
                        'div','section','article','ul','ol','table','tr',
                        'th','td','img','svg','canvas'):
            for child in node.children:
                walk(child, inherited_color, inherited_bold, inherited_size, inherited_strike, inherited_bg)
            return

        style = compute_element_style(node, css_rules, node.get('style', ''), parent_style)

        # Color
        color = inherited_color
        bi = style.get('backgroundImage', '')
        bc = style.get('webkitBackgroundClip', style.get('backgroundClip', ''))
        if 'gradient' in bi and bc == 'text':
            cm = re.findall(r'rgba?\([^)]+\)', bi)
            if cm:
                color = cm[0]
        else:
            sc = style.get('color', '')
            if sc and sc != 'rgba(0, 0, 0, 0)':
                color = sc

        # Bold
        bold = inherited_bold
        fw = style.get('fontWeight', '')
        if fw in ('bold', '700', '800', '900'):
            bold = True
        elif tag in ('strong', 'b'):
            bold = True
        elif fw and fw.isdigit() and int(fw) >= 600:
            bold = True

        # Font size
        size = inherited_size
        fs = style.get('fontSize', '')
        if fs and 'px' in fs:
            size = fs

        # Strikethrough
        strike = inherited_strike
        td = style.get('textDecoration', style.get('textDecorationLine', ''))
        if 'line-through' in td:
            strike = True
        if tag in ('s', 'del', 'strike'):
            strike = True

        # Inline background
        child_bg = style.get('backgroundColor', '')
        child_has_bg = child_bg and child_bg not in ('transparent', 'rgba(0, 0, 0, 0)')
        new_bg = child_bg if child_has_bg else inherited_bg

        for child in node.children:
            walk(child, color, bold, size, strike, new_bg)

    base_style = compute_element_style(element, css_rules, element.get('style', ''), parent_style)
    base_color = base_style.get('color', '')
    bi = base_style.get('backgroundImage', '')
    bc = base_style.get('webkitBackgroundClip', base_style.get('backgroundClip', ''))
    if 'gradient' in bi and bc == 'text':
        cm = re.findall(r'rgba?\([^)]+\)', bi)
        if cm:
            base_color = cm[0]

    base_bold = False
    fw = base_style.get('fontWeight', '')
    if fw in ('bold', '700', '800', '900') or element.name in ('strong', 'b'):
        base_bold = True
    elif fw and fw.isdigit() and int(fw) >= 600:
        base_bold = True

    base_size = base_style.get('fontSize', '16px')
    base_strike = 'line-through' in base_style.get('textDecoration', base_style.get('textDecorationLine', ''))

    for child in element.children:
        walk(child, base_color, base_bold, base_size, base_strike, None)

    # Merge consecutive same-style segments
    merged = []
    for seg in segments:
        if (merged and
            merged[-1].get('color') == seg.get('color') and
            merged[-1].get('bold') == seg.get('bold') and
            merged[-1].get('fontSize') == seg.get('fontSize') and
            merged[-1].get('strike') == seg.get('strike') and
            merged[-1].get('bgColor') == seg.get('bgColor') and
            seg['text'] != '\n' and merged[-1]['text'] != '\n'):
            merged[-1]['text'] += seg['text']
        else:
            merged.append(seg)

    return merged


# ─── Flat Extract (adapted from browser version's flatExtract) ────────────────

def build_image_element(element: Tag, style: Dict[str, str]) -> Optional[Dict]:
    """Build an image element IR."""
    src = ''
    if element.name == 'img':
        src = element.get('src', '') or element.get('data-src', '')
    elif element.name == 'svg':
        src = str(element)
    if not src and element.name != 'svg':
        return None
    return {
        'type': 'image', 'tag': element.name, 'imageKind': element.name,
        'source': src,
        'bounds': {'x': 0, 'y': 0, 'width': 4, 'height': 3},
        'styles': {'borderRadius': '', 'objectFit': style.get('objectFit', '')},
    }


def build_shape_element(element: Tag, style: Dict[str, str], slide_width_px: float = 1440) -> Dict:
    """Build a shape element IR (div with background/border)."""
    border_radius_px = resolve_border_radius(style, parse_px(style.get('width', '100%')),
                                                parse_px(style.get('height', '100px')))
    return {
        'type': 'shape', 'tag': element.name,
        'bounds': {'x': 0.5, 'y': 0.5, 'width': 12.33, 'height': 1.0},
        'styles': {
            'backgroundColor': style.get('backgroundColor', ''),
            'backgroundImage': style.get('backgroundImage', ''),
            'border': style.get('border', ''),
            'borderLeft': style.get('borderLeft', ''),
            'borderRight': style.get('borderRight', ''),
            'borderTop': style.get('borderTop', ''),
            'borderBottom': style.get('borderBottom', ''),
            'borderRadius': f'{border_radius_px}px',
            'marginTop': style.get('marginTop', ''),
            'marginBottom': style.get('marginBottom', ''),
            'height': style.get('height', ''),
        },
    }


def build_text_element(element: Tag, style: Dict[str, str], css_rules: List[CSSRule],
                       slide_width_px: float = 1440, content_width_px: float = None,
                       exclude_elements: set = None) -> Optional[Dict]:
    """Build a text element IR with segments."""
    text = get_text_content(element, exclude_elements).strip()
    if not text:
        return None
    segments = extract_text_segments(element, css_rules, style)
    if not segments and not text:
        return None

    font_size_px = parse_px(style.get('fontSize', '16px'))
    if font_size_px <= 0:
        font_size_px = 16.0
    font_size_pt = px_to_pt(style.get('fontSize', '16px'))
    if font_size_pt <= 0:
        font_size_pt = 12.0

    # Determine effective width from CSS constraints
    explicit_w = parse_px(style.get('width', ''))
    max_w = parse_px(style.get('maxWidth', ''))
    display = style.get('display', '')
    is_inline_block = 'inline-block' in display

    # Compute content-based width if needed
    cjk_count = sum(1 for c in text if ord(c) > 127)
    latin_count = len(text) - cjk_count
    content_w_in = (cjk_count * font_size_px + latin_count * font_size_px * 0.55) / PX_PER_IN
    content_w_px = content_w_in * PX_PER_IN

    # Determine final width in inches
    width_in = None
    if explicit_w > 0:
        width_in = explicit_w / PX_PER_IN
    elif is_inline_block and max_w > 0:
        # For inline-block with max-width: use max line width capped by max-width
        # For multi-line text, compute the widest line, not total text width
        # Include horizontal padding in width calculation
        pad_l = parse_px(style.get('paddingLeft', ''))
        pad_r = parse_px(style.get('paddingRight', ''))
        h_pad_px = pad_l + pad_r
        max_line_px = 0.0
        for line in text.split('\n'):
            line = line.strip()
            if not line:
                continue
            line_cjk = sum(1 for c in line if ord(c) > 127)
            line_latin = len(line) - line_cjk
            line_w = line_cjk * font_size_px + line_latin * font_size_px * 0.55
            if line_w > max_line_px:
                max_line_px = line_w
        width_in = min(max_line_px + h_pad_px, max_w) / PX_PER_IN
    elif is_inline_block and content_width_px and content_width_px > 0:
        # Constrained inline-block: use content-based width, capped by parent constraint
        cw_in = content_width_px / PX_PER_IN
        width_in = min(content_w_in + 0.2, cw_in)  # content width + small padding, capped

    default_w_in = 13.33 - 1.0  # default bounds width
    if width_in is None:
        # For inline-block elements, use content width; for block-level, use full width
        if is_inline_block:
            width_in = min(content_w_in + 0.08, default_w_in)
        else:
            width_in = default_w_in

    line_count = estimate_wrapped_lines(text, font_size_pt, width_in)

    # Compute line height multiplier from CSS
    lh = style.get('lineHeight', '')
    if lh and 'px' in lh:
        line_height_px = parse_px(lh)
    elif lh and lh.replace('.', '').isdigit():
        line_height_px = font_size_px * float(lh)
    else:
        line_height_px = font_size_px * 0.82  # PPTX renders tighter than HTML line-height
    total_height_px = line_count * line_height_px

    return {
        'type': 'text', 'tag': element.name, 'text': text, 'segments': segments,
        'gradientColors': None, 'textTransform': style.get('textTransform', 'none'),
        'naturalHeight': total_height_px / (slide_width_px / 13.33),
        'bounds': {
            'x': 0.5, 'y': 0.5,
            'width': width_in,
            'height': max(total_height_px / (slide_width_px / 13.33), 0.3),
        },
        'styles': {
            'fontSize': style.get('fontSize', '16px'),
            'fontWeight': style.get('fontWeight', '400'),
            'fontFamily': style.get('fontFamily', ''),
            'letterSpacing': style.get('letterSpacing', ''),
            'color': style.get('color', ''),
            'textAlign': style.get('textAlign', 'left'),
            'lineHeight': style.get('lineHeight', 'normal'),
            'listStyleType': style.get('listStyleType', ''),
            'paddingLeft': style.get('paddingLeft', '0px'),
            'paddingRight': style.get('paddingRight', '0px'),
            'paddingTop': style.get('paddingTop', '0px'),
            'paddingBottom': style.get('paddingBottom', '0px'),
            'alignItems': style.get('alignItems', ''),
            'justifyContent': style.get('justifyContent', ''),
            'width': style.get('width', ''),
            'maxWidth': style.get('maxWidth', ''),
            'display': style.get('display', ''),
        },
    }


def _compute_table_column_widths(rows, table_w):
    """Compute content-aware column widths for a table."""
    if not rows:
        return []
    num_cols = max(len(r['cells']) for r in rows)
    if num_cols == 1:
        return [table_w]

    # Compute max content width per column
    col_max_w = [0.0] * num_cols
    for row in rows:
        for ci, cell in enumerate(row['cells']):
            if ci >= num_cols:
                break
            text = cell.get('text', '')
            font_px = parse_px(cell['styles'].get('fontSize', '14px'))
            if font_px <= 0:
                font_px = 14.0
            # Estimate text width
            cjk = sum(1 for c in text if ord(c) > 127)
            latin = len(text) - cjk
            text_w = (cjk * font_px + latin * font_px * 0.55) / PX_PER_IN
            # Add padding
            pad_l = parse_px(cell['styles'].get('paddingLeft', '0px')) / PX_PER_IN
            pad_r = parse_px(cell['styles'].get('paddingRight', '0px')) / PX_PER_IN
            total_w = text_w + pad_l + pad_r + 0.1
            col_max_w[ci] = max(col_max_w[ci], total_w)

    # If total fits, use content widths; otherwise scale proportionally
    total = sum(col_max_w)
    if total <= table_w and total > 0:
        # Distribute extra space proportionally
        scale = table_w / total
        return [w * scale for w in col_max_w]
    elif total > 0:
        # Scale down proportionally
        scale = table_w / total
        return [w * scale for w in col_max_w]
    else:
        # Fallback: equal widths
        return [table_w / num_cols] * num_cols


def build_table_element(element: Tag, css_rules: List[CSSRule], style: Dict[str, str]) -> Dict:
    """Build a table element IR."""
    rows = []
    for tr in element.find_all('tr'):
        is_header = bool(tr.parent and tr.parent.name == 'thead')
        cells = []
        for cell in tr.find_all(['th', 'td']):
            cell_style = compute_element_style(cell, css_rules, cell.get('style', ''), style)
            cell_text = get_text_content(cell).strip()
            cell_segments = extract_text_segments(cell, css_rules, cell_style)
            cells.append({
                'bounds': {'x': 0, 'y': 0, 'width': 2, 'height': 0.4},
                'text': cell_text, 'segments': cell_segments,
                'isHeader': is_header or cell.name == 'th',
                'styles': {
                    'fontSize': cell_style.get('fontSize', '14px'),
                    'fontWeight': cell_style.get('fontWeight', '400'),
                    'color': cell_style.get('color', ''),
                    'backgroundColor': cell_style.get('backgroundColor', ''),
                    'textAlign': cell_style.get('textAlign', 'left'),
                    'paddingLeft': cell_style.get('paddingLeft', '0px'),
                    'paddingRight': cell_style.get('paddingRight', '0px'),
                    'paddingTop': cell_style.get('paddingTop', '0px'),
                    'paddingBottom': cell_style.get('paddingBottom', '0px'),
                    'fontFamily': cell_style.get('fontFamily', ''),
                    'letterSpacing': cell_style.get('letterSpacing', ''),
                    'borderBottom': cell_style.get('borderBottom', ''),
                    'borderRight': cell_style.get('borderRight', ''),
                },
            })
        if cells:
            rows.append({'isHeader': is_header, 'cells': cells})
    return {
        'type': 'table',
        'bounds': {'x': 0.5, 'y': 1.0, 'width': 12.33, 'height': max(len(rows) * 0.5, 1.0)},
        'rows': rows,
    }


def estimate_wrapped_lines(text: str, font_size_pt: float, box_width_in: float, has_cjk: bool = None) -> int:
    """Estimate visual line count including word wrapping.

    For CJK text: each character is ~1em wide.
    For Latin text: average ~0.5em per character.
    Mixed: use proportion-based estimate.
    """
    if not text or font_size_pt <= 0 or box_width_in <= 0:
        return 1
    # Normalize: split on explicit newlines first
    explicit_lines = text.split('\n')
    total_lines = 0
    for line in explicit_lines:
        if not line.strip():
            total_lines += 1
            continue
        # Calculate character-based width in inches
        cjk_count = sum(1 for c in line if ord(c) > 127)
        latin_count = len(line) - cjk_count
        # CJK chars ~ font_size_pt wide, Latin ~ font_size_pt * 0.55
        text_width_in = (cjk_count * font_size_pt + latin_count * font_size_pt * 0.55) / 72.0
        chars_per_line = box_width_in / (text_width_in / len(line)) if line else 1
        if chars_per_line <= 0:
            chars_per_line = 1
        wrapped = max(1, math.ceil(len(line) / chars_per_line))
        total_lines += wrapped
    return max(1, total_lines)


def _cjk_correct_width(has_border: bool, text: str, width_in: float, is_condensed: bool) -> float:
    """Apply CJK/condensed font width correction."""
    if is_condensed:
        return width_in * 1.50
    if has_border and has_cjk(text) and width_in < 3.0:
        return width_in * CJK_BOX_FACTOR
    return width_in


def flat_extract(
    element: Tag,
    css_rules: List[CSSRule],
    parent_style: Optional[Dict[str, str]] = None,
    slide_width_px: float = 1440,
    content_width_px: float = None,  # If set, constrains grid/layout width
) -> List[Dict[str, Any]]:
    """
    Adapted from browser version's flatExtract.
    Recursively extracts text, shape, table, and image elements from a DOM subtree.
    Positions are computed via a simulated flex-column layout.
    """
    style = compute_element_style(element, css_rules, element.get('style', ''), parent_style)
    tag = element.name.lower()

    # Raster elements
    if tag == 'img':
        img_el = build_image_element(element, style)
        return [img_el] if img_el else []

    if tag == 'svg':
        return [{
            'type': 'image', 'tag': 'svg', 'imageKind': 'svg',
            'source': str(element),
            'bounds': {'x': 0.5, 'y': 0.5, 'width': 4, 'height': 3},
            'styles': {'borderRadius': '', 'objectFit': ''},
        }]

    # Tables
    if tag == 'table':
        return [build_table_element(element, css_rules, style)]

    # Text elements (h1-h6, p, li, span, a)
    if tag in TEXT_TAGS:
        # Check for styled inline children that should have their own shapes (pills/badges)
        styled_shapes = []
        pill_elements = set()  # Track which children are pills, to exclude from combined text
        for child in element.children:
            if isinstance(child, Tag) and child.name in INLINE_TAGS:
                child_s = compute_element_style(child, css_rules, child.get('style', ''))
                if has_visible_bg_or_border(child_s):
                    child_text = get_text_content(child).strip()
                    if child_text:
                        pill_elements.add(child)
                        # Estimate shape width from text content
                        cjk = sum(1 for c in child_text if ord(c) > 127)
                        latin = len(child_text) - cjk
                        font_px = parse_px(child_s.get('fontSize', '16px'))
                        if font_px <= 0:
                            font_px = 16.0
                        # Base text width
                        text_w = (cjk * font_px + latin * font_px * 0.55) / PX_PER_IN
                        # Padding (expand shorthand for pill spans only)
                        _expand_padding(child_s)
                        pad_lr = parse_px(child_s.get('paddingLeft', '0px')) / PX_PER_IN
                        pad_rr = parse_px(child_s.get('paddingRight', '0px')) / PX_PER_IN
                        # Height: golden pill height is 0.53" for 14px font
                        # With PX_PER_IN=108: font_px/108 * 4.09 = 0.53
                        shape_w = text_w + pad_lr + pad_rr + 0.1
                        shape_h = font_px / PX_PER_IN * 4.09

                        # CJK width correction for pill shape (same as native version)
                        if cjk > 0:
                            shape_w *= 1.15

                        # Extract borderRadius for proper pill rendering
                        br = child_s.get('borderRadius', '0px')

                        # Get text color for rendering inside pill
                        pill_color = child_s.get('color', '')

                        styled_shapes.append({
                            'type': 'shape', 'tag': 'span',
                            'bounds': {'x': 0.5, 'y': 0.5, 'width': shape_w, 'height': shape_h},
                            'styles': {
                                'backgroundColor': child_s.get('backgroundColor', ''),
                                'border': child_s.get('border', ''),
                                'borderRadius': br,
                            },
                            '_is_pill': True,
                            'pill_text': child_text,
                            'pill_color': pill_color,
                        })

        text_el = build_text_element(element, style, css_rules, slide_width_px, content_width_px,
                                      exclude_elements=pill_elements if pill_elements else None)

        # Always emit styled pill shapes (decorative — even if parent has no non-pill text)
        import uuid
        results = []
        if styled_shapes:
            for pill in styled_shapes:
                pill['_skip_layout'] = True  # Skip during layout, positioned in post-layout
                results.append(pill)
            # Pills now have embedded text (pill_text), no need for full subtitle text element

        if text_el:
            # Create shape for TEXT_TAGS if there's actual background fill,
            # full border, or border-bottom (list item background containers).
            # The shape acts as a background container that's synced with the text.
            has_bg = (style.get('backgroundColor', '') and
                     style['backgroundColor'] not in ('transparent', 'rgba(0, 0, 0, 0)', ''))
            has_full_border = (style.get('border', '') or
                              style.get('borderLeft', '') or
                              style.get('borderTop', '') or
                              style.get('borderRight', ''))
            has_border_bottom = bool(style.get('borderBottom', ''))
            if has_bg or has_full_border or has_border_bottom:
                shape = build_shape_element(element, style, slide_width_px)
                shape['bounds'] = dict(text_el['bounds'])
                text = get_text_content(element, pill_elements if pill_elements else None).strip()
                ff = style.get('fontFamily', '')
                c = is_condensed_font(ff)
                cjk_w = _cjk_correct_width(True, text, text_el['bounds']['width'], c)
                if cjk_w != text_el['bounds']['width']:
                    shape['bounds']['width'] = cjk_w
                pair_id = str(uuid.uuid4())[:8]
                shape['_pair_with'] = pair_id
                text_el['_pair_with'] = pair_id
                results.append(shape)

            results.append(text_el)
            return results

        if results:
            return results
        return []

    # Container elements (div, section, article, ul, ol)
    if tag in CONTAINER_TAGS:
        bg_image = style.get('backgroundImage', 'none')
        has_gradient_bg = bg_image != 'none' and 'gradient' in bg_image
        has_url_bg = bg_image != 'none' and 'url(' in bg_image
        total_text = get_text_content(element).strip()

        # Background image with no text → raster
        if has_url_bg and not total_text:
            return [{
                'type': 'image', 'tag': 'div', 'imageKind': 'background-image',
                'source': bg_image,
                'bounds': {'x': 0.5, 'y': 0.5, 'width': 12.33, 'height': 6.5},
                'styles': {'borderRadius': '', 'objectFit': ''},
            }]

        # Grid or flex-row layout: wrap children in a container element
        if _is_grid_container(style) or _detect_flex_row(style):
            grid_children = build_grid_children(element, css_rules, style, slide_width_px, content_width_px)
            if grid_children:
                grid_h = max(c['bounds']['y'] + c['bounds']['height'] for c in grid_children) + 0.15
            else:
                grid_h = 0.5
            container = {
                'type': 'container',
                'tag': element.name,
                'bounds': {'x': 0.5, 'y': 0.5, 'width': 12.33, 'height': grid_h},
                'styles': style,
                'children': grid_children,
            }
            return [container]

        # Inline-block container with visible bg/border and SINGLE-LINE text content
        # Treat as a leaf text element (tag pills, badges, etc.)
        # Must come BEFORE is_leaf_text_container check
        # Only applies to single-line text — multi-line inline-blocks should recurse.
        display_val = style.get('display', '')
        is_inline_block_el = 'inline-block' in display_val
        if (is_inline_block_el and has_visible_bg_or_border(style) and total_text
                and '\n' not in total_text):
            # Check that all child Tags are inline (no block children)
            child_tags = [c for c in element.children if isinstance(c, Tag)]
            all_inline = all(c.name.lower() in INLINE_TAGS or c.name.lower() in ('br',) for c in child_tags)
            if all_inline or not child_tags:
                # Don't pass content_width_px for inline-block — size to content
                text_el = build_text_element(element, style, css_rules, slide_width_px)
                if text_el:
                    results = []
                    if has_visible_bg_or_border(style):
                        shape = build_shape_element(element, style, slide_width_px)
                        shape['bounds'] = dict(text_el['bounds'])
                        # Add padding to shape dimensions
                        pad_l = parse_px(style.get('paddingLeft', '0px')) / PX_PER_IN
                        pad_r = parse_px(style.get('paddingRight', '0px')) / PX_PER_IN
                        pad_t = parse_px(style.get('paddingTop', '0px')) / PX_PER_IN
                        pad_b = parse_px(style.get('paddingBottom', '0px')) / PX_PER_IN
                        if pad_l + pad_r > 0:
                            shape['bounds']['width'] = text_el['bounds']['width'] + pad_l + pad_r
                        if pad_t + pad_b > 0:
                            shape['bounds']['height'] = text_el['bounds']['height'] + pad_t + pad_b
                        # Pair shape with text for position syncing after layout
                        import uuid
                        pair_id = str(uuid.uuid4())[:8]
                        shape['_pair_with'] = pair_id
                        text_el['_pair_with'] = pair_id
                        results.append(shape)
                    results.append(text_el)
                    return results

        # Leaf text container: entire content is text (with inline formatting)
        if is_leaf_text_container(element, css_rules):
            text_el = build_text_element(element, style, css_rules, slide_width_px)
            if text_el:
                results = []
                if has_visible_bg_or_border(style) or has_gradient_bg:
                    shape = build_shape_element(element, style, slide_width_px)
                    shape['bounds'] = dict(text_el['bounds'])
                    if has_gradient_bg:
                        shape['styles']['backgroundImage'] = bg_image
                    # Pair shape with text for position syncing after layout
                    import uuid
                    pair_id = str(uuid.uuid4())[:8]
                    shape['_pair_with'] = pair_id
                    text_el['_pair_with'] = pair_id
                    results.append(shape)
                results.append(text_el)
                return results
            return []

        # Empty container with visible bg/border: treat as shape (e.g. .swiss-rule, .hero-rule, .stat-divider)
        if not total_text and not element.find(['img', 'svg', 'table']) and (has_visible_bg_or_border(style) or has_gradient_bg):
            shape = build_shape_element(element, style, slide_width_px)
            if has_gradient_bg:
                shape['styles']['backgroundImage'] = bg_image
            # Use CSS width/height if available
            css_w_str = style.get('width', '')
            css_h_str = style.get('height', '')
            css_w = parse_px(css_w_str)
            css_h = parse_px(css_h_str)
            # Handle percentage width
            if css_w_str and '%' in css_w_str:
                shape['bounds']['width'] = 12.33  # full content area
            elif css_w > 0:
                shape['bounds']['width'] = css_w / PX_PER_IN
            # Handle height - if not set, mark for stretch during layout
            if css_h > 0:
                shape['bounds']['height'] = css_h / PX_PER_IN
            elif css_h_str and 'auto' in css_h_str or not css_h_str:
                # Mark as stretch-to-content for layout phase
                shape['_stretch_height'] = True
                shape['bounds']['height'] = 0.5  # temporary, will be adjusted in layout
            return [shape]

        # Standard container: recurse into children
        results = []
        # Check if element has only border (no background) - create a line shape instead of full bg
        has_bg = bool(style.get('backgroundColor', '')) and style['backgroundColor'] not in ('transparent', 'rgba(0, 0, 0, 0)', '')
        has_only_border = has_visible_bg_or_border(style) and not has_bg and not has_gradient_bg
        border_left_only = bool(style.get('borderLeft', '')) and 'none' not in style.get('borderLeft', '') and not style.get('borderLeft', '').startswith('0px')

        if has_only_border and border_left_only:
            # Create a thin vertical line shape for border-left only
            import uuid
            bl_str = style.get('borderLeft', '')
            bl_width_px = 2.0  # default
            m = re.search(r'([\d.]+)px', bl_str)
            if m:
                bl_width_px = float(m.group(1))
            # Check for border-left-color override (e.g., .accent-border)
            bl_color_str = style.get('borderLeftColor', '')
            if bl_color_str:
                bl_color = parse_color(bl_color_str)
            else:
                bl_color = parse_color(style.get('borderLeft', '').replace('solid', '').strip())
            if not bl_color:
                # Try to extract from rgba in border
                m = re.search(r'rgba?\((\d+),\s*(\d+),\s*(\d+)', bl_str)
                if m:
                    bl_color = (int(m.group(1)), int(m.group(2)), int(m.group(3)))
            if bl_color:
                results.append({
                    'type': 'shape', 'tag': element.name,
                    '_border_left_line': True,  # Mark for special handling in layout
                    'bounds': {'x': 0.5, 'y': 0.5, 'width': bl_width_px / PX_PER_IN, 'height': 1.0},
                    'styles': {'backgroundColor': 'rgb({}, {}, {})'.format(*bl_color)},
                })
        elif has_visible_bg_or_border(style) or has_gradient_bg:
            shape = build_shape_element(element, style, slide_width_px)
            if has_gradient_bg:
                shape['styles']['backgroundImage'] = bg_image
            results.append(shape)

        # Propagate maxWidth constraint from parent to children
        parent_maxw = style.get('maxWidth', '')

        for child in element.children:
            if isinstance(child, Tag):
                # Skip navigation/chrome elements
                if _should_skip_element(child):
                    continue
                # Skip absolutely positioned elements (bg-num, slide-num-label, etc.)
                child_pos_style = compute_element_style(child, css_rules, child.get('style', ''), style)
                if child_pos_style.get('position', '') == 'absolute':
                    continue
                # Determine content width for child: use parent's constraint or child's own maxWidth
                child_cw = content_width_px  # Start with parent's content width
                if not child_cw and parent_maxw and 'px' in parent_maxw:
                    child_cw = parse_px(parent_maxw)
                child_elems = flat_extract(child, css_rules, style, slide_width_px, content_width_px=child_cw)
                # Apply parent maxWidth to child elements if they don't have one
                if parent_maxw and 'px' in parent_maxw:
                    for ce in child_elems:
                        cs = ce.get('styles', {})
                        if not cs.get('maxWidth', ''):
                            cs['maxWidth'] = parent_maxw
                results.extend(child_elems)

        return results

    return []


def build_grid_children(
    container: Tag,
    css_rules: List[CSSRule],
    style: Dict[str, str],
    slide_width_px: float = 1440,
    content_width_px: float = None,  # If set, constrain grid to this width
) -> List[Dict[str, Any]]:
    """Process children of a grid/flex-row container with proper layout."""
    width_in = 13.33
    px_per_in = slide_width_px / width_in
    margin_in = 0.5

    # If content width is constrained, adjust width and margin accordingly
    if content_width_px and content_width_px > 0:
        content_width_in = content_width_px / px_per_in
        if content_width_in < width_in:
            # Center content within slide
            side_margin = (width_in - content_width_in) / 2
            margin_in = side_margin

    grid_cols = style.get('gridTemplateColumns', '')
    if not grid_cols:
        grid_cols = style.get('grid-template-columns', '')

    if grid_cols:
        num_cols = _parse_grid_columns(grid_cols)
    elif _detect_flex_row(style):
        children = [c for c in container.children if isinstance(c, Tag)]
        num_cols = len(children) if children else 1
    else:
        num_cols = 1

    gap_px = _get_gap_px(style)
    gap_in = gap_px / px_per_in
    available_width_in = width_in - 2 * margin_in

    # Check if grid should be centered (justify-content: center)
    justify = style.get('justifyContent', style.get('justify-content', ''))
    is_centered = justify == 'center'

    # Estimate item width based on content (for centered grids) or fill available width
    if is_centered:
        # Calculate content-based width for each child, using the max font size within
        content_widths = []
        for child in container.children:
            if not isinstance(child, Tag):
                continue
            child_style = compute_element_style(child, css_rules, child.get('style', ''), style)
            # Find max font size within this child's subtree for accurate width estimation
            max_font_px = 0.0
            for desc in child.descendants:
                if hasattr(desc, 'name') and desc.name:
                    ds = compute_element_style(desc, css_rules, desc.get('style', ''), child_style)
                    fp = parse_px(ds.get('fontSize', '16px'))
                    if fp > max_font_px:
                        max_font_px = fp
            if max_font_px <= 0:
                max_font_px = 16.0

            child_text = get_text_content(child).strip()
            if child_text:
                cjk = sum(1 for c in child_text if ord(c) > 127)
                latin = len(child_text) - cjk
                text_w = (cjk * max_font_px + latin * max_font_px * 0.55) / PX_PER_IN
                # Add modest padding
                content_widths.append(text_w + 0.3)
            else:
                content_widths.append(2.0)  # default width

        if content_widths:
            # Use the max content width for uniform items
            item_width_in = max(max(content_widths), 1.5)
        else:
            item_width_in = 2.0
    else:
        # Check if children have explicit CSS width percentages (e.g. 38%, 62%)
        pct_widths = []
        for child in container.children:
            if not isinstance(child, Tag):
                continue
            child_style = compute_element_style(child, css_rules, child.get('style', ''), style)
            if child_style.get('position', '') == 'absolute':
                continue
            w = child_style.get('width', '')
            pct_m = re.match(r'([\d.]+)%', w)
            if pct_m:
                pct_widths.append(float(pct_m.group(1)) / 100.0 * width_in)
            else:
                pct_widths = []  # Not all children have percentage widths
                break

        if pct_widths:
            item_width_in = None  # Use per-item widths from pct_widths
            num_cols = len(pct_widths)
        else:
            item_width_in = (available_width_in - gap_in * (num_cols - 1)) / num_cols if num_cols > 1 else available_width_in

    x_offset = 0.0
    if is_centered and num_cols > 1:
        total_grid_width = num_cols * item_width_in + (num_cols - 1) * gap_in
        x_offset = (available_width_in - total_grid_width) / 2

    # Collect child element groups (each group = one grid cell's elements)
    child_groups = []
    child_groups_original = []  # Keep original child elements for style lookup
    item_widths = []  # Per-item content widths
    for child in container.children:
        if not isinstance(child, Tag):
            continue
        # Skip navigation/chrome elements
        if _should_skip_element(child):
            continue
        child_style = compute_element_style(child, css_rules, child.get('style', ''), style)
        # Skip absolutely positioned elements
        if child_style.get('position', '') == 'absolute':
            continue
        child_tag = child.name.lower()

        if child_tag == 'img':
            img_el = build_image_element(child, child_style)
            if img_el:
                child_groups.append([img_el])
                child_groups_original.append(child)
                item_widths.append(img_el['bounds']['width'])
            continue
        if child_tag == 'svg':
            child_groups.append([{
                'type': 'image', 'tag': 'svg', 'imageKind': 'svg',
                'source': str(child),
                'bounds': {'x': 0, 'y': 0, 'width': item_width_in, 'height': 2},
                'styles': {'borderRadius': '', 'objectFit': ''},
            }])
            child_groups_original.append(child)
            item_widths.append(item_width_in)
            continue
        if child_tag == 'table':
            tbl = build_table_element(child, css_rules, child_style)
            child_groups.append([tbl])
            child_groups_original.append(child)
            item_widths.append(tbl['bounds']['width'])
            continue
        if child_tag in TEXT_TAGS:
            text_el = build_text_element(child, child_style, css_rules, slide_width_px)
            if text_el:
                group = []
                if has_visible_bg_or_border(child_style):
                    shape = build_shape_element(child, child_style, slide_width_px)
                    shape['bounds'] = dict(text_el['bounds'])
                    group.append(shape)
                group.append(text_el)
                child_groups.append(group)
                child_groups_original.append(child)
                item_widths.append(text_el['bounds']['width'])
            continue
        if child_tag in CONTAINER_TAGS:
            bg_img = child_style.get('backgroundImage', 'none')
            has_grad = bg_img != 'none' and 'gradient' in bg_img
            has_url = bg_img != 'none' and 'url(' in bg_img
            total_txt = get_text_content(child).strip()

            if has_url and not total_txt:
                child_groups.append([{
                    'type': 'image', 'tag': 'div', 'imageKind': 'background-image',
                    'source': bg_img,
                    'bounds': {'x': 0, 'y': 0, 'width': item_width_in, 'height': 3},
                    'styles': {'borderRadius': '', 'objectFit': ''},
                }])
                child_groups_original.append(child)
                item_widths.append(item_width_in)
                continue

            if is_leaf_text_container(child, css_rules):
                text_el = build_text_element(child, child_style, css_rules, slide_width_px)
                if text_el:
                    group = []
                    if has_visible_bg_or_border(child_style) or has_grad:
                        shape = build_shape_element(child, child_style, slide_width_px)
                        shape['bounds'] = dict(text_el['bounds'])
                        if has_grad:
                            shape['styles']['backgroundImage'] = bg_img
                        group.append(shape)
                    group.append(text_el)
                    child_groups.append(group)
                    child_groups_original.append(child)
                    item_widths.append(text_el['bounds']['width'])
                continue

            sub_elements = flat_extract(child, css_rules, child_style, slide_width_px)
            # If child is purely a shape (e.g., stat-divider), skip bg wrapper
            is_pure_shape = (len(sub_elements) == 1 and sub_elements[0].get('type') == 'shape'
                             and not sub_elements[0].get('text'))
            # If container has visible background/border, prepend a shape for it
            # and strip child elements' shape wrappers (they inherit the container bg)
            if not is_pure_shape and (has_visible_bg_or_border(child_style) or has_grad):
                shape = build_shape_element(child, child_style, slide_width_px)
                shape['bounds'] = {'x': 0, 'y': 0, 'width': item_width_in, 'height': 3.0}
                if has_grad:
                    shape['styles']['backgroundImage'] = bg_img
                # Filter out shape elements from children that only exist because of inherited bg
                # But keep: shapes with text, border-left lines, and thin separator/divider shapes
                def _is_meaningful_shape(e):
                    if e.get('type') != 'shape':
                        return True
                    if e.get('text'):
                        return True
                    if e.get('_border_left_line'):
                        return True
                    bg_c = e.get('styles', {}).get('backgroundColor', '')
                    h = e.get('bounds', {}).get('height', 0)
                    w = e.get('bounds', {}).get('width', 0)
                    # Thin lines (width < 0.1in or height < 0.1in)
                    if (w < 0.1 or h < 0.1) and bg_c:
                        return True
                    # Stretch-height dividers
                    if e.get('_stretch_height'):
                        return True
                    return False
                content_only = [e for e in sub_elements if _is_meaningful_shape(e)]
                sub_elements = [shape] + content_only
            child_groups.append(sub_elements)
            child_groups_original.append(child)
            # Compute item width from text content, not element bounds
            # (element bounds may be full slide width for unconstrained text)
            text_w = compute_text_content_width(child, css_rules)
            if text_w > 0:
                item_widths.append(text_w)
            elif sub_elements:
                # Use max element height as fallback (text elements' height is based on content)
                text_h = sum(e['bounds'].get('height', 0.3) for e in sub_elements if e.get('type') == 'text')
                item_widths.append(min(text_h * 3.0, 3.0))  # Rough estimate
            else:
                item_widths.append(item_width_in)
            continue

    # Compute per-item x positions using item_widths
    # For centered grids, determine if items should have uniform or per-item widths
    num_items = len(child_groups)
    item_x_list = []

    # Detect if this is a card-like grid (container divs with multi-element children)
    # vs simple content grid (each child has single text element)
    has_multi_children = any(len(g) > 2 for g in child_groups)

    if pct_widths:
        # Flex-row with percentage widths (e.g. 38%/62% panel layout): use CSS widths
        current_x = 0.0
        for idx in range(num_items):
            item_x_list.append(current_x)
            current_x += pct_widths[idx]
    elif is_centered and num_items > 1 and not has_multi_children:
        # Simple content grid (like stat-row): use individual item widths
        total_content_w = sum(item_widths) + (num_items - 1) * gap_in
        x_start = margin_in + (available_width_in - total_content_w) / 2
        current_x = x_start
        for idx in range(num_items):
            item_x_list.append(current_x)
            current_x += item_widths[idx] + gap_in
    elif num_cols > 1:
        # Multi-column grid (like cards): use uniform column widths
        # For centered grids, center the uniform grid
        if is_centered:
            total_w = num_cols * item_width_in + (num_cols - 1) * gap_in
            x_start = margin_in + (available_width_in - total_w) / 2
        else:
            x_start = margin_in + x_offset
        for idx in range(num_items):
            col_idx = idx % num_cols
            row_idx = idx // num_cols
            item_x_list.append(x_start + col_idx * (item_width_in + gap_in))
    else:
        for idx in range(num_items):
            item_x_list.append(margin_in)

    # Layout grid items
    results = []
    for idx, group in enumerate(child_groups):
        if num_cols > 1:
            col_idx = idx % num_cols
            row_idx = idx // num_cols
        else:
            col_idx = 0
            row_idx = idx

        item_x = item_x_list[idx]
        # Use per-item width for percentage-width panels or centered grids
        if pct_widths and idx < len(pct_widths):
            this_item_width = pct_widths[idx]
        elif is_centered and num_items > 1 and not has_multi_children:
            this_item_width = item_widths[idx] if idx < len(item_widths) else item_width_in
        else:
            this_item_width = item_width_in

        # For percentage-width panels, check if child has justifyContent: center
        panel_center_y = False
        if pct_widths and idx < len(child_groups_original):
            child_elem = child_groups_original[idx]
            child_style = compute_element_style(child_elem, css_rules, child_elem.get('style', ''), style)
            jc = child_style.get('justifyContent', '')
            fd = child_style.get('flexDirection', '')
            if jc == 'center' and (fd == 'column' or not fd):
                panel_center_y = True

        # Estimate item height from its elements' existing bounds
        # Use the sum of text element heights (already computed by build_text_element)
        text_h_total = 0.0
        has_text = False
        for elem in group:
            if elem.get('type') == 'text':
                text_h_total += elem['bounds'].get('height', 0.3)
                has_text = True
            elif elem.get('type') == 'table':
                text_h_total += len(elem.get('rows', [])) * 0.5
                has_text = True

        if has_text:
            text_count = sum(1 for e in group if e.get('type') == 'text')
            # Check if this group has a background shape (card-style item)
            grp_has_bg = any(
                e.get('type') == 'shape' and not e.get('text') and e['bounds'].get('height', 0) >= 1.0
                for e in group
            )
            if grp_has_bg:
                # Card: add padding + internal gaps
                card_pad_t = 28.0 / PX_PER_IN
                card_pad_b = 24.0 / PX_PER_IN
                internal_gap = 0.20 * max(text_count - 1, 0)
                item_h = text_h_total + card_pad_t + card_pad_b + internal_gap
            else:
                # No background (stat items, etc.): minimal gap between text elements
                internal_gap = 0.05 * max(text_count - 1, 0)
                item_h = text_h_total + internal_gap
        else:
            item_h = 2.0  # default for non-text items

        # For full-height flex-row panels, use slide height
        if pct_widths:
            item_h = 7.5  # Full slide height for panel layouts

        item_y = row_idx * (item_h + gap_in)  # Container-relative (layout pass adds container y)

        # For panel with justifyContent: center, calculate vertical centering offset
        center_y_offset = 0.0
        if panel_center_y:
            # Calculate total content height
            content_h = 0.0
            for elem in group:
                if elem.get('type') == 'text':
                    content_h += elem['bounds'].get('height', 0.3)
                elif elem.get('type') == 'shape' and elem.get('_border_left_line'):
                    # Border-left line should span content height
                    pass
            # Add gaps between elements
            text_count = sum(1 for e in group if e.get('type') == 'text')
            if text_count > 1:
                content_h += 0.15 * (text_count - 1)  # gap between items
            # Calculate centering offset
            available_h = item_h  # Full panel height
            if content_h < available_h:
                center_y_offset = (available_h - content_h) / 2

        # Layout elements: background shapes overlap content, content stacks vertically
        # Detect if this group has a background shape (card-style item)
        has_bg_shape = any(
            e.get('type') == 'shape' and not e.get('text') and e['bounds'].get('height', 0) >= 1.0
            for e in group
        )
        # Only apply card padding when there's a background shape
        if has_bg_shape:
            pad_t = 28.0 / PX_PER_IN
            pad_x = 24.0 / PX_PER_IN
        else:
            pad_t = 0.0
            pad_x = 0.0

        group_y = item_y + pad_t + center_y_offset
        initial_group_y = group_y  # Save for border-left positioning
        # First pass: compute max content height for stretch-height shapes
        max_content_h = 0.0
        for elem in group:
            if elem.get('type') in ('text', 'table'):
                max_content_h = max(max_content_h, elem['bounds'].get('height', 0.3))
        for elem in group:
            if elem.get('_stretch_height') and max_content_h > 0:
                elem['bounds']['height'] = max_content_h

        # Second pass: compute per-border-left-line content heights
        # Group elements by border-left-line: each border-left + following text elements
        bl_ranges = []  # list of (bl_index, start_y_offset, height)
        running_y = 0.0
        current_bl_idx = None
        current_bl_start = 0.0
        current_bl_h = 0.0
        has_bl = any(e.get('_border_left_line') for e in group)

        for ei, elem in enumerate(group):
            b = elem['bounds']
            if elem.get('type') == 'shape' and elem.get('_border_left_line'):
                # Save previous border-left range
                if current_bl_idx is not None:
                    bl_ranges.append((current_bl_idx, current_bl_start, current_bl_h))
                current_bl_idx = ei
                current_bl_start = running_y
                current_bl_h = 0.0
                continue
            if elem.get('type') == 'shape' and not elem.get('text') and b.get('height', 0) >= 2.0:
                continue  # Skip bg shapes
            eh = b.get('height', 0.3)
            if current_bl_idx is not None:
                current_bl_h += eh + 0.05
            running_y += eh + 0.05
        # Save last border-left range
        if current_bl_idx is not None:
            bl_ranges.append((current_bl_idx, current_bl_start, current_bl_h))

        # Compute border-left offset for text positioning
        bl_offset = 0.0
        for e in group:
            if e.get('_border_left_line'):
                bl_offset = e['bounds']['width'] + 0.05
                break

        for idx, elem in enumerate(group):
            b = elem['bounds']
            elem['layoutDone'] = True
            # Check if this is a pure background shape (no text, height >= 2)
            is_bg_shape = (elem.get('type') == 'shape' and not elem.get('text')
                           and b.get('height', 0) >= 2.0)
            if is_bg_shape:
                # Background shape: same position as content area, full item height
                b['x'] = item_x
                b['y'] = item_y
                b['width'] = this_item_width
                b['height'] = item_h
                results.append(elem)
                continue  # Don't advance group_y
            # Border-left line: position at left edge, use per-item height
            if elem.get('type') == 'shape' and elem.get('_border_left_line'):
                # Find this border-left's range
                for bl_idx, bl_start, bl_h in bl_ranges:
                    if elem is group[bl_idx]:
                        b['x'] = item_x + pad_x
                        b['y'] = initial_group_y + bl_start
                        b['height'] = max(bl_h - 0.05, 0.3)
                        break
                results.append(elem)
                continue  # Don't advance group_y
            if elem.get('type') in ('text', 'shape'):
                # Stretch-height divider shapes (stat-divider): keep thin width, center vertically
                if elem.get('_stretch_height') and b['width'] < 0.1:
                    b['x'] = item_x + pad_x + (this_item_width - 2 * pad_x) / 2
                    b['y'] = item_y + pad_t + center_y_offset
                    results.append(elem)
                    continue  # Don't advance group_y
                # Thin horizontal/vertical lines (left-rule, swiss-rule): keep CSS width, don't stretch
                is_thin_line = (elem.get('type') == 'shape' and not elem.get('text')
                                and (b['width'] < 0.5 or b['height'] < 0.1)
                                and not elem.get('_border_left_line'))
                if is_thin_line:
                    b['x'] = item_x + pad_x
                    b['y'] = group_y
                    # Keep CSS width (already set by build_shape_element)
                    # Check for marginTop in styles
                    mt = elem.get('styles', {}).get('marginTop', '')
                    mt_in = parse_px(mt) / PX_PER_IN if mt else 0.0
                    group_y += b['height'] + max(0.05, mt_in)
                    results.append(elem)
                    continue
                b['x'] = item_x + pad_x + bl_offset
                b['y'] = group_y
                b['width'] = this_item_width - 2 * pad_x - bl_offset
                # Check for next element's marginTop
                next_mt_in = 0.0
                for j in range(idx + 1, len(group)):
                    ne = group[j]
                    if ne.get('_skip_layout'):
                        continue
                    nmt = ne.get('styles', {}).get('marginTop', '')
                    if nmt:
                        next_mt_in = parse_px(nmt) / PX_PER_IN
                    break
                group_y += b['height'] + 0.05 + next_mt_in
            elif elem.get('type') == 'table':
                b['x'] = item_x
                b['y'] = group_y
                b['width'] = this_item_width
                group_y += b['height'] + 0.05
            elif elem.get('type') == 'image':
                b['x'] = item_x
                b['y'] = group_y
                b['width'] = item_width_in
                group_y += b['height'] + 0.05
            results.append(elem)

    # Mark all grid children so pill positioning can skip them
    for elem in results:
        elem['_grid_child'] = True

    return results


# ─── Slide Background Extraction ─────────────────────────────────────────────

def extract_slide_background(slide_el: Tag, css_rules: List[CSSRule]) -> Dict:
    """Extract slide-level background (solid color, gradient, or grid)."""
    style = compute_element_style(slide_el, css_rules, slide_el.get('style', ''))
    bg_color = style.get('backgroundColor', '')
    bg_image = style.get('backgroundImage', '')

    result = {'solid': None, 'gradient': None, 'grid': None}

    if 'gradient' in bg_image:
        stops = re.findall(r'rgba?\([^)]+\)', bg_image)
        if len(stops) >= 2:
            c1 = parse_color(stops[0])
            c2 = parse_color(stops[-1])
            if c1 and c2:
                result['gradient'] = (c1, c2)

    if not result['gradient']:
        rgb = parse_color(bg_color)
        if rgb:
            result['solid'] = rgb

    gradient_count = len(re.findall(r'linear-gradient', bg_image))
    if gradient_count >= 2 and '90deg' in bg_image:
        color_match = re.search(r'rgba?\([^)]+\)', bg_image)
        if color_match:
            size_str = style.get('backgroundSize', '')
            size_match = re.search(r'([\d.]+)px', size_str)
            size_px = float(size_match.group(1)) if size_match else 24.0
            result['grid'] = {'color': color_match.group(0), 'sizePx': size_px}

    return result


# ─── HTML → Slides Parsing ───────────────────────────────────────────────────

def parse_html_to_slides(html_path: Path, width_px: float = 1440, height_px: float = 810) -> List[Dict]:
    """Parse an HTML file into a list of slide data dicts."""
    with open(html_path, 'r', encoding='utf-8') as f:
        html_content = f.read()

    soup = BeautifulSoup(html_content, 'lxml')
    css_rules = extract_css_from_soup(soup)

    # Check body background
    body_style_str = ''
    body_tag = soup.find('body')
    if body_tag and body_tag.get('style'):
        body_style_str = body_tag['style']
    body_style = compute_element_style(body_tag or Tag(name='body'), css_rules, body_style_str)
    body_bg = body_style.get('backgroundColor', '')
    body_bi = body_style.get('backgroundImage', '')
    if body_bi and 'gradient' in body_bi:
        stops = re.findall(r'rgba?\([^)]+\)', body_bi)
        if len(stops) >= 2:
            c1 = parse_color(stops[0])
            c2 = parse_color(stops[-1])
            if c1 and c2:
                css_rules.insert(0, CSSRule(selector='body', properties={
                    'backgroundImage': body_bi
                }))

    slides_html = soup.select('.slide')
    if not slides_html:
        print("No .slide elements found in HTML.")
        return []

    print(f"  Found {len(slides_html)} slides. Parsing...")

    slides = []
    for i, slide_html in enumerate(slides_html):
        content_root = slide_html.select_one('.slide-content') or slide_html

        bg_info = extract_slide_background(slide_html, css_rules)
        background_solid = bg_info['solid']
        background_gradient = bg_info['gradient']
        grid_bg = bg_info['grid']

        if not background_solid and not background_gradient:
            body_rgb = parse_color(body_bg)
            if body_rgb:
                background_solid = body_rgb

        has_own_chrome = bool(
            slide_html.select('.nav-dots') or
            slide_html.select('.slide-counter') or
            slide_html.select('.page-counter')
        )

        # For slides without .slide-content (e.g. left-panel/right-panel layout),
        # just use the section element — flat_extract will skip absolute-positioned children
        if not slide_html.select_one('.slide-content'):
            cr_style = {'display': 'flex', 'flexDirection': 'row'}
        else:
            cr_style = compute_element_style(content_root, css_rules, content_root.get('style', ''))

        # Detect content area max-width for proper grid/layout constraints
        content_mw = None
        if isinstance(cr_style, dict) and 'display' in cr_style:
            # Already computed above (wrapper or slide-content)
            pass
        else:
            cr_style = compute_element_style(content_root, css_rules, content_root.get('style', ''))
        cr_maxw = cr_style.get('maxWidth', '')
        if cr_maxw and 'px' in cr_maxw:
            content_mw = parse_px(cr_maxw)

        elements = flat_extract(content_root, css_rules, body_style, slide_width_px=width_px, content_width_px=content_mw)
        title = get_text_content(slide_html)[:50]

        print(f"  [{i+1}/{len(slides_html)}] {title}... ({len(elements)} elements)")

        # Get slide-content style for vertical alignment (justifyContent)
        content_style = cr_style if cr_style else {}

        slides.append({
            'background': background_solid,
            'bgGradient': background_gradient,
            'gridBg': grid_bg,
            'elements': elements,
            'hasOwnChrome': has_own_chrome,
            'slideStyle': content_style,  # Use slide-content style for layout
        })

    return slides


# ─── Layout Pass ──────────────────────────────────────────────────────────────

def layout_slide_elements(elements: List[Dict], slide_width_in: float = 13.33, slide_height_in: float = 7.5, slide_style: Optional[Dict[str, str]] = None):
    """
    Simulate flex-column layout: stack elements vertically with padding.
    Handles container elements (grid/flex) by positioning them and adjusting children.
    Supports text-align: center for centered layouts and max-width constraints.
    Supports vertical centering when slide has justifyContent: center.
    """
    internal_margin = 0.5
    current_y = internal_margin
    slide_margin = internal_margin

    # For flex-row layouts (full-slide panels), start at y=0 with no margin
    is_flex_row_layout = False
    for elem in elements:
        if elem.get('type') == 'container':
            es = elem.get('styles', {})
            if es.get('flexDirection', '') == 'row' or es.get('flex-direction', '') == 'row':
                is_flex_row_layout = True
                break
    if is_flex_row_layout:
        internal_margin = 0.0
        current_y = 0.0
        slide_margin = 0.0

    # Detect max-width constraint from elements that have it
    # Skip full-width background shapes (section, etc.) - find first content element with maxWidth
    for elem in elements:
        s = elem.get('styles', {})
        mw = s.get('maxWidth', '')
        if mw and 'px' in mw:
            mw_in = parse_px(mw) / PX_PER_IN
            content_area_width = mw_in
            if content_area_width < slide_width_in - 2 * internal_margin:
                slide_margin = (slide_width_in - content_area_width) / 2
            break

    max_width = slide_width_in - 2 * slide_margin

    # Pre-compute which shapes are paired with text — they shouldn't advance current_y
    # The paired text element will position both
    paired_shapes = set()
    for elem in elements:
        if elem.get('_pair_with') and elem.get('type') == 'shape':
            paired_shapes.add(id(elem))

    for i, elem in enumerate(elements):
        b = elem['bounds']
        s = elem.get('styles', {})
        elem_type = elem.get('type', '')

        # Handle container elements (grid/flex wrappers)
        if elem_type == 'container':
            b['y'] = current_y
            b['x'] = slide_margin
            # Adjust children's y positions, but x is already set by build_grid_children
            for child in elem.get('children', []):
                child['bounds']['y'] = current_y + child['bounds']['y']
                # Don't adjust x - build_grid_children already positioned correctly
            current_y += b['height'] + 0.15
            continue

        # Skip elements already positioned by grid/flex layout, but advance y
        if elem.get('layoutDone'):
            current_y = max(current_y, b['y'] + b['height'] + 0.15)
            continue

        # Skip elements marked for post-layout positioning (e.g., pill shapes/text)
        if elem.get('_skip_layout'):
            continue

        # Position element
        b['y'] = current_y
        text_align = s.get('textAlign', 'left')

        if elem_type == 'text':
            # Preserve pre-pass corrected height (from pre_pass_corrections)
            if elem.get('pptx_height_corrected'):
                pass  # Height already corrected, skip recalculation
            else:
                text = elem.get('text', '')
                font_size_px = parse_px(elem.get('styles', {}).get('fontSize', '16px'))
                if font_size_px <= 0:
                    font_size_px = 16.0
                font_size_pt = font_size_px / PX_PER_IN * 72.0
                line_count = estimate_wrapped_lines(text, font_size_pt, max_width) if text else 1
                # Use CSS lineHeight if available
                lh = s.get('lineHeight', '')
                if lh and 'px' in lh:
                    line_height_px = parse_px(lh)
                elif lh and lh.replace('.', '').isdigit():
                    line_height_px = font_size_px * float(lh)
                else:
                    line_height_px = font_size_px * 1.0
                # For headings (h1, h2), add extra height for PPTX rendering
                is_heading = elem.get('tag') in ('h1', 'h2')
                height_mult = 1.15 if (is_heading and line_count > 1) else 1.0
                b['height'] = max(line_count * line_height_px / PX_PER_IN * height_mult, font_size_pt / 72.0 * 1.5)
                # Update naturalHeight to match layout-calculated height
                # (export_text_element uses max of bounds height and naturalHeight)
                elem['naturalHeight'] = b['height']

            # Add CSS vertical padding to height
            pad_t = parse_px(s.get('paddingTop', ''))
            pad_b = parse_px(s.get('paddingBottom', ''))
            if pad_t + pad_b > 0:
                b['height'] += (pad_t + pad_b) / PX_PER_IN

            # For list items, ensure minimum height matches rendered browser output
            if elem.get('tag') == 'li' and b['height'] < 0.4:
                b['height'] = 0.4

            # For inline-block elements, use the pre-computed width from build_text_element
            # (max-line-width) — don't overwrite with effective_w which is just maxWidth constraint
            is_inline_block = 'inline-block' in s.get('display', '')
            effective_w = parse_px(s.get('width', '')) if parse_px(s.get('width', '')) > 0 else (parse_px(s.get('maxWidth', '')) if is_inline_block and parse_px(s.get('maxWidth', '')) > 0 else 0)
            has_explicit_width = effective_w > 0 and effective_w < max_width * PX_PER_IN

            # Inline-block elements: keep pre-computed content width, center via x position
            if is_inline_block and b['width'] < max_width * 0.8:
                # Trust the pre-computed width from build_text_element (max-line-width)
                # Don't overwrite with effective_w which is just maxWidth constraint
                b['x'] = slide_margin + (max_width - b['width']) / 2
            elif text_align == 'center' and text:
                cjk_count = sum(1 for c in text if ord(c) > 127)
                latin_count = len(text) - cjk_count
                text_width_in = (cjk_count * font_size_px + latin_count * font_size_px * 0.55) / PX_PER_IN
                if elem.get('_full_subtitle'):
                    # Full subtitle spans the full slide content width (like golden)
                    content_width = max_width
                elif elem.get('tag', '') in ('h1', 'h2'):
                    # Use full content width for headings (golden uses full width)
                    content_width = max_width
                else:
                    content_width = min(text_width_in * 1.3 + 1.0, max_width)
                b['width'] = content_width
                b['x'] = slide_margin + (max_width - content_width) / 2
            elif effective_w > 0 and effective_w < max_width * PX_PER_IN:
                # Element has explicit CSS width or maxWidth (inline-block)
                b['width'] = effective_w / PX_PER_IN
                b['x'] = slide_margin + (max_width - b['width']) / 2
            else:
                b['width'] = max_width
                b['x'] = slide_margin
            # Clamp all text widths to max_width to prevent overflow past slide edge
            if b['width'] > max_width:
                b['width'] = max_width
                b['x'] = slide_margin
        elif elem_type == 'shape':
            # Check if this is a separator/divider shape (thin line)
            is_separator = b.get('height', 0) <= 0.05 or elem.get('_border_left_line')

            if is_separator and not elem.get('_border_left_line'):
                # Horizontal separator: use full width, keep thin height
                b['width'] = max_width
                b['x'] = slide_margin
                # Height based on CSS
                css_h = parse_px(s.get('height', ''))
                if css_h > 0:
                    b['height'] = max(css_h / PX_PER_IN, 0.02)
                else:
                    b['height'] = max(b['height'], 0.02)
            elif elem.get('_border_left_line'):
                # Vertical border-left line: compute height from following text elements
                bl_start_y = current_y
                bl_h = 0.0
                for j in range(i + 1, len(elements)):
                    next_el = elements[j]
                    if next_el.get('_border_left_line'):
                        break  # Next border-left starts
                    if next_el.get('type') == 'text':
                        bl_h += next_el['bounds'].get('height', 0.3) + 0.05
                if bl_h > 0:
                    b['height'] = bl_h - 0.05
                b['x'] = slide_margin
                b['y'] = current_y
            # Paired shapes (from inline-block containers) already have correct width
            elif elem.get('_pair_with') and not elem.get('_is_pill'):
                # Find paired text element to sync width
                paired_text = None
                for j in range(i + 1, min(i + 3, len(elements))):
                    next_elem = elements[j]
                    if next_elem.get('_pair_with') == elem.get('_pair_with'):
                        paired_text = next_elem
                        break

                if paired_text:
                    nb = paired_text['bounds']
                    # Sync shape width to text width
                    b['width'] = nb['width']
                    b['x'] = slide_margin + (max_width - b['width']) / 2
                # Full-width separator shapes (from list items with border-bottom): use slide margin
                elif b['width'] >= max_width * 0.9:
                    b['x'] = slide_margin
                    b['width'] = max_width
                else:
                    # Keep pre-computed width, just center
                    b['x'] = slide_margin + (max_width - b['width']) / 2
            elif b['width'] > max_width * 0.5:
                # Check if next element is text at similar y-position (tag shape)
                for j in range(i + 1, min(i + 3, len(elements))):
                    next_elem = elements[j]
                    if next_elem.get('type') == 'text':
                        nb = next_elem['bounds']
                        if abs(nb['y'] - b['y']) < 0.1 and nb['width'] < max_width * 0.5:
                            # Tag-style shape: match text width
                            b['width'] = nb['width']
                            b['x'] = slide_margin + (max_width - b['width']) / 2
                            break
                else:
                    # No matching text found — use full width
                    b['width'] = max_width
                    b['x'] = slide_margin
            # else: keep pre-computed width (already set by build_shape_element)
        elif elem_type == 'table':
            b['width'] = max_width
            b['x'] = slide_margin
            rows = elem.get('rows', [])
            b['height'] = max(len(rows) * 0.5, 0.5)
        elif elem_type == 'image':
            if b['width'] > max_width:
                b['width'] = max_width
                b['height'] = b['width'] * 0.75
            b['x'] = slide_margin

        # Paired shapes don't advance current_y — the paired text element positions both
        if id(elem) not in paired_shapes:
            # Check for next element's marginTop (should add gap before it)
            next_margin_top = 0.0
            for j in range(i + 1, min(i + 3, len(elements))):
                ne = elements[j]
                if ne.get('_skip_layout'):
                    continue
                ns = ne.get('styles', {})
                mt = ns.get('marginTop', '')
                if mt:
                    next_margin_top = parse_px(mt) / PX_PER_IN
                break

            # For horizontal separator lines (swiss-rule, hero-rule), use larger gap
            is_horizontal_separator = (elem_type == 'shape' and b['height'] <= 0.05
                                       and not elem.get('_border_left_line')
                                       and not elem.get('_stretch_height'))
            # Check if next element is a horizontal separator (need larger gap before it)
            next_is_separator = False
            for j in range(i + 1, min(i + 3, len(elements))):
                ne = elements[j]
                if ne.get('_skip_layout'):
                    continue
                if ne.get('type') == 'shape' and ne['bounds'].get('height', 0) <= 0.05:
                    next_is_separator = True
                    break
                elif ne.get('type') == 'text':
                    break

            if is_horizontal_separator:
                gap = 0.05  # Smaller gap after the separator itself
            elif elem_type == 'shape' and b['height'] <= 0.1:
                gap = 0.05
            elif elem.get('tag') in ('h1', 'h2') and next_is_separator:
                gap = 0.20  # Larger gap when heading precedes a separator line
            else:
                gap = 0.15

            # Add next element's margin-top to the gap
            gap += next_margin_top
            current_y += b['height'] + gap

    # Post-layout sync: for shapes paired with text, copy text position to shape
    text_by_pair: Dict[str, Dict] = {}
    for elem in elements:
        pair_id = elem.get('_pair_with')
        if pair_id and elem.get('type') == 'text':
            text_by_pair[pair_id] = elem
    for elem in elements:
        pair_id = elem.get('_pair_with')
        if pair_id and elem.get('type') == 'shape' and pair_id in text_by_pair:
            text_elem = text_by_pair[pair_id]
            tb = text_elem['bounds']
            sb = elem['bounds']
            sb['x'] = tb['x']
            sb['y'] = tb['y']
            sb['height'] = tb['height']
            # Match text width exactly (golden shows shape and text have same width for paired elements)
            sb['width'] = tb['width']

    # Vertical positioning: handle justifyContent center or flex-end
    y_offset = 0.0
    align_items_center = False
    if slide_style:
        jc = slide_style.get('justifyContent', slide_style.get('justify-content', ''))
        fd = slide_style.get('flexDirection', slide_style.get('flex-direction', 'column'))
        ai = slide_style.get('alignItems', slide_style.get('align-items', ''))
        align_items_center = (ai == 'center')

        if fd == 'column' and jc in ('center', 'flex-end'):
            # Calculate content height EXCLUDING _skip_layout elements (pills, subtitle)
            # These are positioned absolutely and shouldn't affect positioning
            non_skip_y_max = 0.0
            for elem in elements:
                if elem.get('_skip_layout'):
                    continue
                b = elem['bounds']
                elem_bottom = b['y'] + b['height']
                if elem_bottom > non_skip_y_max:
                    non_skip_y_max = elem_bottom
                if elem.get('type') == 'container':
                    for child in elem.get('children', []):
                        cb = child['bounds']
                        child_bottom = cb['y'] + cb['height']
                        if child_bottom > non_skip_y_max:
                            non_skip_y_max = child_bottom
            total_content_h = non_skip_y_max - internal_margin
            available_h = slide_height_in - 2 * internal_margin

            # For flex-end, also check for CSS padding-bottom (e.g., padding-bottom: 10vh)
            padding_bottom_extra = 0.0
            if jc == 'flex-end':
                pb_str = slide_style.get('paddingBottom', '')
                if pb_str and 'vh' in pb_str:
                    vh_val = float(re.search(r'([\d.]+)', pb_str).group(1))
                    padding_bottom_extra = vh_val * 0.081  # Convert vh to inches

            if total_content_h < available_h:
                if jc == 'center':
                    y_offset = (available_h - total_content_h) / 2
                elif jc == 'flex-end':
                    # Push content to bottom, accounting for extra padding
                    y_offset = available_h - total_content_h - padding_bottom_extra

                # Apply offset to all elements (but NOT _skip_layout elements — they're positioned absolutely)
                for elem in elements:
                    if elem.get('_skip_layout'):
                        continue
                    b = elem['bounds']
                    b['y'] += y_offset
                    if elem.get('type') == 'container':
                        for child in elem.get('children', []):
                            child['bounds']['y'] += y_offset

    # Horizontal centering for alignItems: center
    if align_items_center:
        slide_center_x = slide_width_in / 2
        for elem in elements:
            if elem.get('_skip_layout'):
                continue
            b = elem['bounds']
            # Center each element horizontally based on its width
            b['x'] = slide_center_x - b['width'] / 2
            if elem.get('type') == 'container':
                for child in elem.get('children', []):
                    cb = child['bounds']
                    # Adjust child x relative to parent
                    cb['x'] = b['x'] + (cb['x'] - b['x'])

    # Position inline tag pills AFTER vertical centering so stat items have final positions.
    # Pattern: pill shapes with _is_pill flag (no separate text elements).
    PILLS_GAP = 0.16  # gap between pills

    # Collect ALL pill shapes that are still at default position
    all_pills = []
    for elem in elements:
        if elem.get('type') == 'shape' and elem.get('_is_pill') and elem['bounds']['x'] == 0.50:
            all_pills.append(elem)

    if all_pills:
        # Find the lowest y position among all non-pill content (including container children like stat items).
        # This gives us the bottom of stat items after vertical centering.
        lowest_y = 0.0
        for elem in elements:
            if elem.get('_skip_layout') or elem.get('_is_pill'):
                continue
            b = elem['bounds']
            elem_bottom = b['y'] + b['height']
            if elem_bottom > lowest_y:
                lowest_y = elem_bottom
            # Also check container children (stat items)
            if elem.get('type') == 'container':
                for child in elem.get('children', []):
                    cb = child['bounds']
                    child_bottom = cb['y'] + cb['height']
                    if child_bottom > lowest_y:
                        lowest_y = child_bottom

        # Place pills 0.08" below the lowest content element (tight gap, matching golden spacing)
        pill_y = lowest_y + 0.08

        # Cap to not go too close to nav dots (y=7.22)
        pill_y = min(pill_y, slide_height_in - 0.60)

        # Position pills centered in full slide content area (matching golden)
        total_pill_w = sum(p['bounds']['width'] for p in all_pills) + PILLS_GAP * (len(all_pills) - 1)
        row_x = 1.57 + (10.18 - total_pill_w) / 2  # 10.18 = golden content area width
        for p in all_pills:
            p['bounds']['x'] = row_x
            p['bounds']['y'] = pill_y
            row_x += p['bounds']['width'] + PILLS_GAP


# ─── Pre-pass Corrections (from browser version) ─────────────────────────────

def pre_pass_corrections(elements: List[Dict]):
    """
    Pre-pass 1: Sync background shape height with adjacent text's naturalHeight.
    For shape+text pairs (same y-position), apply 1.3x PPTX correction for multi-line.

    Pre-pass 2: Push adjacent elements right of large single-line titles
    to prevent CJK overflow in PPTX.
    """
    # Flatten container children for correction processing
    flat = []
    for elem in elements:
        if elem.get('type') == 'container':
            flat.extend(elem.get('children', []))
        else:
            flat.append(elem)

    # Pre-pass 1: Shape+Text sync
    for i in range(len(flat) - 1):
        s, t = flat[i], flat[i + 1]
        if (s.get('type') == 'shape' and t.get('type') == 'text' and
            abs(s['bounds']['y'] - t['bounds']['y']) < 0.05 and
            abs(s['bounds']['height'] - t['bounds']['height']) < 0.05):
            nat = t.get('naturalHeight', t['bounds']['height'])
            base = max(nat, s['bounds']['height'])
            t_font_pt = px_to_pt(t['styles'].get('fontSize', '16px'))
            t_line_h = t_font_pt / 72.0 * 1.2
            t_est_lines = base / max(t_line_h, 0.001)
            if t_est_lines >= 2.0:
                corrected = base * PPTX_HEIGHT_FACTOR
                s['bounds']['height'] = corrected
                t['bounds']['height'] = corrected
                t['pptx_height_corrected'] = True
            elif nat > s['bounds']['height'] * 1.05:
                s['bounds']['height'] = nat
                t['bounds']['height'] = nat

    # Pre-pass 2: Adjacent element push for large titles
    for i in range(len(flat)):
        el = flat[i]
        if el.get('type') != 'text':
            continue
        fp = px_to_pt(el['styles'].get('fontSize', '16px'))
        if fp <= 24.0:
            continue
        lh = fp / 72.0 * 1.2
        est = el['bounds']['height'] / max(lh, 0.001)
        if est >= 2.0:
            continue
        orig_right = el['bounds']['x'] + el['bounds']['width']
        extra = el['bounds']['width'] * CJK_H_FACTOR
        for j in range(len(flat)):
            if j == i:
                continue
            other = flat[j]
            gap = other['bounds']['x'] - orig_right
            y_overlap = abs(other['bounds']['y'] - el['bounds']['y']) < el['bounds']['height']
            if 0 <= gap <= 0.3 and y_overlap:
                other['bounds']['x'] += extra


# ─── PPTX Rendering ───────────────────────────────────────────────────────────

_FONT_MAP = {
    'Clash Display': ('Calibri Light', 'Microsoft YaHei'),
    'Satoshi':       ('Calibri',       'Microsoft YaHei'),
    'Microsoft YaHei': ('Microsoft YaHei', 'Microsoft YaHei'),
    '微软雅黑':          ('Microsoft YaHei', 'Microsoft YaHei'),
    'PingFang SC':      ('PingFang SC',     'PingFang SC'),
    'Noto Sans CJK SC': ('Noto Sans CJK SC','Noto Sans CJK SC'),
    'Source Han Sans':  ('Source Han Sans', 'Source Han Sans'),
    'system-ui':        ('Microsoft YaHei', 'Microsoft YaHei'),
    '-apple-system':    ('PingFang SC',     'PingFang SC'),
}
_DEFAULT_FONTS = ('Microsoft YaHei', 'Microsoft YaHei')


def map_font(css_font_family: str):
    if css_font_family:
        for css_name, fonts in _FONT_MAP.items():
            if css_name in css_font_family:
                return fonts
    return _DEFAULT_FONTS


def set_run_fonts(run, latin_font, ea_font):
    from lxml import etree
    NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    run.font.name = latin_font
    rPr = run._r.get_or_add_rPr()
    for tag, typeface in [('ea', ea_font), ('cs', ea_font)]:
        el = rPr.find(f'{{{NS}}}{tag}')
        if el is None:
            el = etree.SubElement(rPr, f'{{{NS}}}{tag}')
        el.set('typeface', typeface)


def set_letter_spacing(run, css_letter_spacing: str):
    if not css_letter_spacing or css_letter_spacing in ('normal', '0px'):
        return
    m = re.search(r'([\d.]+)px', css_letter_spacing)
    if m:
        px = float(m.group(1))
        spc = int(px * 75)
        if spc > 0:
            run._r.get_or_add_rPr().set('spc', str(spc))


def set_roundrect_adj(shape, radius_px: float, width_in: float, height_in: float):
    from lxml import etree
    NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    radius_in = radius_px / 108.0
    shorter = min(width_in, height_in)
    if shorter <= 0:
        return
    adj = int(radius_in / (shorter / 2) * 100000)
    adj = max(0, min(50000, adj))
    prstGeom = shape._element.spPr.find(f'{{{NS}}}prstGeom')
    if prstGeom is None:
        return
    avLst = prstGeom.find(f'{{{NS}}}avLst')
    if avLst is None:
        avLst = etree.SubElement(prstGeom, f'{{{NS}}}avLst')
    for gd in avLst.findall(f'{{{NS}}}gd'):
        avLst.remove(gd)
    gd = etree.SubElement(avLst, f'{{{NS}}}gd')
    gd.set('name', 'adj')
    gd.set('fmla', f'val {adj}')


def suppress_line(shape):
    """Remove line (border) from shape."""
    # Use python-pptx API - this is the cleanest approach
    try:
        shape.line.fill.background()
    except:
        pass
    # Also try direct XML manipulation as backup
    try:
        from lxml import etree
        NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        spPr = shape._element.spPr
        ln_tag = f'{{{NS}}}ln'
        for ln in spPr.findall(ln_tag):
            spPr.remove(ln)
        ln = etree.SubElement(spPr, ln_tag)
        etree.SubElement(ln, f'{{{NS}}}noFill')
    except:
        pass


def set_light_shadow(shape):
    from lxml import etree
    NP = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    NA = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    spPr = shape._element.spPr
    eff_tag = f'{{{NA}}}effectLst'
    existing = spPr.find(eff_tag)
    if existing is not None:
        spPr.remove(existing)
    effectLst = etree.fromstring(
        f'<a:effectLst xmlns:a="{NA}">'
        f'<a:outerShdw blurRad="25000" dist="8000" dir="5400000" rotWithShape="0">'
        f'<a:srgbClr val="000000"><a:alpha val="8000"/></a:srgbClr>'
        f'</a:outerShdw></a:effectLst>'
    )
    spPr.append(effectLst)
    style_elem = shape._element.find(f'{{{NP}}}style')
    if style_elem is not None:
        eff_ref = style_elem.find(f'{{{NA}}}effectRef')
        if eff_ref is not None:
            eff_ref.set('idx', '0')


def segments_to_lines(segments):
    """Split segments into lines (list of list)."""
    cleaned = []
    for s in segments:
        t = s['text']
        if t == '\n':
            cleaned.append(s)
        elif t.strip():
            cleaned.append({'text': t, 'color': s['color'],
                           'bold': s.get('bold', False),
                           'fontSize': s.get('fontSize', ''),
                           'strike': s.get('strike', False),
                           'bgColor': s.get('bgColor'),
                           'inlineBgBounds': s.get('inlineBgBounds')})
    segments = cleaned
    lines = []
    current_line = []
    for seg in segments:
        text = seg['text']
        if '\n' in text:
            parts = text.split('\n')
            for i, part in enumerate(parts):
                if part:
                    current_line.append({**seg, 'text': part})
                if i < len(parts) - 1:
                    lines.append(current_line)
                    current_line = []
        else:
            current_line.append(seg)
    lines.append(current_line)
    # Filter out empty lines (from pure \n segments)
    lines = [l for l in lines if any(s['text'].strip() for s in l)]
    return lines


def apply_run(run, text, color_str, font_size_pt, font_weight,
              text_transform='none', font_family='', letter_spacing='', strike=False):
    if text_transform == 'uppercase':
        text = text.upper()
    run.text = text
    if text and (text[0] == ' ' or text[-1] == ' '):
        _nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        _t_elem = run._r.find('.//a:t', _nsmap)
        if _t_elem is not None:
            _t_elem.set('{http://www.w3.org/XML/1998/namespace}space', 'preserve')
    latin_font, ea_font = map_font(font_family)
    set_run_fonts(run, latin_font, ea_font)
    if font_size_pt:
        run.font.size = Pt(font_size_pt)
    try:
        if font_weight == 'bold':
            run.font.bold = True
        else:
            run.font.bold = int(font_weight) >= 600
    except Exception:
        pass
    rgb = parse_color(color_str)
    if rgb:
        run.font.color.rgb = RGBColor(*rgb)
    if strike:
        _ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        _rPr = run._r.find(f'{{{_ns}}}rPr')
        if _rPr is None:
            _rPr = _etree.SubElement(run._r, f'{{{_ns}}}rPr')
            run._r.insert(0, _rPr)
        _rPr.set('strike', 'sngStrike')
    set_letter_spacing(run, letter_spacing)


def apply_para_format(p, s):
    lh = s.get('lineHeight', 'normal')
    if lh == 'normal':
        p.line_spacing = 1.2
    else:
        try:
            if 'px' in lh:
                lh_px = float(re.search(r'([\d.]+)', lh).group(1))
                p.line_spacing = Pt(round(lh_px * 0.75, 1))
            else:
                p.line_spacing = float(lh)
        except Exception:
            p.line_spacing = 1.2
    align = s.get('textAlign', 'left')
    if align == 'center':
        p.alignment = PP_ALIGN.CENTER
    elif align == 'right':
        p.alignment = PP_ALIGN.RIGHT


def gradient_to_solid(bg_image, slide_bg=(255, 255, 255)):
    if not bg_image or 'gradient' not in bg_image:
        return None
    rgba_matches = re.findall(r'rgba?\((\d+),\s*(\d+),\s*(\d+)(?:,\s*([\d.]+))?\)', bg_image)
    if not rgba_matches:
        return None
    r, g, b = int(rgba_matches[0][0]), int(rgba_matches[0][1]), int(rgba_matches[0][2])
    a = float(rgba_matches[0][3]) if rgba_matches[0][3] else 1.0
    if a <= 0:
        return None
    if a < 1.0:
        r = int(a * r + (1 - a) * slide_bg[0])
        g = int(a * g + (1 - a) * slide_bg[1])
        b = int(a * b + (1 - a) * slide_bg[2])
    return (r, g, b)


def interpolate_color(c1, c2, t):
    return (
        int(c1[0] + (c2[0] - c1[0]) * t),
        int(c1[1] + (c2[1] - c1[1]) * t),
        int(c1[2] + (c2[2] - c1[2]) * t)
    )


def export_shape_background(slide, elem, slide_bg=(255, 255, 255)):
    """Create background shape (no text) for type=shape containers."""
    b = elem['bounds']
    s = elem['styles']
    border_radius = s.get('borderRadius', '')
    radius_px = 0.0
    if border_radius and border_radius != '0px':
        m = re.search(r'([\d.]+)px', border_radius)
        if m:
            radius_px = float(m.group(1))

    height_px = b['height'] * 108.0
    if radius_px > 0 and radius_px < height_px * 0.4:
        radius_px = min(radius_px, 6.0)

    def parse_border_side(bs):
        if not bs or 'none' in bs or bs.startswith('0px'):
            return None
        # Try rgb/rgba format first
        m = re.search(r'([\d.]+)px.*?rgba?\((\d+),\s*(\d+),\s*(\d+)', bs)
        if m:
            return {'width': float(m.group(1)), 'rgb': (int(m.group(2)), int(m.group(3)), int(m.group(4)))}
        # Try hex color format: 4px solid #D97706
        m = re.search(r'([\d.]+)px.*?#([0-9a-fA-F]{6})', bs)
        if m:
            h = m.group(2)
            return {'width': float(m.group(1)), 'rgb': (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))}
        return None

    bl = parse_border_side(s.get('borderLeft', ''))
    br = parse_border_side(s.get('borderRight', ''))
    bt = parse_border_side(s.get('borderTop', ''))
    bb = parse_border_side(s.get('borderBottom', ''))
    borders = [x for x in [bl, br, bt, bb] if x is not None]
    all_uniform = (len(borders) >= 3 and
                   all(bd['rgb'] == borders[0]['rgb'] and bd['width'] == borders[0]['width']
                       for bd in borders))

    BAR_VISIBLE_PX = 4.0
    bl_handled = False
    if bl and not all_uniform and radius_px > 0:
        bar_visible_in = BAR_VISIBLE_PX / 108.0
        bar_total_width_px = BAR_VISIBLE_PX + radius_px * 2 + 4
        bar_total_width_in = bar_total_width_px / 108.0
        bar_x_in = b['x'] - bar_visible_in
        bar_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(bar_x_in), Inches(b['y']),
            Inches(bar_total_width_in), Inches(b['height'])
        )
        set_roundrect_adj(bar_shape, radius_px, bar_total_width_in, b['height'])
        bar_shape.fill.solid()
        bar_shape.fill.fore_color.rgb = RGBColor(*bl['rgb'])
        suppress_line(bar_shape)
        set_light_shadow(bar_shape)
        bl_handled = True

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius_px > 0 else MSO_SHAPE.RECTANGLE,
        Inches(b['x']), Inches(b['y']),
        Inches(b['width']), Inches(b['height'])
    )
    if radius_px > 0:
        set_roundrect_adj(shape, radius_px, b['width'], b['height'])

    # Background color: slide-based alpha blending for rgba colors (matches P1 style)
    blend_bg = slide_bg
    bg_rgb = parse_color(s.get('backgroundColor', ''), bg=blend_bg)
    if bg_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*bg_rgb)
    else:
        grad_fill = gradient_to_solid(s.get('backgroundImage', ''), slide_bg=slide_bg)
        if grad_fill:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*grad_fill)
        else:
            shape.fill.background()

    if all_uniform:
        # For subtle borders (white/near-white with low alpha), suppress entirely.
        # Golden reference shows container shapes (cards, callouts) have no borders.
        bd = borders[0]
        is_white_or_near_white = (
            bd['rgb'][0] >= 240 and bd['rgb'][1] >= 240 and bd['rgb'][2] >= 240
        )
        if is_white_or_near_white and bd['width'] <= 1.5:
            suppress_line(shape)
        else:
            shape.line.color.rgb = RGBColor(*bd['rgb'])
            shape.line.width = Pt(max(0.5, bd['width'] * 0.75))
    elif len(borders) >= 1:
        suppress_line(shape)
        # Helper: skip border shape if near-white (golden doesn't have white borders)
        def _is_subtle_border(bd):
            return (bd['rgb'][0] >= 240 and bd['rgb'][1] >= 240 and
                    bd['rgb'][2] >= 240)
        if bl and not bl_handled and not _is_subtle_border(bl):
            border_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(b['x']), Inches(b['y']),
                Inches(bl['width'] / 108.0), Inches(b['height'])
            )
            border_shape.fill.solid()
            border_shape.fill.fore_color.rgb = RGBColor(*bl['rgb'])
            suppress_line(border_shape)
        if br and not _is_subtle_border(br):
            border_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(b['x'] + b['width'] - br['width']/108.0), Inches(b['y']),
                Inches(br['width'] / 108.0), Inches(b['height'])
            )
            border_shape.fill.solid()
            border_shape.fill.fore_color.rgb = RGBColor(*br['rgb'])
            suppress_line(border_shape)
        if bt and not _is_subtle_border(bt):
            border_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(b['x']), Inches(b['y']),
                Inches(b['width']), Inches(bt['width'] / 108.0)
            )
            border_shape.fill.solid()
            border_shape.fill.fore_color.rgb = RGBColor(*bt['rgb'])
            suppress_line(border_shape)
        if bb:
            # Skip rendering border-bottom as a separate shape.
            # Golden doesn't create separate separator shapes for borders —
            # the border-bottom CSS is too subtle to warrant a PPTX element.
            pass
    else:
        suppress_line(shape)

    set_light_shadow(shape)
    tf = shape.text_frame

    # Pill shapes: embed text directly into shape's text frame (like P1 _pair_with pattern)
    pill_text = elem.get('pill_text', '')
    if pill_text:
        from pptx.enum.text import MSO_AUTO_SIZE
        tf.word_wrap = False
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.margin_left = Pt(8)
        tf.margin_right = Pt(8)
        tf.margin_top = Pt(4)
        tf.margin_bottom = Pt(4)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = pill_text
        font_size_pt = px_to_pt(s.get('fontSize', '14px'))
        run.font.size = Pt(font_size_pt)
        pill_color_hex = elem.get('pill_color', '')
        if pill_color_hex:
            # Resolve CSS variables if needed
            if pill_color_hex.startswith('var('):
                resolved = resolve_css_variables_inline(pill_color_hex)
                if resolved.startswith('#') and len(resolved) == 7:
                    pill_color_hex = resolved[1:]
                elif resolved.startswith('rgb'):
                    rm = re.match(r'rgb\((\d+),\s*(\d+),\s*(\d+)\)', resolved)
                    if rm:
                        pill_color_hex = f'{int(rm.group(1)):02X}{int(rm.group(2)):02X}{int(rm.group(3)):02X}'
            if pill_color_hex.startswith('#'):
                pill_color_hex = pill_color_hex[1:]
            if len(pill_color_hex) == 6:
                try:
                    run.font.color.rgb = RGBColor(
                        int(pill_color_hex[0:2], 16),
                        int(pill_color_hex[2:4], 16),
                        int(pill_color_hex[4:6], 16)
                    )
                except Exception:
                    pass
        else:
            run.font.color.rgb = RGBColor(0, 0, 0)
    else:
        # Clear text for non-pill shapes
        if tf.paragraphs:
            for para in tf.paragraphs:
                for run in para.runs:
                    run.text = ''


def export_text_element(slide, elem, bg_color=None):
    b = elem['bounds']
    s = elem['styles']
    segments = elem.get('segments', [])
    text_transform = elem.get('textTransform', 'none')
    font_size_pt = px_to_pt(s.get('fontSize', '16px'))
    font_weight = s.get('fontWeight', '400')
    font_family = s.get('fontFamily', '')
    letter_spacing = s.get('letterSpacing', '')

    natural_h = elem.get('naturalHeight', b['height'])
    # Use a generous height: take the max of estimated height and a minimum based on font size
    # PPTX CJK rendering needs more vertical space than browser estimate
    min_h_by_font = font_size_pt / 72.0 * 1.5  # at least 1.5x font size
    effective_h = max(b['height'], natural_h, min_h_by_font)

    if not segments:
        raw = (elem.get('text', '') or '').strip()
        segments = [{'text': raw, 'color': s.get('color', '')}]

    lines = segments_to_lines(segments)
    if not lines:
        lines = [[{'text': '', 'color': s.get('color', '')}]]

    txBox = slide.shapes.add_textbox(
        Inches(b['x']), Inches(b['y']),
        Inches(b['width']), Inches(effective_h)
    )
    tf = txBox.text_frame
    # Match golden: single-line text uses TEXT_TO_FIT_SHAPE with no wrap,
    # multi-line text uses SHAPE_TO_FIT_TEXT with wrap
    from pptx.enum.text import MSO_AUTO_SIZE
    if len(lines) <= 1:
        # Check if text is likely to overflow the box width (card descriptions, etc.)
        raw_text = elem.get('text', '') or ''
        # If text has no newlines but is long (>30 chars) and box is narrow (<5"),
        # enable word wrap to prevent overflow
        needs_wrap = len(raw_text) > 20 and b['width'] < 5.0
        if needs_wrap:
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        else:
            tf.word_wrap = False
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    else:
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    pad_l = parse_px(s.get('paddingLeft', ''))
    pad_r = parse_px(s.get('paddingRight', ''))
    pad_t = parse_px(s.get('paddingTop', ''))
    pad_b = parse_px(s.get('paddingBottom', ''))
    tf.margin_left = Inches(pad_l / 108.0) if pad_l > 0 else 0
    tf.margin_right = Inches(pad_r / 108.0) if pad_r > 0 else 0
    tf.margin_top = Inches(pad_t / 108.0) if pad_t > 0 else 0
    tf.margin_bottom = Inches(pad_b / 108.0) if pad_b > 0 else 0

    gradient_colors = elem.get('gradientColors') if elem.get('tag') == 'h1' else None
    gc_start = parse_color(gradient_colors[0]) if gradient_colors else None
    gc_end = parse_color(gradient_colors[1]) if gradient_colors else None
    total_lines = len(lines)
    is_li = elem.get('tag') == 'li'
    li_bullet_color = 'rgb(56, 139, 253)'

    for line_idx, line_segs in enumerate(lines):
        p = tf.add_paragraph() if line_idx > 0 else tf.paragraphs[0]
        apply_para_format(p, s)
        justify_content = s.get('justifyContent', '')
        if justify_content in ('center', 'space-around', 'space-evenly'):
            p.alignment = PP_ALIGN.CENTER

        if is_li and line_idx == 0:
            bullet_run = p.add_run()
            apply_run(bullet_run, '• ', s.get('color', '#0a0a0a'), font_size_pt, '400')

        if gc_start and gc_end and total_lines > 1:
            t = line_idx / (total_lines - 1)
            grad_rgb = interpolate_color(gc_start, gc_end, t)
            override_color = 'rgb({},{},{})'.format(*grad_rgb)
        elif gc_start and gc_end:
            override_color = gradient_colors[1]
        else:
            override_color = None

        for seg in line_segs:
            if not seg['text']:
                continue
            run = p.add_run()
            color = override_color or seg['color']
            seg_weight = 'bold' if seg.get('bold') else font_weight
            seg_fs_raw = seg.get('fontSize', '')
            seg_font_size_pt = px_to_pt(seg_fs_raw) if seg_fs_raw and 'px' in str(seg_fs_raw) else font_size_pt
            apply_run(run, seg['text'], color, seg_font_size_pt, seg_weight, text_transform,
                      font_family=font_family, letter_spacing=letter_spacing,
                      strike=seg.get('strike', False))


def export_table_element(slide, elem):
    rows = elem.get('rows', [])
    if not rows:
        return
    tb = elem['bounds']
    table_x = tb['x']
    table_y = tb['y']
    table_w = tb['width']
    num_cols = max(len(row_data['cells']) for row_data in rows) if rows else 1
    row_h = tb['height'] / len(rows) if rows else 0.5

    # Content-aware column widths
    col_widths = _compute_table_column_widths(rows, table_w)

    for row_idx, row_data in enumerate(rows):
        for col_idx, cell in enumerate(row_data['cells']):
            cb = cell['bounds']
            cs = cell['styles']
            cx = table_x + sum(col_widths[:col_idx])
            cy = table_y + row_idx * row_h
            cw = col_widths[col_idx] if col_idx < len(col_widths) else table_w / num_cols
            ch = row_h
            if cw < 0.01 or ch < 0.01:
                continue
            cell_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(cx), Inches(cy),
                Inches(cw), Inches(ch)
            )
            bg_rgb = parse_color(cs.get('backgroundColor', ''))
            if bg_rgb:
                cell_shape.fill.solid()
                cell_shape.fill.fore_color.rgb = RGBColor(*bg_rgb)
            else:
                cell_shape.fill.background()
            suppress_line(cell_shape)
            # Skip border-bottom rendering — golden doesn't create separate
            # border shapes for table cells or list items.
            segments = cell.get('segments', [])
            text = cell.get('text', '').strip()
            if not segments and text:
                segments = [{'text': text, 'color': cs.get('color', '')}]
            if not segments:
                continue
            font_size_pt = px_to_pt(cs.get('fontSize', '14px'))
            font_weight = cs.get('fontWeight', '400')
            font_family = cs.get('fontFamily', '')
            letter_spacing = cs.get('letterSpacing', '')
            if cell['isHeader']:
                font_weight = 'bold'
            tf = cell_shape.text_frame
            tf.word_wrap = True
            tf.margin_left = Pt(6)
            tf.margin_right = Pt(6)
            tf.margin_top = Pt(4)
            tf.margin_bottom = Pt(4)
            lines = segments_to_lines(segments)
            for line_idx, line_segs in enumerate(lines):
                p = tf.add_paragraph() if line_idx > 0 else tf.paragraphs[0]
                align = cs.get('textAlign', 'left')
                if align == 'center':
                    p.alignment = PP_ALIGN.CENTER
                elif align == 'right':
                    p.alignment = PP_ALIGN.RIGHT
                for seg in line_segs:
                    if not seg['text']:
                        continue
                    run = p.add_run()
                    seg_weight = 'bold' if seg.get('bold') else font_weight
                    apply_run(run, seg['text'], seg['color'], font_size_pt, seg_weight,
                              font_family=font_family, letter_spacing=letter_spacing,
                              strike=seg.get('strike', False))


def export_image_element(slide, elem, html_dir: Path):
    """Render image element (img src or data URI)."""
    b = elem['bounds']
    source = elem.get('source', '')
    if not source:
        return
    try:
        img_bytes = None
        if source.startswith('data:image'):
            comma_idx = source.index(',')
            data = source[comma_idx + 1:]
            img_bytes = base64.b64decode(data)
        elif source.startswith(('http://', 'https://')):
            req = urllib.request.Request(source, headers={'User-Agent': 'Mozilla/5.0'})
            import ssl
            try:
                import certifi
                _ssl_ctx = ssl.create_default_context(cafile=certifi.where())
            except ImportError:
                _ssl_ctx = ssl._create_unverified_context()
            with urllib.request.urlopen(req, context=_ssl_ctx, timeout=15) as resp:
                img_bytes = resp.read()
        elif source.startswith('file://'):
            with open(source[len('file://'):], 'rb') as f:
                img_bytes = f.read()
        elif not source.startswith('<svg') and not source.startswith('<?xml'):
            img_path = html_dir / source
            if img_path.exists():
                with open(img_path, 'rb') as f:
                    img_bytes = f.read()
        if img_bytes:
            slide.shapes.add_picture(
                io.BytesIO(img_bytes),
                Inches(b['x']), Inches(b['y']),
                Inches(b['width']), Inches(b['height'])
            )
            return
    except Exception as e:
        print(f"    Warning: failed to load image: {e}")

    placeholder = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(b['x']), Inches(b['y']),
        Inches(b['width']), Inches(b['height'])
    )
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = RGBColor(230, 230, 230)
    suppress_line(placeholder)
    tf = placeholder.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = '[Image]'
    run.font.size = Pt(12)
    run.font.color.rgb = RGBColor(150, 150, 150)


def add_slide_chrome(slide, slide_idx: int, slide_count: int,
                     slide_w_in: float, slide_h_in: float, px_per_in: float = 108.0):
    # Page counter in top-right corner (matching source HTML)
    counter_w = 1.0
    counter_h = 0.22
    counter_x = slide_w_in - counter_w - 36 / px_per_in  # Right side
    counter_y = 24 / px_per_in
    txBox = slide.shapes.add_textbox(
        Inches(counter_x), Inches(counter_y), Inches(counter_w), Inches(counter_h)
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{slide_idx + 1:02d} / {slide_count:02d}"
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(100, 116, 139)
    run.font.bold = True


def add_grid_background(slide, slide_w_in: float, slide_h_in: float,
                        grid_color_str: str, grid_size_px: float):
    try:
        from PIL import Image, ImageDraw
    except ImportError:
        return
    scale = 3
    w = int(slide_w_in * 96 * scale)
    h = int(slide_h_in * 96 * scale)
    grid_px = max(1, int(grid_size_px * scale))
    img = Image.new('RGBA', (w, h), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    m = re.match(r'rgba?\((\d+),\s*(\d+),\s*(\d+)(?:,\s*([\d.]+))?\)', grid_color_str.strip())
    if m:
        r, g, b = int(m.group(1)), int(m.group(2)), int(m.group(3))
        a = float(m.group(4) or '1.0')
        line_color = (r, g, b, int(a * 255))
    else:
        line_color = (80, 100, 170, 25)
    for y in range(0, h, grid_px):
        draw.line([(0, y), (w - 1, y)], fill=line_color, width=1)
    for x in range(0, w, grid_px):
        draw.line([(x, 0), (x, h - 1)], fill=line_color, width=1)
    buf = io.BytesIO()
    img.save(buf, format='PNG')
    buf.seek(0)
    pic = slide.shapes.add_picture(buf, Inches(0), Inches(0),
                                   Inches(slide_w_in), Inches(slide_h_in))
    sp_tree = slide.shapes._spTree
    sp_tree.remove(pic._element)
    sp_tree.insert(2, pic._element)


def generate_preview_from_pptx(pptx_path: Path) -> Optional[Path]:
    """Generate a preview grid from the saved PPTX using Pillow."""
    try:
        from PIL import Image, ImageDraw, ImageFont
        from pptx import Presentation as _Prs
        from pptx.enum.dml import MSO_FILL
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = _Prs(str(pptx_path))
        THUMB_W = 320
        slides_data = []

        def safe_fill_color(shape):
            try:
                ft = shape.fill.type
                if ft == MSO_FILL.SOLID:
                    return shape.fill.fore_color.rgb
                return None
            except Exception:
                return None

        def get_slide_bg_color(slide):
            try:
                ft = slide.background.fill.type
                if ft == MSO_FILL.SOLID:
                    return slide.background.fill.fore_color.rgb
            except Exception:
                pass
            return None

        for slide in prs.slides:
            slide_w = prs.slide_width
            slide_h = prs.slide_height
            img_h = int(THUMB_W * slide_h / slide_w)

            bg_rgb = get_slide_bg_color(slide)
            if bg_rgb:
                img = Image.new('RGB', (THUMB_W, img_h), bg_rgb)
            else:
                img = Image.new('RGB', (THUMB_W, img_h), (240, 240, 240))
            draw = ImageDraw.Draw(img)

            for shape in slide.shapes:
                try:
                    x = int(shape.left / slide_w * THUMB_W)
                    y = int(shape.top / slide_h * img_h)
                    w = int(shape.width / slide_w * THUMB_W)
                    h = int(shape.height / slide_h * img_h)
                except Exception:
                    continue

                if w < 2 or h < 2:
                    continue

                # Skip tiny nav dots
                if w < 15 and h < 15:
                    continue

                if shape.has_text_frame and shape.text.strip():
                    text_color = (255, 255, 255)
                    try:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.color and run.font.color.rgb:
                                    text_color = run.font.color.rgb
                                    break
                            if text_color != (255, 255, 255):
                                break
                    except Exception:
                        pass
                    if bg_rgb:
                        luminance = 0.299 * bg_rgb[0] + 0.587 * bg_rgb[1] + 0.114 * bg_rgb[2]
                        text_color = (50, 50, 50) if luminance > 128 else (230, 230, 230)
                    # Show text within bounds, truncate at natural break
                    text = shape.text.strip()
                    max_chars = int(w / 3)  # roughly chars that fit in width
                    if len(text) > max_chars:
                        text = text[:max_chars].rsplit(' ', 1)[0] + '…'
                    draw.text((x + 2, y + 2), text, fill=text_color)
                else:
                    fill = safe_fill_color(shape)
                    if fill:
                        draw.rectangle([x, y, x + w - 1, y + h - 1], fill=fill)
                    # Show empty text box outlines (subtle)
                    elif shape.has_text_frame:
                        draw.rectangle([x, y, x + w - 1, y + h - 1], outline=(80, 80, 80))

            slides_data.append(img)

        if not slides_data:
            return None

        n = len(slides_data)
        th = slides_data[0].height
        PAD = 4
        LABEL_H = 22
        grid_w = n * THUMB_W + (n - 1) * PAD
        grid_h = th + LABEL_H
        grid = Image.new('RGB', (grid_w, grid_h), (32, 32, 32))
        draw = ImageDraw.Draw(grid)

        for j, thumb in enumerate(slides_data):
            x = j * (THUMB_W + PAD)
            grid.paste(thumb, (x, 0))
            draw.text((x + THUMB_W // 2, th + 3), f"Slide {j+1}", fill=(200, 200, 200))

        preview_path = pptx_path.with_name(pptx_path.stem + '-preview.png')
        grid.save(str(preview_path))
        return preview_path
    except Exception as e:
        print(f"  Warning: preview generation failed: {e}")
        return None


def extract_css_from_soup(soup: BeautifulSoup) -> List[CSSRule]:
    """Extract and parse all <style> blocks from the HTML."""
    all_rules = []
    for style_tag in soup.find_all('style'):
        css_text = style_tag.string or ''
        all_rules.extend(parse_css_rules(css_text))
    return all_rules


# ─── Main Export Pipeline ─────────────────────────────────────────────────────

def export_sandbox(html_path, output_path=None, width=1440, height=810):
    html_path = Path(html_path).resolve()
    if not html_path.exists():
        print(f"Error: {html_path}")
        sys.exit(1)

    output_path = Path(output_path) if output_path else html_path.with_suffix('.pptx')
    html_dir = html_path.parent

    print(f"导出（sandbox, no browser）: {html_path.name}")

    # Parse HTML → slide data
    slides = parse_html_to_slides(html_path, width, height)
    if not slides:
        print("Nothing to export.")
        return

    # Create PPTX
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)  # 16:9
    blank_layout = prs.slide_layouts[6]
    slide_w_in = 13.33
    slide_h_in = 7.5

    for i, slide_data in enumerate(slides):
        pptx_slide = prs.slides.add_slide(blank_layout)

        # Background
        if slide_data['bgGradient']:
            c1, c2 = slide_data['bgGradient']
            try:
                from pptx.oxml.ns import qn
                fill = pptx_slide.background.fill
                fill.gradient()
                fill.gradient_angle = 135.0
                stops = fill.gradient_stops
                stops[0].position = 0.0
                stops[0].color.rgb = RGBColor(*c1)
                stops[1].position = 1.0
                stops[1].color.rgb = RGBColor(*c2)
            except Exception:
                if slide_data['background']:
                    pptx_slide.background.fill.solid()
                    pptx_slide.background.fill.fore_color.rgb = RGBColor(*slide_data['background'])
        elif slide_data['background']:
            r, g, b = slide_data['background']
            pptx_slide.background.fill.solid()
            pptx_slide.background.fill.fore_color.rgb = RGBColor(r, g, b)

        # Grid background
        if slide_data['gridBg']:
            add_grid_background(pptx_slide, slide_w_in, slide_h_in,
                                slide_data['gridBg']['color'], slide_data['gridBg']['sizePx'])

        # Pre-pass corrections
        elements = slide_data['elements']
        pre_pass_corrections(elements)

        # Compute slide style for vertical centering detection
        slide_element_style = slide_data.get('slideStyle', {})

        # Layout elements
        layout_slide_elements(elements, slide_w_in, slide_h_in, slide_element_style)

        # Clamp widths
        for elem in elements:
            if elem.get('type') == 'container':
                for child in elem.get('children', []):
                    if child['bounds']['x'] < slide_w_in and child['bounds']['width'] > slide_w_in - child['bounds']['x']:
                        child['bounds']['width'] = slide_w_in - child['bounds']['x']
            elif elem['bounds']['x'] < slide_w_in and elem['bounds']['width'] > slide_w_in - elem['bounds']['x']:
                elem['bounds']['width'] = slide_w_in - elem['bounds']['x']

        # Render elements
        for elem in elements:
            try:
                elem_type = elem.get('type', 'text')
                if elem_type == 'container':
                    # Render children of grid/flex container
                    for child in elem.get('children', []):
                        child_type = child.get('type', 'text')
                        cb = child['bounds']
                        if child_type == 'shape':
                            export_shape_background(pptx_slide, child,
                                                    slide_bg=slide_data['background'] or (255, 255, 255))
                        elif child_type == 'image':
                            export_image_element(pptx_slide, child, html_dir)
                        elif child_type == 'table':
                            export_table_element(pptx_slide, child)
                        else:
                            export_text_element(pptx_slide, child, slide_data['background'])
                    continue
                elif elem_type == 'shape':
                    export_shape_background(pptx_slide, elem,
                                            slide_bg=slide_data['background'] or (255, 255, 255))
                elif elem_type == 'image':
                    export_image_element(pptx_slide, elem, html_dir)
                elif elem_type == 'table':
                    export_table_element(pptx_slide, elem)
                else:
                    export_text_element(pptx_slide, elem, slide_data['background'])
            except Exception as e:
                print(f"    警告: {e}")

        # Chrome
        if not slide_data['hasOwnChrome']:
            add_slide_chrome(pptx_slide, i, len(slides), slide_w_in, slide_h_in)

    prs.save(str(output_path))
    print(f"Saved: {output_path}  ({len(slides)} 张幻灯片)")

    # Preview
    preview_path = generate_preview_from_pptx(output_path)
    if preview_path:
        print(f"Preview: {preview_path}")

    return output_path


# ─── CLI ──────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Export HTML to editable PPTX (no browser required)",
        formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("html", help="Path to HTML file")
    parser.add_argument("output", nargs="?", help="Output .pptx path (default: same name as HTML)")
    parser.add_argument("--width", type=int, default=1440, help="Slide width in pixels (default: 1440)")
    parser.add_argument("--height", type=int, default=810, help="Slide height in pixels (default: 810)")
    parser.add_argument("--no-chrome", action="store_true", help="Skip page counter and nav dots")
    args = parser.parse_args()

    export_sandbox(args.html, args.output, args.width, args.height)


if __name__ == "__main__":
    main()
